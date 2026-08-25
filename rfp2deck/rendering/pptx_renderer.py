from __future__ import annotations

import math
import re
from io import BytesIO
from pathlib import Path
from typing import Any, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from rfp2deck.core.schemas import DeckPlan, DiagramSpec
from rfp2deck.core.logging import get_logger

log = get_logger(__name__)

# Layout constants (legacy defaults, retained for the no-coordinate fallback path)
LEFT_MARGIN_IN = 0.75
RIGHT_MARGIN_IN = 0.75
TOP_MARGIN_IN = 0.45
BOTTOM_MARGIN_IN = 0.45

TITLE_X_IN = LEFT_MARGIN_IN
TITLE_Y_IN = TOP_MARGIN_IN
TITLE_H_IN = 0.70

# Fonts
# HCLTech brand typeface. Roobert renders on machines that have it (HCLTech
# laptops); elsewhere PowerPoint substitutes a default sans (Calibri-class).
FONT_NAME = "HCLTech Roobert Light"        # body text
FONT_NAME_HEAVY = "HCLTech Roobert SemiBold"  # titles / emphasis / headings
FONT_NAME_MONO = "Consolas"                # code / diagram captions
FONT_TITLE_PT = 28
FONT_TITLE_MIN_PT = 18
FONT_TITLE_SLIDE_PT = 44
FONT_BODY_START_PT = 16
FONT_BODY_MIN_PT = 11

EMU_PER_INCH = 914400

# ------------------------------------------------------------------
# Brand palette (sampled from the official HCLTech template + the
# SATS reference deck). The renderer draws every shape with explicit
# RGB values so the look is independent of the .pptx theme.
# ------------------------------------------------------------------
COLOR_BRAND = "5F1EBE"        # HCLTech purple — cover / section identity
COLOR_BRAND_DK = "411482"     # deep purple
COLOR_NAVY = "132A45"         # primary heading + body text on light
COLOR_NAVY_DK = "0A2540"      # banner / key-message fill
COLOR_TEAL = "12B5C9"         # teal accent (key messages, highlights)
COLOR_BLUE = "1E6FB0"         # medium blue (bullets, "solution")
COLOR_CORAL = "FF6B5C"        # coral (challenge / warning)
COLOR_GREEN = "1FA97A"        # green (outcomes / value)
COLOR_PURPLE = "5A2D9C"       # purple accent (why HCLTech)
COLOR_CARD_BG = "FFFFFF"      # card fill
COLOR_CARD_TINT = "F4F7FB"    # pale blue card / chip background
COLOR_CARD_LINE = "E2E7F0"    # hairline card border
COLOR_WHITE = "FFFFFF"
COLOR_BODY = "132A45"         # body text (navy, per reference)
COLOR_BODY_MUTED = "44566B"   # secondary / supporting body text
COLOR_BODY_LIGHT = "DCE6F2"   # body text on dark backgrounds
COLOR_TINT = "CADCFC"         # pale blue (decorative dot grid on dark slides)
COLOR_FOOTER = "8A93A8"       # muted slate for the footer line

# Back-compat aliases for the legacy code paths.
COLOR_PRIMARY = COLOR_BRAND
COLOR_ACCENT = COLOR_TEAL
COLOR_ACCENT_ALT = COLOR_GREEN
COLOR_HEADER_TEXT = COLOR_NAVY  # titles now sit on white, not a navy bar

# Semantic accent keys (from card/comparison specs) -> brand colours.
ACCENT_MAP = {
    "challenge": COLOR_CORAL,
    "problem": COLOR_CORAL,
    "risk": COLOR_CORAL,
    "pain": COLOR_CORAL,
    "solution": COLOR_BLUE,
    "approach": COLOR_BLUE,
    "info": COLOR_TEAL,
    "goal": COLOR_TEAL,
    "why": COLOR_PURPLE,
    "differentiator": COLOR_PURPLE,
    "outcome": COLOR_GREEN,
    "value": COLOR_GREEN,
    "benefit": COLOR_GREEN,
    "neutral": COLOR_BLUE,
}
# Default colour rotation when a card carries no accent.
ACCENT_CYCLE = [COLOR_BLUE, COLOR_TEAL, COLOR_PURPLE, COLOR_GREEN]


def _resolve_accent(accent: Optional[str], idx: int = 0) -> str:
    """Map a semantic accent key or hex string to a brand colour."""
    if accent:
        a = accent.strip()
        if a.lower() in ACCENT_MAP:
            return ACCENT_MAP[a.lower()]
        hexcand = a.lstrip("#")
        if len(hexcand) == 6:
            try:
                int(hexcand, 16)
                return hexcand.upper()
            except ValueError:
                pass
    return ACCENT_CYCLE[idx % len(ACCENT_CYCLE)]


FOOTER_TEXT = "© HCLTech  |  Confidential"


def _theme_for(archetype: str) -> dict:
    """Return a per-archetype colour/style theme.

    Dark slides (cover + closing) sit on navy; everything else is a white
    content slide that carries a full-width navy header bar at the top.
    """
    a = (archetype or "").strip().lower()
    if a == "title":
        return {
            "kind": "title",
            "bg": COLOR_PRIMARY,
            "title": COLOR_WHITE,
            "body": COLOR_BODY_LIGHT,
            "accent": COLOR_ACCENT,
        }
    if a == "next steps":
        return {
            "kind": "section",
            "bg": COLOR_PRIMARY,
            "title": COLOR_WHITE,
            "body": COLOR_BODY_LIGHT,
            "accent": COLOR_ACCENT,
        }
    if a == "agenda":
        return {
            "kind": "agenda",
            "bg": COLOR_WHITE,
            "title": COLOR_HEADER_TEXT,
            "body": COLOR_BODY,
            "accent": COLOR_ACCENT,
        }
    return {
        "kind": "content",
        "bg": COLOR_WHITE,
        "title": COLOR_HEADER_TEXT,
        "body": COLOR_BODY,
        "accent": COLOR_ACCENT,
    }


def _set_slide_background(slide, hex_color: str) -> None:
    """Fill the whole slide background with a solid colour."""
    try:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hex_color)
    except Exception:
        # Background fill is cosmetic; never fail the render over it.
        pass


def _set_run_font(run, name: str = FONT_NAME) -> None:
    """Force a run's typeface (latin + complex script) so it never falls back."""
    try:
        run.font.name = name
        rPr = run._r.get_or_add_rPr()  # pylint: disable=protected-access
        for tag in ("a:latin", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {})
                rPr.append(el)
            el.set("typeface", name)
    except Exception:
        pass


def _estimate_title_lines(title: str, width_in: float, font_pt: int) -> int:
    """Estimate how many lines a bold title wraps to inside a box of width_in."""
    title = (title or "").strip()
    if not title:
        return 1
    avg_char_pt = max(1.0, font_pt * 0.52)
    chars_per_line = max(1, int((width_in * 72.0) / avg_char_pt))
    return max(1, math.ceil(len(title) / chars_per_line))


def _fit_title_font(title: str, width_in: float, max_lines: int = 2,
                    start_pt: int = FONT_TITLE_PT, min_pt: int = FONT_TITLE_MIN_PT) -> int:
    """Shrink the title font (within bounds) until it fits within max_lines."""
    for pt in range(start_pt, min_pt - 1, -1):
        if _estimate_title_lines(title, width_in, pt) <= max_lines:
            return pt
    return min_pt


def _add_rect(slide, x_in, y_in, w_in, h_in, fill_hex, *, shape=MSO_SHAPE.RECTANGLE,
              line_hex: Optional[str] = None, line_w_pt: float = 0.0):
    """Add a filled (optionally outlined) auto-shape, shadow disabled."""
    try:
        shp = slide.shapes.add_shape(shape, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        if fill_hex is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
        if line_hex is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = RGBColor.from_string(line_hex)
            shp.line.width = Pt(line_w_pt or 1.0)
        shp.shadow.inherit = False
        return shp
    except Exception:
        return None


def _dot_grid(slide, x_in: float, y_in: float, cols: int, rows: int,
              dot_in: float = 0.07, gap_in: float = 0.45, color_hex: str = COLOR_TINT) -> None:
    """Draw a decorative grid of small squares (cover-slide motif)."""
    for r in range(rows):
        for c in range(cols):
            _add_rect(
                slide,
                x_in + c * gap_in,
                y_in + r * gap_in,
                dot_in,
                dot_in,
                color_hex,
            )


# ------------------------------------------------------------------
# Header bar + footer (content slides)
# ------------------------------------------------------------------
def _draw_header(slide, prs: Presentation, title: str, theme: dict) -> float:
    """Draw a clean content-slide header (navy title + short accent rule).

    Modern light style (per the reference deck): the title sits as bold navy
    text on the white slide, underlined by a short accent rule — no heavy
    navy bar. Returns the content-top (in).
    """
    w_in = float(prs.slide_width) / EMU_PER_INCH
    h_in = float(prs.slide_height) / EMU_PER_INCH

    pad = max(0.4, w_in * 0.04)
    title_w = w_in - 2 * pad

    title_pt = _fit_title_font(title, title_w, max_lines=2,
                               start_pt=FONT_TITLE_PT, min_pt=FONT_TITLE_MIN_PT)
    lines = _estimate_title_lines(title, title_w, title_pt)
    title_h = 0.66 if lines <= 1 else 1.05
    top = h_in * 0.05

    tb = slide.shapes.add_textbox(Inches(pad), Inches(top), Inches(title_w), Inches(title_h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(COLOR_NAVY)
    _set_run_font(run, FONT_NAME_HEAVY)
    p.alignment = PP_ALIGN.LEFT

    # Short accent rule under the title.
    rule_y = top + title_h + 0.04
    _add_rect(slide, pad, rule_y, min(title_w, 1.6), 0.055, theme["accent"])

    return rule_y + 0.055 + h_in * 0.04


def _draw_footer(slide, prs: Presentation, page_no: Optional[int], dark: bool = False) -> float:
    """Draw the footer rule + label; return the content-bottom (in)."""
    w_in = float(prs.slide_width) / EMU_PER_INCH
    h_in = float(prs.slide_height) / EMU_PER_INCH
    pad = max(0.4, w_in * 0.04)

    footer_y = h_in - max(0.34, h_in * 0.065)
    rule_color = COLOR_ACCENT if dark else "E2E7F0"
    text_color = COLOR_BODY_LIGHT if dark else COLOR_FOOTER

    _add_rect(slide, pad, footer_y, w_in - 2 * pad, 0.012, rule_color)

    tb = slide.shapes.add_textbox(
        Inches(pad), Inches(footer_y + 0.02), Inches(w_in - 2 * pad), Inches(0.3)
    )
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = FOOTER_TEXT
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor.from_string(text_color)
    _set_run_font(run)
    p.alignment = PP_ALIGN.LEFT

    if page_no is not None:
        # Page number in its own right-aligned box (a tab in the same frame is fragile).
        pb = slide.shapes.add_textbox(
            Inches(w_in - pad - 1.0), Inches(footer_y + 0.02), Inches(1.0), Inches(0.3)
        )
        ptf = pb.text_frame
        ptf.clear()
        pp = ptf.paragraphs[0]
        prun = pp.add_run()
        prun.text = str(page_no)
        prun.font.size = Pt(9)
        prun.font.color.rgb = RGBColor.from_string(text_color)
        _set_run_font(prun)
        pp.alignment = PP_ALIGN.RIGHT

    return footer_y - h_in * 0.02


# ------------------------------------------------------------------
# Content layout (within the header/footer band)
# ------------------------------------------------------------------
def _content_box(prs: Presentation, content_top: float, content_bottom: float) -> dict:
    """Full-width body box between the header and footer."""
    w_in = float(prs.slide_width) / EMU_PER_INCH
    pad = max(0.4, w_in * 0.04)
    return {
        "body": (pad, content_top, w_in - 2 * pad, max(0.6, content_bottom - content_top)),
    }


def _split_box(prs: Presentation, content_top: float, content_bottom: float) -> dict:
    """Two-column body: image LEFT, text RIGHT, between header and footer."""
    w_in = float(prs.slide_width) / EMU_PER_INCH
    pad = max(0.4, w_in * 0.04)
    inner_w = w_in - 2 * pad
    gutter = w_in * 0.03
    img_w = (inner_w - gutter) * 0.46
    body_w = (inner_w - gutter) * 0.54
    h = max(0.6, content_bottom - content_top)
    return {
        "image": (pad, content_top, img_w, h),
        "body": (pad + img_w + gutter, content_top, body_w, h),
    }


def _style_bullet_paragraph(p, color_hex: str, bullet_hex: str, level: int = 0) -> None:
    """Apply consulting-style bullet formatting: bullet glyph, colour, spacing."""
    p.font.color.rgb = RGBColor.from_string(color_hex)
    if level <= 0:
        mar_l, indent, glyph, spc_pts = 0.30, -0.22, "▪", "600"
    else:
        mar_l, indent, glyph, spc_pts = 0.62, -0.20, "–", "300"
    try:
        pPr = p._p.get_or_add_pPr()  # pylint: disable=protected-access
        pPr.set("marL", str(int(Inches(mar_l))))
        pPr.set("indent", str(int(Inches(indent))))

        for tag in ("a:spcBef", "a:buClr", "a:buFont", "a:buChar", "a:buNone", "a:buAutoNum"):
            for el in pPr.findall(qn(tag)):
                pPr.remove(el)

        def_rpr = pPr.find(qn("a:defRPr"))

        def _place(el):
            if def_rpr is not None:
                def_rpr.addprevious(el)
            else:
                pPr.append(el)

        spc_bef = pPr.makeelement(qn("a:spcBef"), {})
        spc_val = pPr.makeelement(qn("a:spcPts"), {"val": spc_pts})
        spc_bef.append(spc_val)
        _place(spc_bef)

        bu_clr = pPr.makeelement(qn("a:buClr"), {})
        srgb = pPr.makeelement(qn("a:srgbClr"), {"val": bullet_hex})
        bu_clr.append(srgb)
        _place(bu_clr)

        bu_font = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        _place(bu_font)

        bu_char = pPr.makeelement(qn("a:buChar"), {"char": glyph})
        _place(bu_char)
    except Exception:
        pass


def _fit_font_for_box(lines, w_in: float, h_in: float, min_pt: int = FONT_BODY_MIN_PT,
                      start_pt: int = FONT_BODY_START_PT) -> int:
    """Estimate wrapped line usage while preserving the readability floor.

    ``min_pt`` lets dense components (e.g. cards) shrink below the body floor
    rather than overflowing their box. The wrapping estimate is intentionally
    conservative so text does not spill past the placeholder in PowerPoint.
    """
    content = [str(text or "").strip() for _, text in (lines or []) if str(text or "").strip()]
    if not content:
        return start_pt
    min_pt = max(1, min(min_pt, start_pt))
    for size in range(start_pt, min_pt - 1, -1):
        average_char_width_pt = max(1.0, size * 0.55)
        chars_per_line = max(8, int((w_in * 72.0) / average_char_width_pt))
        wrapped_lines = sum(max(1, math.ceil(len(text) / chars_per_line)) for text in content)
        required_height = wrapped_lines * (size / 72.0) * 1.30
        if required_height <= h_in:
            return size
    return min_pt


def _estimated_required_height_in(lines, w_in: float, font_pt: int = FONT_BODY_MIN_PT) -> float:
    """Estimate text height with PowerPoint-like wrapping."""
    content = [str(text or "").strip() for _, text in (lines or []) if str(text or "").strip()]
    if not content:
        return 0.0
    average_char_width_pt = max(1.0, font_pt * 0.55)
    chars_per_line = max(8, int((w_in * 72.0) / average_char_width_pt))
    wrapped_lines = sum(max(1, math.ceil(len(text) / chars_per_line)) for text in content)
    return wrapped_lines * (font_pt / 72.0) * 1.30


def _text_fits_box(lines, w_in: float, h_in: float, margin: float = 0.92) -> bool:
    """Return whether text can fit at the minimum readable font size."""
    return _estimated_required_height_in(lines, w_in, FONT_BODY_MIN_PT) <= max(0.1, h_in * margin)


def _normalize_body_lines(bullets, detailed_points) -> list[tuple[int, str]]:
    """Flatten slide body into ``(level, text)`` lines."""
    lines: list[tuple[int, str]] = []
    if detailed_points:
        for point in detailed_points:
            text = (getattr(point, "text", "") or "").strip()
            if text:
                lines.append((0, text))
            for sub in getattr(point, "sub_points", []) or []:
                sub_txt = (sub or "").strip()
                if sub_txt:
                    lines.append((1, sub_txt))
        if lines:
            return lines
    for b in bullets or []:
        txt = (b or "").strip()
        if txt:
            lines.append((0, txt))
    return lines


def _add_body(
    slide,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    lines: list[tuple[int, str]],
    font_pt: int,
    color_hex: str = COLOR_BODY,
    bullet_hex: str = COLOR_ACCENT,
):
    """Render a list of ``(level, text)`` lines as styled, possibly nested bullets."""
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    for level, txt in lines:
        if not txt:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = txt
        _set_run_font(run)
        p.level = 1 if level >= 1 else 0
        run.font.size = Pt(font_pt if level <= 0 else max(FONT_BODY_MIN_PT, font_pt - 2))
        p.alignment = PP_ALIGN.LEFT
        _style_bullet_paragraph(p, color_hex=color_hex, bullet_hex=bullet_hex, level=level)


def _normalize_table(table) -> tuple[list[str], list[list[str]]]:
    if not table:
        return [], []
    headers = [str(h or "").strip() for h in (table.get("headers") or [])]
    rows = []
    for row in table.get("rows") or []:
        if isinstance(row, dict):
            vals = [str(row.get(h, "") or "").strip() for h in headers]
        else:
            vals = [str(v or "").strip() for v in (row or [])]
        rows.append(vals)
    max_cols = max([len(headers)] + [len(r) for r in rows] + [0])
    if max_cols <= 0:
        return [], []
    if not headers:
        headers = [f"Column {i + 1}" for i in range(max_cols)]
    headers = (headers + [""] * max_cols)[:max_cols]
    rows = [(r + [""] * max_cols)[:max_cols] for r in rows]
    return headers, rows


def _render_table_at(slide, headers: list[str], rows: list[list[str]], x_in: float, y_in: float,
                     w_in: float, h_in: float):
    if not headers:
        return None
    row_count = max(1, len(rows) + 1)
    col_count = max(1, len(headers))
    table_shape = slide.shapes.add_table(
        row_count, col_count, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    table = table_shape.table
    widths = [1.0 / col_count] * col_count
    if col_count == 7:
        widths = [0.16, 0.12, 0.19, 0.13, 0.12, 0.14, 0.14]
    elif col_count == 6:
        widths = [0.20, 0.15, 0.11, 0.20, 0.15, 0.19]
    elif col_count == 5:
        widths = [0.20, 0.16, 0.26, 0.20, 0.18]
    elif col_count == 4:
        widths = [0.24, 0.20, 0.34, 0.22]
    for idx, width in enumerate(widths):
        table.columns[idx].width = int(Inches(w_in * width))

    # Main-slide tables must remain readable. Large tables are paginated before
    # rendering, so five-column SBOM tables can stay at 9 pt.
    font_pt = 8 if len(rows) > 10 else 9
    header_h = 0.42
    body_h = min(0.58, max(0.30, (h_in - header_h) / max(1, len(rows))))
    table.rows[0].height = Inches(header_h)
    for row_idx in range(1, len(table.rows)):
        table.rows[row_idx].height = Inches(body_h)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(COLOR_NAVY)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(font_pt)
                run.font.bold = True
                run.font.color.rgb = RGBColor.from_string(COLOR_WHITE)
                _set_run_font(run, FONT_NAME_HEAVY)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF" if r % 2 else "F4F7FB")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(font_pt)
                    run.font.color.rgb = RGBColor.from_string(COLOR_BODY)
                    _set_run_font(run)
    return table_shape


def _render_table(slide, prs: Presentation, table_data, x_in: Optional[float] = None,
                  y_in: Optional[float] = None, w_in: Optional[float] = None,
                  h_in: Optional[float] = None):
    headers, rows = _normalize_table(table_data)
    if not headers:
        return None
    if x_in is None or y_in is None or w_in is None or h_in is None:
        x_in, y_in, w_in, h_in = _native_content_box(prs)
    slide_w = float(prs.slide_width) / EMU_PER_INCH
    slide_h = float(prs.slide_height) / EMU_PER_INCH
    w_in = min(w_in, max(1.0, slide_w - x_in - 0.7))
    h_in = min(h_in, max(1.0, slide_h - y_in - 0.75))
    return _render_table_at(slide, headers, rows, x_in, y_in, w_in, h_in)


def _place_image_contain(
    slide,
    img_source: Union[Path, bytes],
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    inset_in: float = 0.0,
):
    """Place an image contained within the bounding box (no overflow)."""
    from PIL import Image

    if inset_in > 0:
        x_in = x_in + inset_in
        y_in = y_in + inset_in
        w_in = max(0.01, w_in - (2 * inset_in))
        h_in = max(0.01, h_in - (2 * inset_in))

    if isinstance(img_source, bytes):
        img = Image.open(BytesIO(img_source))
    else:
        img = Image.open(img_source)
    iw, ih = img.size
    img.close()

    box_ratio = w_in / max(h_in, 1e-6)
    img_ratio = iw / max(ih, 1e-6)

    if img_ratio > box_ratio:
        new_w = w_in
        new_h = w_in / img_ratio
    else:
        new_h = h_in
        new_w = h_in * img_ratio

    px = x_in + (w_in - new_w) / 2.0
    py = y_in + (h_in - new_h) / 2.0

    if isinstance(img_source, bytes):
        stream = BytesIO(img_source)
        stream.seek(0)
        return slide.shapes.add_picture(
            stream, Inches(px), Inches(py), width=Inches(new_w), height=Inches(new_h)
        )
    return slide.shapes.add_picture(
        str(img_source), Inches(px), Inches(py), width=Inches(new_w), height=Inches(new_h)
    )


# ------------------------------------------------------------------
# Modern card-based components (key message, cards, KPI chips)
# ------------------------------------------------------------------
def _add_key_message(slide, x_in: float, y_in: float, w_in: float, text: str) -> float:
    """Navy rounded banner carrying a single emphasised 'so what' line.

    Returns the y-coordinate just below the banner.
    """
    text = (text or "").strip()
    if not text:
        return y_in
    h = 0.62
    _add_rect(slide, x_in, y_in, w_in, h, COLOR_NAVY_DK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = slide.shapes.add_textbox(Inches(x_in + 0.22), Inches(y_in), Inches(w_in - 0.44), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(COLOR_TEAL)
    _set_run_font(run, FONT_NAME_HEAVY)
    p.alignment = PP_ALIGN.LEFT
    return y_in + h + 0.18


def _add_card(slide, x_in: float, y_in: float, w_in: float, h_in: float,
              heading: str, body: str = "", bullets=None, accent_hex: str = COLOR_BLUE) -> None:
    """Render one content card: white panel, coloured left stripe, heading + body."""
    bullets = bullets or []
    # Panel + hairline border.
    _add_rect(slide, x_in, y_in, w_in, h_in, COLOR_CARD_BG,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_hex=COLOR_CARD_LINE, line_w_pt=0.75)
    # Coloured left accent stripe.
    _add_rect(slide, x_in, y_in, 0.11, h_in, accent_hex, shape=MSO_SHAPE.ROUNDED_RECTANGLE)

    pad_l = 0.30
    tb = slide.shapes.add_textbox(
        Inches(x_in + pad_l), Inches(y_in + 0.12),
        Inches(max(0.4, w_in - pad_l - 0.18)), Inches(max(0.3, h_in - 0.22)),
    )
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    hr = p.add_run()
    hr.text = (heading or "").strip()
    hr.font.size = Pt(13)
    hr.font.bold = True
    hr.font.color.rgb = RGBColor.from_string(accent_hex)
    _set_run_font(hr, FONT_NAME_HEAVY)
    p.alignment = PP_ALIGN.LEFT

    body = (body or "").strip()
    if body:
        bp = tf.add_paragraph()
        bp.space_before = Pt(4)
        br = bp.add_run()
        br.text = body
        br.font.size = Pt(11)
        br.font.color.rgb = RGBColor.from_string(COLOR_BODY)
        _set_run_font(br)
        bp.alignment = PP_ALIGN.LEFT

    for b in bullets:
        b = (b or "").strip()
        if not b:
            continue
        lp = tf.add_paragraph()
        lp.space_before = Pt(3)
        lr = lp.add_run()
        lr.text = b
        lr.font.size = Pt(11)
        lr.font.color.rgb = RGBColor.from_string(COLOR_BODY)
        _set_run_font(lr)
        lp.alignment = PP_ALIGN.LEFT
        _style_bullet_paragraph(lp, color_hex=COLOR_BODY, bullet_hex=accent_hex, level=0)


def _add_kpi_row(slide, x_in: float, y_in: float, w_in: float, items, accent_hex: str = COLOR_BLUE) -> None:
    """Render a horizontal row of KPI/stat chips."""
    items = [i.strip() for i in (items or []) if (i or "").strip()]
    if not items:
        return
    n = min(len(items), 4)
    items = items[:n]
    gap = 0.18
    chip_w = (w_in - gap * (n - 1)) / n
    chip_h = 0.52
    for i, text in enumerate(items):
        cx = x_in + i * (chip_w + gap)
        _add_rect(slide, cx, y_in, chip_w, chip_h, COLOR_CARD_TINT,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb = slide.shapes.add_textbox(Inches(cx + 0.08), Inches(y_in), Inches(chip_w - 0.16), Inches(chip_h))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(accent_hex)
        _set_run_font(run, FONT_NAME_HEAVY)
        p.alignment = PP_ALIGN.CENTER


def _add_icon_bullets(slide, x_in: float, y_in: float, w_in: float, h_in: float,
                      items, accent_hex: str = COLOR_BLUE) -> None:
    """Render a vertical list with circular accent markers beside each item."""
    items = [i.strip() for i in (items or []) if (i or "").strip()]
    if not items:
        return
    n = len(items)
    row_h = min(0.92, h_in / n)
    dot = min(0.42, row_h * 0.55)
    y = y_in
    for text in items:
        _add_rect(slide, x_in, y + (row_h - dot) / 2.0, dot, dot, accent_hex, shape=MSO_SHAPE.OVAL)
        tb = slide.shapes.add_textbox(
            Inches(x_in + dot + 0.22), Inches(y), Inches(w_in - dot - 0.22), Inches(row_h)
        )
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor.from_string(COLOR_BODY)
        _set_run_font(run)
        p.alignment = PP_ALIGN.LEFT
        y += row_h


def _set_speaker_notes(slide, notes: Optional[str]) -> None:
    """Write speaker notes onto the slide's notes page (best-effort)."""
    text = (notes or "").strip()
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


# ------------------------------------------------------------------
# Native HCLTech layout support
# ------------------------------------------------------------------
def _layout_name(layout) -> str:
    return (getattr(layout, "name", None) or "").strip()


def _is_instruction_layout(layout) -> bool:
    name = _layout_name(layout).lower()
    return (
        not name
        or "do not use" in name
        or name.startswith("instructions")
        or "instructions" in name
    )


def _is_hcltech_template(prs: Presentation) -> bool:
    names = [_layout_name(layout).lower() for layout in prs.slide_layouts]
    return len(names) > 50 and any("the beam" in name for name in names)


def _find_layout_by_name(prs: Presentation, name: Optional[str]):
    if not name:
        return None
    target = name.strip().lower()
    for layout in prs.slide_layouts:
        if _is_instruction_layout(layout):
            continue
        if _layout_name(layout).lower() == target:
            return layout
    return None


def _find_layout_with_tokens(prs: Presentation, *tokens: str):
    wanted = [t.strip().lower() for t in tokens if t and t.strip()]
    if not wanted:
        return None
    for layout in prs.slide_layouts:
        if _is_instruction_layout(layout):
            continue
        name = _layout_name(layout).lower()
        if all(token in name for token in wanted):
            return layout
    return None


def _find_text_layout_with_tokens(prs: Presentation, *tokens: str):
    """Find a token-matched layout that does not reserve space for imagery."""
    wanted = [token.strip().lower() for token in tokens if token and token.strip()]
    if not wanted:
        return None
    for layout in prs.slide_layouts:
        if _is_instruction_layout(layout):
            continue
        name = _layout_name(layout).lower()
        if "image" in name or "picture" in name:
            continue
        if not all(token in name for token in wanted):
            continue
        has_picture_placeholder = False
        for placeholder in getattr(layout, "placeholders", []):
            try:
                if str(placeholder.placeholder_format.type).upper().startswith("PICTURE"):
                    has_picture_placeholder = True
                    break
            except Exception:
                continue
        if not has_picture_placeholder:
            return layout
    return None


_COUNT_WORDS = {
    "one": 1,
    "two": 2,
    "three": 3,
    "four": 4,
    "five": 5,
    "six": 6,
    "seven": 7,
    "eight": 8,
    "nine": 9,
    "ten": 10,
}


def _layout_item_count(layout) -> Optional[int]:
    name = _layout_name(layout).lower()
    for word, count in _COUNT_WORDS.items():
        if re.search(rf"\b{word}\b", name):
            return count
    return None


# ------------------------------------------------------------------
# Layout variety rotation.
#
# The template ships many stylistically-different layouts that hold the same
# number of boxes (numbered / gradient-top / boxes-numbered / dark variants).
# Picking the first token match every time collapses a whole deck onto a
# handful of layouts (one real render put 73% of its content slides on a
# single "Two key points" skin). These families list only siblings verified
# (by inspecting their placeholder shapes) to hold the same box count with
# the same generic Number/Text-placeholder pattern the fill helpers already
# handle, so rotating between them is a pure style swap with no content risk.
# ------------------------------------------------------------------
_TWO_KEY_POINT_LAYOUTS = [
    "Two key points – Numbered, Boxes",
]
# "Two key points – Numbered sidebars" was in this pool briefly. Its accent
# color lives in decorative "Sidebar 1"/"Sidebar 2" shapes stacked vertically
# on the right edge of the slide, while the two content boxes it holds sit
# side by side on the left — there is no visual correspondence between which
# sidebar/number goes with which column, so a real two-column card/point
# slide renders as text on the left with two unrelated-looking colored blocks
# on the right. Confirmed on a real generated deck (title "Service measures
# connect operational performance to outcomes" and others) — the "Numbered,
# Boxes" variant keeps color and text in the same box and reads clearly, so
# it's the only variant in rotation for now. Re-add "Two key points – Numbered
# sidebars" only alongside a fill that actually places two side-by-side
# sidebars beside their matching column, not this one.
_THREE_KEY_POINT_LAYOUTS = [
    "Three key points – Gradient Top (Light)",
    "Three key points – Numbered boxes",
    "Three key points – Boxes, numbered",
    "Three key points – Gradient Top (Dark)",
    "Three key points – Boxes (Dark)",
    "Three key points – Gradient, boxes (Dark)",
]
_FOUR_KEY_POINT_LAYOUTS = [
    "Four key points – Gradient Top (Light)",
    "Four key points – Numbered boxes",
    "Four key points – Boxes, numbered",
    "Four key points – Numbered, Gradient, boxes (Light)",
    "Four key points – Gradient Top (Dark)",
    "Four key points – Boxes (Dark)",
    "Four key points – Gradient, boxes (Dark)",
]
# Diagram pages that carry a short caption (a key_message riding beside the
# picture, not the full companion explanation) get a native split layout
# instead of floating the caption over a full-bleed image. Both hold exactly
# one extra text box beside Title/Subtitle, mirrored left/right, so a single
# caption line always has a designed home. An *uncaptioned* image page (the
# common two-slide diagram pattern) is deliberately excluded from this list —
# see ``_find_safe_diagram_layout``.
_DIAGRAM_CAPTIONED_LAYOUTS = [
    "Diagram (space right, flexible) – Title, Subtitle",
    "Wide margins – Diagram (space left) – Title, Subtitle, Text",
]
# Alternated so a long paginated table doesn't run as one unbroken wall of
# identical-looking slides.
_TABLE_LAYOUTS = [
    "Table – Title",
    "Table w/t Sidebar",
]
# A 2-point slide that also carries numeric outcome highlights (kpis) gets a
# stat-forward layout instead of another plain two-box grid. Its own "Subhead
# N" / "Text Placeholder N" pair is filled by the existing generic card path
# unchanged; only the extra "Data N" / "Data description N" slots are new.
_INFOGRAPHIC_TWO_STAT_LAYOUT = "Wide margins – Infographics – Numbers (2), List, Two key points"


def _pick_varied_layout(prs: Presentation, usage: Optional[dict], names: list[str]):
    """Return the least-used layout among ``names`` that exists in ``prs``.

    ``usage`` is a per-render ``{layout_name: use_count}`` dict shared across
    one ``render_deck_from_template`` call. Ties fall back to ``names`` order,
    so with no usage history yet this reproduces the old "first match" pick.
    Pass ``usage=None`` (e.g. from a caller that doesn't track rotation) to
    always get the first available candidate.
    """
    by_name = {layout.name: layout for layout in prs.slide_layouts}
    candidates = [by_name[name] for name in names if name in by_name]
    if not candidates:
        return None
    if usage is None:
        return candidates[0]
    choice = min(candidates, key=lambda layout: (usage.get(layout.name, 0), names.index(layout.name)))
    usage[choice.name] = usage.get(choice.name, 0) + 1
    return choice


def _has_approved_diagram(slide_spec) -> bool:
    diagram = getattr(slide_spec, "diagram", None)
    return bool(diagram and getattr(diagram, "approved", False))


def _has_renderable_diagram(slide_spec, diagram_images: Optional[dict[str, bytes]] = None) -> bool:
    if not _has_approved_diagram(slide_spec):
        return False
    if diagram_images and _diagram_image_for_slide(slide_spec, diagram_images):
        return True
    diagram = getattr(slide_spec, "diagram", None)
    image_path = getattr(diagram, "image_path", None) if diagram else None
    return bool(image_path and Path(str(image_path)).exists())


def _diagram_image_for_slide(slide_spec, diagram_images: Optional[dict[str, bytes]] = None) -> Optional[bytes]:
    if not diagram_images:
        return None
    slide_id = getattr(slide_spec, "slide_id", "")
    if slide_id in diagram_images:
        return diagram_images[slide_id]
    diagram = getattr(slide_spec, "diagram", None)
    prompt = (getattr(diagram, "prompt", "") or "").strip() if diagram else ""
    if prompt and prompt in diagram_images:
        return diagram_images[prompt]
    return None


def _table_has_content(table) -> bool:
    """A table payload only counts when it carries real rows or headers.

    The planner sometimes emits ``table: {}`` on non-table slides; treating that
    empty dict as a table routes the slide to a table layout it cannot fill.
    """
    if not isinstance(table, dict):
        return bool(table)
    return bool(table.get("rows") or table.get("headers"))


def _is_visual_focused_slide(slide_spec) -> bool:
    if not _has_approved_diagram(slide_spec):
        return False
    if _table_has_content(getattr(slide_spec, "table", None)):
        return False
    if getattr(slide_spec, "comparison", None) is not None:
        return False
    cards = [
        c for c in (getattr(slide_spec, "cards", None) or [])
        if getattr(c, "heading", "").strip()
    ]
    if cards:
        return False
    return True


def _hint_is_compatible(layout, slide_spec) -> bool:
    name = _layout_name(layout).lower()
    archetype = (getattr(slide_spec, "archetype", "") or "").strip().lower()
    cards = [
        c for c in (getattr(slide_spec, "cards", None) or [])
        if getattr(c, "heading", "").strip()
    ]
    detailed_points = [
        p for p in (getattr(slide_spec, "detailed_points", None) or [])
        if (getattr(p, "text", "") or "").strip()
    ]

    if archetype == "agenda" and "agenda" not in name:
        return False
    if archetype == "title" and "cover" not in name:
        return False
    if archetype == "next steps":
        return False
    if "diagram" in name and not _has_approved_diagram(slide_spec):
        return False
    if "image" in name and archetype not in {"title"} and not _has_approved_diagram(slide_spec):
        return False
    if not _has_approved_diagram(slide_spec):
        for placeholder in getattr(layout, "placeholders", []):
            try:
                if str(placeholder.placeholder_format.type).upper().startswith("PICTURE"):
                    return False
            except Exception:
                continue
    has_table_payload = _table_has_content(getattr(slide_spec, "table", None))
    has_comparison = getattr(slide_spec, "comparison", None) is not None
    if has_table_payload and "table" not in name:
        return False
    if "table" in name and not has_table_payload and not has_comparison and archetype != "agenda":
        return False
    if _is_visual_focused_slide(slide_spec):
        return "diagram" in name and "space right, flexible" in name
    if "infographics" in name and "two key points" in name and len(_stat_shaped_kpis(slide_spec)) < 2:
        # Continuation-page hint propagation (_lock_continuation_layouts)
        # stamps whatever layout the first page of a split chose onto every
        # sibling page — a later page can easily have fewer/no kpis (they're
        # cleared for page_idx > 0) and would otherwise force this stats
        # layout with one or both number boxes never filled.
        return False
    if cards:
        declared = _layout_item_count(layout)
        if declared is not None and declared != len(cards):
            return False
    if detailed_points:
        declared = _layout_item_count(layout)
        if declared is not None and declared != len(detailed_points):
            return False
    return True


def _clean_kpis(slide_spec) -> list[str]:
    return [k for k in (getattr(slide_spec, "kpis", None) or []) if (k or "").strip()]


_STAT_TOKEN_RE = re.compile(
    r"\$?\d[\d,.]*\s?(?:%|percent|x\b|days?|weeks?|months?|hours?|hrs?)?",
    re.IGNORECASE,
)


def _stat_shaped_kpi_list(kpis) -> list[str]:
    """kpis that actually look like a short numeric stat chip.

    The ``kpis`` field is sometimes authored by the model as a full
    qualitative sentence (e.g. "Catalogue completeness measures to be
    agreed") rather than a stat ("40% faster onboarding") — routing those to
    a big-number box would either overflow a box sized for 2-3 characters or
    show an empty number with no content. Only route/fill when the text is
    short and actually contains a digit-led token.
    """
    return [
        kpi for kpi in (kpis or [])
        if (kpi or "").strip() and len(kpi) <= 60 and _STAT_TOKEN_RE.search(kpi)
    ]


def _stat_shaped_kpis(slide_spec) -> list[str]:
    """Callers routing to the 2-stat infographic layout require
    ``len(...) >= 2`` on this, not just a truthy check: a single incidental
    digit-led string (e.g. "24x7 visibility" among otherwise qualitative
    kpis like "MTTA"/"RCA closure") used to be enough to pass, leaving the
    layout's second number box permanently unfilled and stripped by
    _remove_unused_placeholders — a visibly blank third of the slide for
    content that was never real 2-stat material to begin with.
    """
    return _stat_shaped_kpi_list(getattr(slide_spec, "kpis", None))


def _has_explanatory_body(slide_spec) -> bool:
    return bool(
        getattr(slide_spec, "bullets", None)
        or getattr(slide_spec, "detailed_points", None)
        or (getattr(slide_spec, "key_message", None) or "").strip()
        or getattr(slide_spec, "kpis", None)
    )


def _find_safe_diagram_layout(prs: Presentation, slide_spec=None, usage: Optional[dict] = None):
    """Use the full slide body for generated visuals; explanation paginates separately.

    A slide that carries a short caption/key-message beside the picture gets a
    native split layout (image on one side, caption in a designed column). An
    uncaptioned image page — the common two-slide diagram pattern, where the
    explanation lives on a separate companion slide — stays full-width so a
    split layout never leaves half the slide visibly empty. ("Diagram (space
    center) – Full image" carries a decorative portrait in its master, so it
    is deliberately not used here.)
    """
    if slide_spec is not None and _has_explanatory_body(slide_spec):
        picked = _pick_varied_layout(prs, usage, _DIAGRAM_CAPTIONED_LAYOUTS)
        if picked is not None:
            return picked
    return (
        _find_layout_by_name(prs, "Title Only")
        or _find_blank_layout(prs)
    )


def _find_single_statement_layout(prs: Presentation):
    """A clean title + single body layout for a lone content block.

    A single item must never be forced into a multi-box "key points" layout
    (which leaves empty boxes and an orphan number) or a non-existent
    "one key points" layout (which would resolve to ``None``).
    """
    # Prefer a clean, full-width title+subtitle text layout with the title at
    # the top. Skip left/top-right/image/numbered variants that offset or crowd
    # the text, so a single statement reads as a centred, uncluttered slide.
    skip = (
        "image", "picture", "(left)", "top right", "chart", "diagram", "sidebar",
        "numbered", "key point", "cover", "divider", "agenda", "org", "contacts",
        "case stud", "premium", "quote", "infographic",
    )
    for layout in prs.slide_layouts:
        if _is_instruction_layout(layout):
            continue
        name = _layout_name(layout).lower()
        if "title" not in name or "subtitle" not in name:
            continue
        if any(token in name for token in skip):
            continue
        has_picture = False
        for placeholder in getattr(layout, "placeholders", []):
            try:
                if str(placeholder.placeholder_format.type).upper().startswith("PICTURE"):
                    has_picture = True
                    break
            except Exception:
                continue
        if not has_picture:
            return layout
    return (
        _find_text_layout_with_tokens(prs, "title", "subtitle")
        or _find_layout_with_tokens(prs, "wide margins", "text")
        or _find_blank_layout(prs)
    )


def _choose_hcltech_layout(prs: Presentation, slide_spec, usage: Optional[dict] = None):
    """Map generated deck archetypes to usable HCLTech POTX layouts."""
    # PlainText was an internal overflow escape hatch that rebuilt slides on
    # top of Title Only. HCLTech decks must remain on native content layouts.
    plain_text = False
    hinted = _find_layout_by_name(prs, getattr(slide_spec, "layout_hint", None))
    if not plain_text and hinted is not None and _hint_is_compatible(hinted, slide_spec):
        return hinted

    archetype = (getattr(slide_spec, "archetype", "") or "Content").strip().lower()
    cards = [c for c in (getattr(slide_spec, "cards", None) or []) if getattr(c, "heading", "").strip()]
    detailed_points = [
        p for p in (getattr(slide_spec, "detailed_points", None) or [])
        if (getattr(p, "text", "") or "").strip()
    ]
    has_comparison = getattr(slide_spec, "comparison", None) is not None
    has_table = _table_has_content(getattr(slide_spec, "table", None))
    has_diagram = _has_approved_diagram(slide_spec)

    if archetype == "team" and not has_diagram and (cards or detailed_points):
        # Route to the org-chart grid layout (populated by
        # _fill_org_chart_slide) ahead of the generic cards/detailed_points
        # branches below, which would otherwise always claim this slide for
        # an unrelated key-points family first — the archetype-specific
        # "team" branch further down is only ever reached once cards and
        # detailed_points are both empty, so it can't be the thing that
        # picks this layout when there's real role content to show.
        org_chart = _find_layout_with_tokens(prs, "org chart", "light bg") or _find_layout_with_tokens(
            prs, "org chart"
        )
        if org_chart is not None:
            return org_chart

    if _is_visual_focused_slide(slide_spec):
        return _find_safe_diagram_layout(prs, slide_spec, usage)

    if archetype == "divider":
        return _pick_varied_layout(prs, usage, ["Divider Beam – Light", "Divider Beam – Dark"])
    if archetype == "win theme":
        return (
            _find_layout_with_tokens(prs, "quote", "gradient", "light")
            or _find_layout_with_tokens(prs, "quote", "big image", "light")
            or _find_layout_with_tokens(prs, "quote")
        )
    if archetype == "title":
        return (
            _find_layout_with_tokens(prs, "cover", "the beam 3", "woman with device")
            or _find_layout_with_tokens(prs, "cover", "the beam", "intense")
            or _find_layout_with_tokens(prs, "cover")
        )
    if archetype == "next steps":
        item_count = min(6, max(3, len(getattr(slide_spec, "bullets", None) or [])))
        word = next((key for key, value in _COUNT_WORDS.items() if value == item_count), "five")
        return (
            _find_layout_with_tokens(
                prs, word, "key points", "numbered", "gradient", "boxes", "light"
            )
            or _find_layout_with_tokens(prs, word, "key points", "boxes", "numbered")
            or _find_layout_with_tokens(prs, word, "key points", "numbered boxes")
            or _find_layout_with_tokens(prs, word, "key points", "numbered")
            or _find_layout_with_tokens(prs, "five key points", "numbered boxes")
        )
    if archetype == "agenda":
        return (
            _find_layout_by_name(prs, "Agenda \u2013 Table")
            or _find_layout_with_tokens(prs, "agenda", "table")
            or _find_layout_with_tokens(prs, "agenda")
        )
    if has_table:
        return _find_layout_with_tokens(prs, "table", "title") or _find_layout_with_tokens(prs, "table")
    if has_comparison:
        return (
            _find_text_layout_with_tokens(prs, "two key points", "numbered", "boxes")
            or _find_text_layout_with_tokens(prs, "two key points")
            or _find_layout_with_tokens(prs, "two key points")
            or _find_layout_with_tokens(prs, "table", "sidebar")
        )
    if cards and not plain_text:
        # The Executive Summary stays on one slide, so give it the largest boxes
        # (numbered-boxes variant) — its cards then render at 14pt rather than
        # shrinking to fit the smaller gradient-top boxes. That's the right
        # shape for the common 2-4 roughly-even-card case, but it isn't the
        # only shape a win thesis can take — a single strong statement or a
        # 2-card stat-forward summary reads better on the same layouts the
        # generic cards path below already uses for those shapes, so check
        # those first rather than always forcing a fixed box count.
        if _is_exec_summary_spec(slide_spec):
            if len(cards) == 1:
                layout = _find_single_statement_layout(prs)
                if layout is not None:
                    return layout
            elif len(cards) == 2 and len(_stat_shaped_kpis(slide_spec)) >= 2:
                infographic = _find_layout_by_name(prs, _INFOGRAPHIC_TWO_STAT_LAYOUT)
                if infographic is not None:
                    return infographic
            count = min(4, max(2, len(cards)))
            word = next((w for w, v in _COUNT_WORDS.items() if v == count), "three")
            layout = (
                _find_layout_with_tokens(prs, word, "key points", "numbered boxes")
                or _find_layout_with_tokens(prs, word, "key points", "boxes", "numbered")
            )
            if layout is not None:
                return layout
        if len(cards) >= 4:
            return (
                _pick_varied_layout(prs, usage, _FOUR_KEY_POINT_LAYOUTS)
                or _find_layout_with_tokens(prs, "four key points", "gradient top", "light")
                or _find_layout_with_tokens(prs, "four key points", "numbered boxes")
            )
        if len(cards) == 3:
            return (
                _pick_varied_layout(prs, usage, _THREE_KEY_POINT_LAYOUTS)
                or _find_layout_with_tokens(prs, "three key points", "gradient top", "light")
                or _find_layout_with_tokens(prs, "three key points", "numbered boxes")
            )
        if len(cards) == 2:
            if len(_stat_shaped_kpis(slide_spec)) >= 2:
                infographic = _find_layout_by_name(prs, _INFOGRAPHIC_TWO_STAT_LAYOUT)
                if infographic is not None:
                    return infographic
            return (
                _pick_varied_layout(prs, usage, _TWO_KEY_POINT_LAYOUTS)
                or _find_text_layout_with_tokens(prs, "two key points", "numbered", "boxes")
                or _find_text_layout_with_tokens(prs, "two key points")
                or _find_layout_with_tokens(prs, "two key points")
            )
        if len(cards) == 1:
            return _find_single_statement_layout(prs)
    if detailed_points and not plain_text:
        if len(detailed_points) == 1:
            return _find_single_statement_layout(prs)
        count = min(5, len(detailed_points))
        word = next((word for word, value in _COUNT_WORDS.items() if value == count), "three")
        if count == 2 and len(_stat_shaped_kpis(slide_spec)) >= 2:
            infographic = _find_layout_by_name(prs, _INFOGRAPHIC_TWO_STAT_LAYOUT)
            if infographic is not None:
                return infographic
        family = {2: _TWO_KEY_POINT_LAYOUTS, 3: _THREE_KEY_POINT_LAYOUTS, 4: _FOUR_KEY_POINT_LAYOUTS}.get(count)
        if family:
            layout = _pick_varied_layout(prs, usage, family)
            if layout is not None:
                return layout
        if count in {3, 4}:
            layout = _find_layout_with_tokens(prs, word, "key points", "gradient top", "light")
            if layout is not None:
                return layout
        return (
            _find_layout_with_tokens(prs, word, "key points", "boxes", "numbered")
            or _find_layout_with_tokens(prs, word, "key points", "numbered boxes")
            or _find_layout_with_tokens(prs, word, "key points", "numbered")
            or _find_layout_with_tokens(prs, word, "key points")
        )
    flat_items = [
        item for item in (getattr(slide_spec, "bullets", None) or [])
        if (item or "").strip()
    ]
    if flat_items and not plain_text and archetype not in {"title", "agenda"}:
        if len(flat_items) == 1:
            return _find_single_statement_layout(prs)
        count = min(5, len(flat_items))
        word = next((word for word, value in _COUNT_WORDS.items() if value == count), "three")
        family = {2: _TWO_KEY_POINT_LAYOUTS, 3: _THREE_KEY_POINT_LAYOUTS, 4: _FOUR_KEY_POINT_LAYOUTS}.get(count)
        if family:
            layout = _pick_varied_layout(prs, usage, family)
            if layout is not None:
                return layout
        return (
            _find_text_layout_with_tokens(
                prs, word, "key points", "numbered", "gradient", "boxes", "light"
            )
            or _find_text_layout_with_tokens(prs, word, "key points", "numbered boxes")
            or _find_text_layout_with_tokens(prs, word, "key points", "numbered")
            or _find_text_layout_with_tokens(prs, word, "key points")
        )
    if archetype in {"timeline", "delivery plan"} and has_diagram:
        return _find_safe_diagram_layout(prs, slide_spec, usage)
    # A timeline/delivery-plan slide reaching here with no diagram already has
    # no cards/detailed_points/bullets either — those are checked generically
    # above this point and would have returned already if populated — so
    # forcing "four key points – numbered boxes" here (an earlier version of
    # this branch did) always produced boxes with nothing to put in them
    # (see slide-22-class content-loss cases). Falling through to the
    # generic catch-all below (same path architecture/solution overview use
    # without a diagram) beats forcing a layout guaranteed to render near-empty.
    if archetype == "deployment architecture" and has_diagram:
        return _find_safe_diagram_layout(prs, slide_spec, usage)
    if archetype == "high availability & dr":
        return (
            (_find_safe_diagram_layout(prs, slide_spec, usage) if has_diagram else None)
            or _find_layout_with_tokens(prs, "three key points", "numbered boxes")
        )
    if archetype == "software bill of materials":
        return _find_layout_with_tokens(prs, "table", "title") or _find_layout_with_tokens(prs, "table")
    if archetype == "team" and has_diagram:
        return _find_safe_diagram_layout(prs, slide_spec, usage)
    # A "team" slide with no diagram AND no cards/detailed_points (the
    # remaining case here — see the early org-chart check above for when
    # role content exists) has nothing to put on the org-chart grid; falling
    # through to the generic catch-all below (same path architecture/solution
    # overview use without a diagram) beats forcing a layout guaranteed to
    # render near-empty.
    if archetype == "case studies":
        return _find_layout_with_tokens(prs, "case studies", "intro") or _find_layout_with_tokens(prs, "case studies")
    if archetype in {"risks", "commercials"}:
        return _find_layout_with_tokens(prs, "table", "sidebar") or _find_layout_with_tokens(prs, "table")
    if archetype in {"architecture", "solution overview"} and has_diagram:
        return _find_safe_diagram_layout(prs, slide_spec, usage)
    if archetype == "content" or plain_text:
        return (
            _find_layout_with_tokens(prs, "wide margins", "text")
            or _find_layout_with_tokens(prs, "title", "text")
            or _find_layout_with_tokens(prs, "content")
            or _find_blank_layout(prs)
        )
    # Nothing structured matched. Pick by how much body content exists so a
    # near-empty slide (e.g. a diagram archetype whose image was not approved)
    # never lands on an over-sized multi-box grid — and never returns None.
    block_count = len(cards) + len(detailed_points) + len(flat_items)
    if block_count <= 1:
        return _find_single_statement_layout(prs)
    return (
        _find_layout_with_tokens(prs, "three key points", "numbered boxes")
        or _find_layout_with_tokens(prs, "four key points", "numbered boxes")
        or _find_single_statement_layout(prs)
    )


def _placeholder_type_name(shape) -> str:
    try:
        return str(shape.placeholder_format.type).upper()
    except Exception:
        return ""


def _placeholder_base_type(shape) -> str:
    return _placeholder_type_name(shape).split()[0]


def _placeholder_name(shape) -> str:
    return (getattr(shape, "name", "") or "").strip()


def _is_placeholder(shape, *type_tokens: str) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    base = _placeholder_base_type(shape)
    return any(token.upper() == base for token in type_tokens)


def _set_text(shape, text: str, *, font_pt: Optional[int] = None, bold: Optional[bool] = None) -> None:
    text = (text or "").strip()
    try:
        tf = shape.text_frame
    except Exception:
        try:
            shape.text = text
        except Exception:
            pass
        return

    tf.clear()
    tf.word_wrap = True
    paragraphs = text.splitlines() or [""]
    for idx, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        clean = line.strip()
        p.text = clean
        p.space_before = Pt(0)
        p.space_after = Pt(4)
        p.line_spacing = 1.05
        if clean.startswith("- "):
            p.text = clean[2:]
            p.level = 0
        elif clean.startswith("  - "):
            p.text = clean[4:]
            p.level = 1
        for run in p.runs:
            if font_pt is not None:
                run.font.size = Pt(font_pt)
            if bold is not None:
                run.font.bold = bold
            _set_run_font(run, FONT_NAME_HEAVY if bold else FONT_NAME)


def _fill_footer_placeholders(slide, page_no: Optional[int]) -> None:
    for shape in slide.placeholders:
        kind = _placeholder_type_name(shape)
        name = _placeholder_name(shape).lower()
        if "SLIDE_NUMBER" in kind and page_no is not None:
            _set_text(shape, str(page_no), font_pt=8)
        elif "FOOTER" in kind:
            # Preserve the template/master footer. Adding another text value
            # duplicates the HCLTech confidentiality line on native layouts.
            continue
        elif "DATE" in kind or "presenter or date" in name:
            _set_text(shape, "")


def _title_placeholders(slide) -> list:
    return [
        shape
        for shape in slide.placeholders
        if _is_placeholder(shape, "TITLE", "CENTER_TITLE")
    ]


def _body_placeholders(slide) -> list:
    skip_names = ("footer", "slide number", "date", "presenter")
    bodies = []
    for shape in slide.placeholders:
        name = _placeholder_name(shape).lower()
        kind = _placeholder_type_name(shape)
        base = _placeholder_base_type(shape)
        if any(skip in name for skip in skip_names):
            continue
        if any(skip in kind for skip in ("FOOTER", "SLIDE_NUMBER", "DATE", "PICTURE")):
            continue
        if base in {"TITLE", "CENTER_TITLE"}:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if base in {"BODY", "OBJECT", "CONTENT", "SUBTITLE"}:
            bodies.append(shape)
    return sorted(bodies, key=lambda s: (int(s.top), int(s.left)))


def _picture_placeholders(slide) -> list:
    return [
        shape
        for shape in slide.placeholders
        if _is_placeholder(shape, "PICTURE")
    ]


def _table_placeholders(slide) -> list:
    return [
        shape
        for shape in slide.placeholders
        if _is_placeholder(shape, "TABLE")
    ]


def _placeholder_box(slide, *type_tokens: str) -> Optional[tuple[float, float, float, float]]:
    """Return a placeholder box from the slide or, if needed, its layout.

    PowerPoint retains specialized placeholders such as TABLE on some slides,
    while python-pptx may expose them only on the layout in other templates.
    Using both sources keeps native HCLTech frame geometry authoritative.
    """
    candidates = [
        shape for shape in slide.placeholders
        if _is_placeholder(shape, *type_tokens)
    ]
    if not candidates:
        candidates = [
            shape for shape in slide.slide_layout.placeholders
            if _is_placeholder(shape, *type_tokens)
        ]
    if not candidates:
        return None
    target = max(candidates, key=lambda s: int(s.width) * int(s.height))
    return (
        float(target.left) / EMU_PER_INCH,
        float(target.top) / EMU_PER_INCH,
        float(target.width) / EMU_PER_INCH,
        float(target.height) / EMU_PER_INCH,
    )


def _set_native_title(slide, title: str, archetype: str) -> None:
    placeholders = _title_placeholders(slide)
    if not placeholders:
        return
    shape = placeholders[0]
    width_in = max(0.5, float(shape.width) / EMU_PER_INCH)
    height_in = max(0.4, float(shape.height) / EMU_PER_INCH)
    archetype_l = (archetype or "").strip().lower()

    if archetype_l == "agenda":
        title = "Agenda"
        start_pt, min_pt, max_lines = 44, 30, 2
    elif archetype_l == "title":
        start_pt, min_pt = 40, 28
        max_lines = max(2, min(4, int(height_in / 0.55)))
    else:
        start_pt = 24 if height_in < 0.75 else 28 if width_in < 5.5 else 32
        min_pt = 20
        max_lines = max(1, min(4, int(height_in / 0.52)))

    font_pt = _fit_title_font(
        title,
        width_in,
        max_lines=max_lines,
        start_pt=start_pt,
        min_pt=min_pt,
    )
    _set_text(shape, title, font_pt=font_pt, bold=None)


def _render_agenda_table(slide, items: list[str]) -> bool:
    box = _placeholder_box(slide, "TABLE")
    if box is None or not items:
        return False
    x, y, w, h = box
    items = [(item or "").strip() for item in items if (item or "").strip()][:9]
    if not items:
        return False

    # A contents list needs no header row ("Agenda item" reads like a tracker).
    shape = slide.shapes.add_table(
        len(items), 2, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = shape.table
    table.columns[0].width = Inches(w * 0.18)
    table.columns[1].width = Inches(w * 0.82)
    row_h = max(0.42, min(0.7, h / len(items)))
    for row_idx in range(len(table.rows)):
        table.rows[row_idx].height = Inches(row_h)

    for idx, item in enumerate(items, start=1):
        for col, value in enumerate((f"{idx:02d}", item)):
            cell = table.cell(idx - 1, col)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF" if idx % 2 else "F4F7FB")
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(14 if col else 13)
                    run.font.bold = col == 0
                    run.font.color.rgb = RGBColor.from_string(COLOR_BLUE if col == 0 else COLOR_BODY)
                    _set_run_font(run, FONT_NAME_HEAVY if col == 0 else FONT_NAME)
    return True


def _native_content_box(prs: Presentation) -> tuple[float, float, float, float]:
    w_in = float(prs.slide_width) / EMU_PER_INCH
    h_in = float(prs.slide_height) / EMU_PER_INCH
    return 0.75, 1.35, w_in - 1.5, h_in - 2.0


def _format_comparison_text(comparison) -> list[str]:
    if comparison is None:
        return []
    lines: list[str] = []
    for col in (getattr(comparison, "left", None), getattr(comparison, "right", None)):
        if col is None:
            continue
        heading = (getattr(col, "heading", "") or "").strip()
        if heading:
            lines.append(heading)
        for item in getattr(col, "items", []) or []:
            if (item or "").strip():
                lines.append(f"- {item.strip()}")
        lines.append("")
    return lines


def _comparison_to_cards(comparison) -> list:
    """Turn a two-column comparison into two card-like blocks for key-point boxes."""
    from types import SimpleNamespace

    cards = []
    sides = [
        (getattr(comparison, "left", None), "challenge"),
        (getattr(comparison, "right", None), "goal"),
    ]
    for col, default_accent in sides:
        if col is None:
            continue
        heading = (getattr(col, "heading", "") or "").strip()
        items = [str(item).strip() for item in (getattr(col, "items", None) or []) if str(item).strip()]
        if not heading and not items:
            continue
        cards.append(
            SimpleNamespace(
                heading=heading or "—",
                body="",
                bullets=items,
                accent=getattr(col, "accent", None) or default_accent,
            )
        )
    return cards


def _format_card_text(card) -> str:
    parts = [(getattr(card, "heading", "") or "").strip()]
    body = (getattr(card, "body", "") or "").strip()
    if body:
        parts.append(body)
    for bullet in getattr(card, "bullets", []) or []:
        if (bullet or "").strip():
            parts.append(f"- {bullet.strip()}")
    return "\n".join(p for p in parts if p)


def _native_slide_lines(bullets, detailed_points=None, key_message=None, comparison=None) -> list[str]:
    lines = []
    if key_message and (key_message or "").strip():
        lines.append(key_message.strip())
        lines.append("")
    if comparison is not None:
        lines.extend(_format_comparison_text(comparison))
    else:
        for level, text in _normalize_body_lines(bullets, detailed_points):
            if not text:
                continue
            prefix = "  - " if level else "- "
            lines.append(f"{prefix}{text}")
    return [line for line in lines if line is not None]


def _fill_numbered_or_keypoint_placeholders(slide, items: list[str]) -> int:
    if not items:
        return 0
    text_placeholders = []
    number_placeholders = []
    for shape in _body_placeholders(slide):
        name = _placeholder_name(shape).lower()
        if "number" in name:
            number_placeholders.append(shape)
        elif any(token in name for token in ("key point", "text placeholder", "subhead")):
            text_placeholders.append(shape)

    _clear_point_placeholders(slide, text_placeholders)
    pairs = list(zip(text_placeholders, items))
    uniform = _uniform_card_font(pairs)
    count = 0
    for shape, item in pairs:
        _set_text(shape, item, font_pt=uniform)
        count += 1
    for idx, shape in enumerate(number_placeholders[:count], start=1):
        _set_text(shape, f"{idx:02d}", font_pt=18, bold=True)
    for shape in number_placeholders[count:]:
        _set_text(shape, "")
    return count


def _placeholder_area_in(shape) -> float:
    return (float(shape.width) / EMU_PER_INCH) * (float(shape.height) / EMU_PER_INCH)


def _key_message_placeholder(slide):
    candidates = [
        shape for shape in _body_placeholders(slide)
        if shape.placeholder_format.idx == 12
    ]
    return candidates[0] if candidates else None


def _semantic_placeholder_name(slide, shape) -> str:
    try:
        idx = shape.placeholder_format.idx
        for layout_shape in slide.slide_layout.placeholders:
            if layout_shape.placeholder_format.idx == idx:
                return _placeholder_name(layout_shape).lower()
    except Exception:
        pass
    return _placeholder_name(shape).lower()


def _card_placeholders(slide) -> list:
    candidates = []
    for shape in _body_placeholders(slide):
        name = _semantic_placeholder_name(slide, shape)
        if shape.placeholder_format.idx <= 12:
            continue
        if "number" in name or "subhead" in name or "sidebar" in name:
            continue
        width_in = float(shape.width) / EMU_PER_INCH
        height_in = float(shape.height) / EMU_PER_INCH
        if width_in < 1.5 or height_in < 1.0:
            continue
        candidates.append(shape)
    return sorted(
        candidates,
        key=lambda shape: (
            round(float(shape.top) / EMU_PER_INCH, 1),
            int(shape.left),
        ),
    )


def _number_placeholders(slide) -> list:
    return sorted(
        [
            shape for shape in _body_placeholders(slide)
            if "number" in _semantic_placeholder_name(slide, shape)
        ],
        key=lambda shape: (
            round(float(shape.top) / EMU_PER_INCH, 1),
            int(shape.left),
        ),
    )


def _clear_point_placeholders(slide, placeholders: Optional[list] = None) -> None:
    for shape in placeholders or _card_placeholders(slide):
        _set_text(shape, "")
    for shape in _number_placeholders(slide):
        _set_text(shape, "")


def _card_fit_pt(shape, text: str) -> int:
    """Largest readable font (<=13, >=9) at which ``text`` fits this card box."""
    # Subtract the placeholder's internal text inset so a dense card shrinks to
    # fit instead of spilling past the box (and off the slide).
    width_in = max(0.5, float(shape.width) / EMU_PER_INCH - 0.30)
    height_in = max(0.5, float(shape.height) / EMU_PER_INCH - 0.25)
    lines = [(0, line) for line in (text or "").splitlines() if line.strip()]
    # Prefer 14pt where the box can hold it; the uniform-per-slide pass still
    # steps every card down together when the densest box needs less.
    return min(14, _fit_font_for_box(lines, width_in, height_in, min_pt=9))


def _set_card_text(shape, text: str, font_pt: Optional[int] = None) -> None:
    if font_pt is None:
        font_pt = _card_fit_pt(shape, text)
    _set_text(shape, text, font_pt=font_pt)
    try:
        first = shape.text_frame.paragraphs[0]
        for run in first.runs:
            run.font.bold = True
            _set_run_font(run, FONT_NAME_HEAVY)
        for paragraph in shape.text_frame.paragraphs[1:]:
            value = (paragraph.text or "").strip()
            if not value.startswith("-"):
                continue
            # A literal marker remains visible in native placeholders whose
            # inherited paragraph XML suppresses PowerPoint bullet properties.
            paragraph.text = "• " + value[1:].strip()
            for run in paragraph.runs:
                _set_run_font(run)
                run.font.size = Pt(font_pt)
    except Exception:
        pass


def _uniform_card_font(pairs: list[tuple]) -> Optional[int]:
    """One font size for every card on a slide (the smallest that fits all).

    Cards on the same slide must read as one system; letting each box auto-fit
    independently produces a distracting 10/12/13pt mix within a single slide.
    """
    sizes = [_card_fit_pt(shape, text) for shape, text in pairs if (text or "").strip()]
    return min(sizes) if sizes else None


def _fill_card_placeholders(slide, cards, key_message: Optional[str] = None) -> int:
    if key_message:
        key_ph = _key_message_placeholder(slide)
        if key_ph is not None:
            _set_text(key_ph, key_message, font_pt=11, bold=True)

    placeholders = _card_placeholders(slide)
    _clear_point_placeholders(slide, placeholders)
    pairs = [(shape, _format_card_text(card)) for shape, card in zip(placeholders, cards)]
    uniform = _uniform_card_font(pairs)
    for shape, text in pairs:
        _set_card_text(shape, text, font_pt=uniform)
    return len(pairs)


def _format_detailed_point(point) -> str:
    parts = [(getattr(point, "text", "") or "").strip()]
    for sub_point in getattr(point, "sub_points", []) or []:
        if (sub_point or "").strip():
            parts.append(f"- {sub_point.strip()}")
    return "\n".join(part for part in parts if part)


def _fill_detailed_point_placeholders(
    slide,
    detailed_points,
    key_message: Optional[str] = None,
) -> int:
    if key_message:
        key_ph = _key_message_placeholder(slide)
        if key_ph is not None:
            _set_text(key_ph, key_message, font_pt=11, bold=True)

    placeholders = _card_placeholders(slide)
    _clear_point_placeholders(slide, placeholders)
    number_shapes = _number_placeholders(slide)
    for idx, (number_shape, text_shape) in enumerate(
        zip(number_shapes, placeholders), start=1
    ):
        # Some native layouts place the number marker inside the content box.
        # Leaving that marker blank avoids collisions with pre-read copy.
        if _overlap_ratio(number_shape, text_shape) < 0.02:
            _set_text(number_shape, f"{idx:02d}", font_pt=16, bold=True)

    # Fill every box at one uniform font (the smallest that fits all points)
    # rather than bailing to a single-box dump or mixing sizes across the slide.
    pairs = [(shape, _format_detailed_point(point)) for shape, point in zip(placeholders, detailed_points)]
    uniform = _uniform_card_font(pairs)
    for shape, text in pairs:
        _set_card_text(shape, text, font_pt=uniform)
    return len(pairs)


def _is_org_chart_layout(slide) -> bool:
    return "org chart" in _layout_name(slide.slide_layout).lower()


def _org_chart_role_boxes(slide) -> tuple[Optional[Any], list]:
    """Return (lead_box, grid_boxes) for the org-chart layout's role placeholders.

    Detected purely from geometry (the topmost box, if it sits alone above
    everything else, is the lead/coordinator slot; the rest form the grid)
    rather than hardcoded placeholder indices, so this keeps working if the
    template's box count or arrangement ever changes.
    """
    boxes = []
    for shape in _body_placeholders(slide):
        name = _semantic_placeholder_name(slide, shape)
        if "text placeholder" not in name:
            continue
        boxes.append((round(float(shape.top) / EMU_PER_INCH, 1), float(shape.left), shape))
    if not boxes:
        return None, []
    boxes.sort(key=lambda b: (b[0], b[1]))
    top_row = boxes[0][0]
    lead_candidates = [b for b in boxes if b[0] == top_row]
    rest = [b for b in boxes if b[0] != top_row]
    if len(lead_candidates) == 1 and rest:
        return lead_candidates[0][2], [b[2] for b in rest]
    return None, [b[2] for b in boxes]


def _set_org_box_text(shape, lines: list[str], font_pt: Optional[int] = None) -> None:
    """Fill an org-chart role box's outline levels (0/1/2) with ``lines`` in
    order, leaving any unused trailing level blank. Proposals rarely have
    named individuals to place pre-award, so level 0 (the box's bold headline
    slot, prompted "Full name" in the template) carries the role title, not a
    fabricated person's name.

    A freshly added slide placeholder starts with a single empty paragraph —
    the layout's own 3-line "Full name / Role / Short description" preview
    text is a layout-authoring aid, not something ``add_slide`` copies onto
    the new instance — so paragraphs beyond the first have to be created
    here, at the matching outline level, for their list-style formatting
    (size/weight per level) to resolve the same way the layout's preview did.
    ``font_pt`` lets the caller apply one uniform size across every box on
    the slide (see ``_fill_org_chart_slide``) instead of each box fitting
    independently, which the rest of this renderer avoids for the same
    "distracting mixed sizes on one slide" reason (see ``_uniform_card_font``).
    """
    tf = shape.text_frame
    lines = [(line or "").strip() for line in lines]
    if font_pt is None:
        font_pt = _card_fit_pt(shape, "\n".join(line for line in lines if line))
    for level, text in enumerate(lines):
        paragraph = tf.paragraphs[level] if level < len(tf.paragraphs) else tf.add_paragraph()
        paragraph.level = level
        for run in list(paragraph.runs)[1:]:
            run._r.getparent().remove(run._r)  # pylint: disable=protected-access
        if paragraph.runs:
            paragraph.runs[0].text = text
        elif text:
            paragraph.add_run().text = text
        for run in paragraph.runs:
            run.font.size = Pt(font_pt)
            run.font.bold = level == 0
            _set_run_font(run, FONT_NAME_HEAVY if level == 0 else FONT_NAME)


def _org_box_description_cap(text: str, max_len: int = 65) -> str:
    """A role description short enough to read as a finished thought in its
    box, or nothing at all.

    The org-chart grid can't paginate the way key-point layouts do (it's a
    fixed set of boxes, not a list the renderer can split across more
    slides), and this box is genuinely small — its own template prompt calls
    it a "Short description". Cutting an arbitrary sentence down to size
    reliably risks a dangling fragment ("...minimises", "...and continuous")
    even with word/clause-boundary cuts, which reads worse than showing no
    description at all — so a description that doesn't already fit is
    dropped rather than truncated; the role title alone still carries the box.
    """
    text = re.sub(r"\s+", " ", (text or "").strip())
    return text if len(text) <= max_len else ""


def _org_chart_role_lines(source) -> tuple[str, str]:
    """(role title, short description) from a card or a detailed_point."""
    heading = (getattr(source, "heading", "") or getattr(source, "text", "") or "").strip()
    body = (getattr(source, "body", "") or "").strip()
    if not body:
        sub_points = getattr(source, "sub_points", None) or []
        body = next((s.strip() for s in sub_points if (s or "").strip()), "")
    if not body:
        bullets = getattr(source, "bullets", None) or []
        body = next((b.strip() for b in bullets if (b or "").strip()), "")
    return heading, _org_box_description_cap(body)


def _fill_org_chart_slide(slide, cards, detailed_points) -> int:
    """Populate the native org-chart layout's role grid from whatever generic
    role content the slide already carries (cards or detailed_points) — no
    fabricated names, no engagement-specific assumptions. Returns the number
    of boxes filled so the caller can fall back when there's nothing to show.
    """
    roles = list(cards) if cards else list(detailed_points)
    if not roles:
        return 0
    lead_box, grid_boxes = _org_chart_role_boxes(slide)
    remaining = list(roles)
    assignments = []
    if lead_box is not None and remaining:
        title, desc = _org_chart_role_lines(remaining.pop(0))
        assignments.append((lead_box, [title, "", desc]))
    for shape, role in zip(grid_boxes, remaining):
        title, desc = _org_chart_role_lines(role)
        assignments.append((shape, [title, desc]))
    if not assignments:
        return 0
    font_pt = min(
        _card_fit_pt(shape, "\n".join(line for line in lines if line))
        for shape, lines in assignments
    )
    for shape, lines in assignments:
        _set_org_box_text(shape, lines, font_pt=font_pt)
    return len(assignments)


def _fill_flat_point_placeholders(
    slide,
    items: list[str],
    key_message: Optional[str] = None,
) -> int:
    """Distribute flat bullets across a native multi-point layout."""
    placeholders = _card_placeholders(slide)
    if len(placeholders) < 2:
        return 0
    clean_items = [(item or "").strip() for item in items if (item or "").strip()]
    if not clean_items:
        return 0

    if key_message:
        key_ph = _key_message_placeholder(slide)
        if key_ph is not None:
            _set_text(key_ph, key_message, font_pt=11, bold=True)

    _clear_point_placeholders(slide, placeholders)
    number_shapes = _number_placeholders(slide)

    groups = [[] for _ in placeholders]
    for idx, item in enumerate(clean_items):
        group_idx = (
            idx
            if len(clean_items) <= len(placeholders)
            else min(len(placeholders) - 1, idx * len(placeholders) // len(clean_items))
        )
        groups[group_idx].append(item)
    group_pairs = [(shape, "\n".join(group)) for shape, group in zip(placeholders, groups) if group]
    uniform = _uniform_card_font(group_pairs)
    for shape, text in group_pairs:
        _set_text(shape, text, font_pt=uniform, bold=None)
    filled_count = sum(bool(group) for group in groups)
    for idx, shape in enumerate(number_shapes[:filled_count], start=1):
        _set_text(shape, f"{idx:02d}", font_pt=16, bold=True)
    for shape in number_shapes[filled_count:]:
        _set_text(shape, "")
    return filled_count


def _fill_diagram_side_text(
    slide,
    inserted_diagram,
    lines: list[str],
    key_message: Optional[str],
) -> None:
    """Populate split-layout text without overflowing the diagram-side column."""
    candidates = [
        shape for shape in _body_placeholders(slide)
        if _overlap_ratio(shape, inserted_diagram) < 0.05
    ]
    if not candidates:
        return

    subtitle = next(
        (
            shape for shape in candidates
            if "subtitle" in _semantic_placeholder_name(slide, shape)
            or shape.placeholder_format.idx == 12
        ),
        None,
    )
    body_candidates = [shape for shape in candidates if shape is not subtitle]
    body = max(body_candidates, key=_placeholder_area_in) if body_candidates else subtitle

    clean_lines = [
        re.sub(r"\s+", " ", (line or "").strip())
        for line in lines
        if (line or "").strip()
    ]
    subtitle_text = (key_message or "").strip()
    if subtitle_text:
        clean_lines = [line for line in clean_lines if line.lower() != subtitle_text.lower()]
    elif clean_lines:
        subtitle_text = clean_lines.pop(0)

    if subtitle is not None and subtitle_text:
        _set_text(subtitle, subtitle_text, font_pt=12, bold=True)

    if body is not None and body is not subtitle and clean_lines:
        body_text = "\n".join(clean_lines)
        box_w = float(body.width) / EMU_PER_INCH
        box_h = float(body.height) / EMU_PER_INCH
        fit_lines = [(0, line.replace("- ", "").strip()) for line in clean_lines]
        font_pt = _fit_font_for_box(fit_lines, box_w, box_h)
        _set_text(body, body_text, font_pt=font_pt, bold=None)


def _fill_native_body(slide, prs: Presentation, lines: list[str]) -> None:
    body = "\n".join(line for line in lines if line is not None).strip()
    if not body:
        return

    body_lines = [(0, line[2:].strip()) if line.startswith("- ") else (0, line.strip()) for line in lines]
    placeholders = _body_placeholders(slide)
    if placeholders:
        largest = max(placeholders, key=lambda s: int(s.width) * int(s.height))
        box_w = max(0.5, float(largest.width) / EMU_PER_INCH)
        box_h = max(0.5, float(largest.height) / EMU_PER_INCH)
        if _text_fits_box(body_lines, box_w, box_h, margin=0.96):
            _set_text(largest, body, font_pt=_fit_font_for_box(body_lines, box_w, box_h))
            return
        _set_text(largest, "")

    # Free-drawn fallback: never start above the title (some native layouts
    # place the title mid-slide, so a fixed top would print body over it).
    x, y, w, h = _native_content_box(prs)
    titles = _title_placeholders(slide)
    if titles:
        title_bottom = max(float(t.top + t.height) / EMU_PER_INCH for t in titles)
        if title_bottom + 0.2 > y:
            new_y = title_bottom + 0.3
            h = max(0.6, h - (new_y - y))
            y = new_y
    font_pt = _fit_font_for_box(body_lines, w, h)
    _add_body(slide, x, y, w, h, body_lines, font_pt, color_hex=COLOR_BODY, bullet_hex=COLOR_BLUE)


def _diagram_target_box(slide, prs: Presentation) -> tuple[float, float, float, float]:
    layout_name = _layout_name(slide.slide_layout).lower()
    w_in = float(prs.slide_width) / EMU_PER_INCH
    h_in = float(prs.slide_height) / EMU_PER_INCH
    # Default to a large visual-first canvas. Text explanation is rendered on
    # the following slide, so the image should be readable in presentation mode.
    title_bottom = 0.0
    titles = _title_placeholders(slide)
    if titles:
        title_bottom = max(
            float(shape.top + shape.height) / EMU_PER_INCH
            for shape in titles
        )
    image_top = max(1.45, title_bottom + 0.16)
    default_box = (0.45, image_top, w_in - 0.9, max(1.0, h_in - image_top - 0.55))

    if "space right" in layout_name:
        bodies = _body_placeholders(slide)
        if bodies:
            target = max(bodies, key=lambda shape: int(shape.width) * int(shape.height))
            return tuple(
                float(value) / EMU_PER_INCH
                for value in (target.left, target.top, target.width, target.height)
            )

    if "space left" in layout_name:
        title_left = min(
            (float(shape.left) / EMU_PER_INCH for shape in _title_placeholders(slide)),
            default=w_in * 0.56,
        )
        return (0.35, 0.35, max(4.5, title_left - 0.7), h_in - 0.8)

    return default_box


def _insert_native_diagram(slide, prs: Presentation, diagram, diagram_bytes: Optional[bytes]):
    if not diagram or not getattr(diagram, "approved", False):
        return None
    image_path = getattr(diagram, "image_path", None)
    image_source = diagram_bytes
    if image_source is None and image_path:
        path = Path(str(image_path))
        if path.exists():
            image_source = path
    if image_source is None:
        return None

    pictures = _picture_placeholders(slide)
    if pictures:
        try:
            source = BytesIO(image_source) if isinstance(image_source, bytes) else str(image_source)
            return pictures[0].insert_picture(source)
        except Exception:
            log.debug("Could not insert image into picture placeholder; using contain placement.", exc_info=True)

    x, y, w, h = _diagram_target_box(slide, prs)
    return _place_image_contain(slide, image_source, x, y, w, h, inset_in=0.08)


def _overlap_ratio(shape_a, shape_b) -> float:
    left = max(int(shape_a.left), int(shape_b.left))
    top = max(int(shape_a.top), int(shape_b.top))
    right = min(int(shape_a.left + shape_a.width), int(shape_b.left + shape_b.width))
    bottom = min(int(shape_a.top + shape_a.height), int(shape_b.top + shape_b.height))
    if right <= left or bottom <= top:
        return 0.0
    overlap = (right - left) * (bottom - top)
    area = max(1, int(shape_a.width) * int(shape_a.height))
    return overlap / area


def _is_generated_diagram_slide(slide_spec) -> bool:
    return bool(
        getattr(slide_spec, "diagram", None)
        and getattr(getattr(slide_spec, "diagram", None), "approved", False)
    )


def _remove_overlapping_generated_pictures(slide, slide_spec=None) -> int:
    if _is_generated_diagram_slide(slide_spec):
        return 0
    text_shapes = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (getattr(shape, "text", "") or "").strip()
        if not text:
            continue
        name = (getattr(shape, "name", "") or "").lower()
        if "footer" in name or "slide number" in name:
            continue
        text_shapes.append(shape)

    removed = 0
    pictures = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not shape.is_placeholder
    ]
    for picture in pictures:
        conflicts = [
            shape for shape in text_shapes
            if _overlap_ratio(shape, picture) > 0.15
        ]
        if not conflicts:
            continue
        element = picture._element  # pylint: disable=protected-access
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)
            removed += 1
            log.warning(
                "Removed generated picture from slide %r because it overlapped %d text shape(s).",
                getattr(slide, "name", ""),
                len(conflicts),
            )
    return removed


def _placeholder_has_inserted_content(shape) -> bool:
    """Return True when a placeholder contains real generated content."""
    if getattr(shape, "has_text_frame", False):
        if (getattr(shape, "text", "") or "").strip():
            return True
    if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
        return True
    try:
        if shape.image is not None:
            return True
    except (AttributeError, ValueError):
        pass
    try:
        if shape._element.xpath(".//a:blip"):  # pylint: disable=protected-access
            return True
    except Exception:
        pass
    return False


# Placeholder name fragments that mark a shape as pure decoration — it carries
# a layout's color/border identity but is never meant to hold text (the fill
# helpers already skip these by the same name check; see _card_placeholders).
# _remove_unused_placeholders must never treat "empty" here as "unused", or it
# strips a layout's entire visual identity (e.g. "Two key points – Numbered
# sidebars" loses its color panels because "Sidebar 1"/"Sidebar 2" hold no text).
_DECORATIVE_PLACEHOLDER_TOKENS = ("sidebar",)


def _is_decorative_placeholder(slide, shape) -> bool:
    name = _semantic_placeholder_name(slide, shape)
    return any(token in name for token in _DECORATIVE_PLACEHOLDER_TOKENS)


def _remove_unused_placeholders(slide) -> int:
    """Delete empty placeholder objects so they do not appear in edit mode."""
    removed = 0
    for shape in list(slide.placeholders):
        if _placeholder_has_inserted_content(shape):
            continue
        if _is_decorative_placeholder(slide, shape):
            continue
        element = shape._element  # pylint: disable=protected-access
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)
            removed += 1
    return removed


def _clear_slide_shapes(slide, *, preserve_title: bool = False) -> int:
    """Remove layout-created shapes before custom rendering."""
    removed = 0
    for shape in list(slide.shapes):
        if preserve_title and shape in _title_placeholders(slide):
            continue
        element = shape._element  # pylint: disable=protected-access
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)
            removed += 1
    return removed


def _xml_local_name(tag: Any) -> str:
    return str(tag).rsplit("}", 1)[-1]


def _remove_presentation_sections(prs: Presentation) -> int:
    """Remove section names inherited from the corporate template package."""
    root = prs.part._element  # pylint: disable=protected-access
    section_lists = [
        element for element in root.iter()
        if _xml_local_name(element.tag) == "sectionLst"
    ]
    for section_list in section_lists:
        parent = section_list.getparent()
        if parent is None:
            continue
        parent.remove(section_list)
        if _xml_local_name(parent.tag) == "ext" and len(parent) == 0:
            grandparent = parent.getparent()
            if grandparent is not None:
                grandparent.remove(parent)
                if _xml_local_name(grandparent.tag) == "extLst" and len(grandparent) == 0:
                    great_grandparent = grandparent.getparent()
                    if great_grandparent is not None:
                        great_grandparent.remove(grandparent)
    if section_lists:
        log.info("Removed %d inherited PowerPoint section list(s).", len(section_lists))
    return len(section_lists)


def _clear_orphan_number_placeholders(slide) -> int:
    """Remove visible number markers that have no corresponding visible point.

    Numbers are assigned by rank at fill time (``zip(number_shapes,
    placeholders)`` in ``_fill_detailed_point_placeholders`` etc., both sorted
    by the same ``(top, left)`` key) — so a number is legitimately "orphaned"
    only when its rank falls beyond how many content boxes actually got
    filled, e.g. a 4-box layout used for a 2-point slide. Checking absolute
    vertical position instead (an earlier version of this function) assumes a
    number sits above/beside its own content at increasing top values, which
    holds for stacked numbered-boxes layouts but not for layouts like "Two
    key points – Numbered sidebars", where two numbers stack in their own
    column (top ~2.3in and ~4.7in) beside content that both sit at one shared
    top (~3.7in) — that mismatch made the second, legitimately-filled number
    look orphaned and blank it.
    """
    numbers = _number_placeholders(slide)
    placeholders = _card_placeholders(slide)
    filled_count = sum(
        1 for shape in placeholders if (getattr(shape, "text", "") or "").strip()
    )
    removed = 0
    for rank, shape in enumerate(numbers):
        if not (getattr(shape, "text", "") or "").strip():
            continue
        if rank < filled_count:
            continue
        _set_text(shape, "")
        removed += 1
    return removed


def _fallback_visible_lines(slide_spec) -> list[str]:
    lines: list[str] = []
    key_message = (getattr(slide_spec, "key_message", None) or "").strip()
    if key_message:
        lines.append(key_message)
    for line in getattr(slide_spec, "bullets", None) or []:
        if (line or "").strip():
            lines.append(str(line).strip())
    for point in getattr(slide_spec, "detailed_points", None) or []:
        text = (getattr(point, "text", "") or "").strip()
        if text:
            lines.append(text)
        for sub in getattr(point, "sub_points", []) or []:
            if (sub or "").strip():
                lines.append(str(sub).strip())
    for card in getattr(slide_spec, "cards", None) or []:
        text = _format_card_text(card)
        if text:
            lines.extend([line for line in text.splitlines() if line.strip()])
    diagram = getattr(slide_spec, "diagram", None)
    prompt = (getattr(diagram, "prompt", "") or "").strip() if diagram else ""
    if prompt:
        lines.append("Diagram explanation: " + prompt)
    if not lines:
        lines.append("Content could not be rendered from the generated plan; please review this slide.")
    return lines[:8]


def _repair_title_only_slide(slide, slide_spec) -> bool:
    """Add visible fallback body text instead of failing on a title-only slide."""
    lines = _fallback_visible_lines(slide_spec)
    left, top, width, height = 0.75, 1.35, 11.8, 4.65
    body_lines = [(0, line) for line in lines if line.strip()]
    if not body_lines:
        return False
    font_pt = _fit_font_for_box(body_lines, width, height)
    _add_body(
        slide,
        left,
        top,
        width,
        height,
        body_lines,
        font_pt,
        color_hex=COLOR_BODY,
        bullet_hex=COLOR_BLUE,
    )
    log.warning(
        "Repaired title-only slide_id=%s with fallback visible body text.",
        getattr(slide_spec, "slide_id", "?"),
    )
    return True


def _validate_rendered_native_slide(slide, slide_spec) -> None:
    """Reject deterministic blank or partially populated template output."""
    archetype = (getattr(slide_spec, "archetype", "") or "").strip().lower()
    has_table = any(getattr(shape, "has_table", False) for shape in slide.shapes)
    has_picture = False
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not shape.is_placeholder:
                has_picture = True
                break
        except (AttributeError, NotImplementedError):
            continue
    body_texts = [
        (getattr(shape, "text", "") or "").strip()
        for shape in _body_placeholders(slide)
        if (getattr(shape, "text", "") or "").strip()
    ]
    generated_body_texts = []
    title_texts = {
        (getattr(shape, "text", "") or "").strip()
        for shape in _title_placeholders(slide)
        if (getattr(shape, "text", "") or "").strip()
    }
    for shape in slide.shapes:
        text = (getattr(shape, "text", "") or "").strip()
        if not text or text in title_texts:
            continue
        if getattr(shape, "is_placeholder", False):
            name = _placeholder_name(shape).lower()
            kind = _placeholder_type_name(shape)
            if any(token in name for token in ("footer", "slide number", "date", "presenter")):
                continue
            if any(token in kind for token in ("FOOTER", "SLIDE_NUMBER", "DATE")):
                continue
            continue
        generated_body_texts.append(text)

    if archetype == "agenda" and not has_table and len(body_texts) < 3:
        raise ValueError("Agenda rendered without populated agenda items")
    if getattr(slide_spec, "table", None) and not has_table:
        raise ValueError(f"Table slide {getattr(slide_spec, 'slide_id', '?')} rendered without a table")
    if (
        archetype not in {"title", "agenda"}
        and not body_texts
        and not generated_body_texts
        and not has_table
        and not has_picture
    ):
        if not _repair_title_only_slide(slide, slide_spec):
            raise ValueError(
                f"Slide {getattr(slide_spec, 'slide_id', '?')} rendered with a title but no visible content"
            )

    removed_orphan_numbers = _clear_orphan_number_placeholders(slide)
    if removed_orphan_numbers:
        log.debug(
            "Removed %d orphan number placeholder(s) from slide_id=%s.",
            removed_orphan_numbers,
            getattr(slide_spec, "slide_id", "?"),
        )

    filled_points = [
        shape for shape in _card_placeholders(slide)
        if (getattr(shape, "text", "") or "").strip()
    ]
    number_count = 0
    for shape in _number_placeholders(slide):
        if not (getattr(shape, "text", "") or "").strip():
            continue
        if any(_overlap_ratio(shape, point) >= 0.02 for point in filled_points):
            continue
        number_count += 1
    point_count = len(filled_points)
    if number_count > point_count:
        log.warning(
            "Slide %s still has %d visible number marker(s) for %d point(s); clearing markers.",
            getattr(slide_spec, "slide_id", "?"),
            number_count,
            point_count,
        )
        for shape in _number_placeholders(slide):
            _set_text(shape, "")


def _fill_divider_slide(
    slide, title: str, subtitle: str, prs: Presentation, page_no: Optional[int] = None
) -> None:
    """Section-break slide: big title, one-line subtitle, no boxes to fill.

    The "Divider Beam" layouts hide all master-inherited shapes
    (``showMasterSp="0"`` in the template), which is also where the
    copyright/confidentiality footer line lives on every other slide — so a
    divider is the one archetype where the footer has to be drawn explicitly
    rather than relying on inheritance. (Toggling ``showMasterSp`` back on was
    tried and made no visible difference, likely because PowerPoint composes
    a slide's master-shape inheritance independently of an automated export
    pass — drawing it directly is unambiguous either way.)
    """
    _set_native_title(slide, title, "divider")
    for shape in _body_placeholders(slide):
        name = _semantic_placeholder_name(slide, shape)
        if "subtitle" in name:
            _set_text(shape, subtitle, font_pt=16)
    dark = "dark" in _layout_name(slide.slide_layout).lower()
    _draw_footer(slide, prs, page_no, dark=dark)


def _fill_quote_slide(slide, quote: str, attribution: str) -> None:
    """A full-bleed statement slide (win theme / positioning), not a fabricated
    customer testimonial — the attribution names HCLTech, never the customer."""
    for shape in _title_placeholders(slide):
        _set_text(shape, "")
    for shape in _body_placeholders(slide):
        name = _semantic_placeholder_name(slide, shape)
        if "quote" in name:
            box_w = max(0.5, float(shape.width) / EMU_PER_INCH)
            box_h = max(0.5, float(shape.height) / EMU_PER_INCH)
            font_pt = _fit_font_for_box([(0, quote)], box_w, box_h, min_pt=16, start_pt=28)
            _set_text(shape, quote, font_pt=font_pt, bold=True)
        elif "name" in name:
            _set_text(shape, attribution, font_pt=12)


def _split_kpi_stat(kpi: str) -> tuple[str, str]:
    """Pull the numeric token out of a stat-shaped kpi string for the big
    number box, leaving the rest as the description. The token can appear
    anywhere in the string, not only at the start ("Save 40% on cycle time")."""
    text = (kpi or "").strip()
    match = _STAT_TOKEN_RE.search(text)
    if not match:
        return "", text
    stat = match.group(0).strip()
    rest = (text[: match.start()] + " " + text[match.end() :]).strip(" -–—:.")
    return stat, rest or text


def _is_two_stat_infographic_layout(slide) -> bool:
    name = _layout_name(slide.slide_layout).lower()
    return "infographics" in name and "two key points" in name


def _fill_infographic_two_key_point_slots(slide, items: list[str], key_message: Optional[str] = None) -> int:
    """Fill the 'Subhead N' + 'Text Placeholder N' pairs on the Infographics
    (2)/List/Two-key-points layout.

    Its label boxes are short (~0.8in tall) — below the generic
    ``_card_placeholders`` height floor used for every other key-point family,
    which exists so a stray small placeholder never gets mistaken for a card
    slot elsewhere. That floor makes this layout's own boxes invisible to the
    generic path, so it needs direct filling instead.
    """
    if key_message:
        key_ph = _key_message_placeholder(slide)
        if key_ph is not None:
            _set_text(key_ph, key_message, font_pt=11, bold=True)

    slots: dict = {}
    for shape in _body_placeholders(slide):
        name = _semantic_placeholder_name(slide, shape)
        if name.startswith("data"):
            continue
        match = re.match(r"subhead\s*(\d+)$", name)
        if match:
            slots.setdefault(int(match.group(1)), {})["heading"] = shape
            continue
        match = re.match(r"text placeholder\s*(\d+)$", name)
        if match:
            slots.setdefault(int(match.group(1)), {})["body"] = shape

    filled = 0
    for idx, text in enumerate(items[:2], start=1):
        slot = slots.get(idx)
        if not slot or not (text or "").strip():
            continue
        heading, _, rest = (text or "").partition("\n")
        body_text = rest.replace("\n", " ").strip() or heading
        if "heading" in slot:
            _set_text(slot["heading"], heading, font_pt=13, bold=True)
        if "body" in slot:
            # This box is short (~0.8in) — a flattened heading+body+bullets
            # string at a fixed 11pt with no fit sizing overflowed straight
            # into the next row's Subhead/Text-Placeholder pair on a real
            # generated slide. Fit it the same way _fill_card_placeholders does.
            _set_text(slot["body"], body_text, font_pt=_card_fit_pt(slot["body"], body_text))
        filled += 1
    return filled


def _fill_infographic_stats(slide, kpis: list[str]) -> int:
    """Fill the extra 'Data N' / 'Data description N' slots an infographic
    layout adds beside its regular key-point boxes (which the normal card /
    detailed-point fill already handles unchanged)."""
    stats = [k for k in (kpis or []) if (k or "").strip()][:2]
    if not stats:
        return 0
    slots: dict = {}
    for shape in _body_placeholders(slide):
        name = _semantic_placeholder_name(slide, shape)
        match = re.match(r"data(?:\s+(description))?\s*(\d+)$", name)
        if not match:
            continue
        n = int(match.group(2))
        slots.setdefault(n, {})["description" if match.group(1) else "value"] = shape
    filled = 0
    for idx, kpi in enumerate(stats, start=1):
        box = slots.get(idx)
        if not box:
            continue
        stat, rest = _split_kpi_stat(kpi)
        # A missing stat token (kpis is a coarse, model-authored field) leaves
        # the number box empty — never dump the full sentence into a box
        # sized for a handful of characters.
        if stat and "value" in box:
            _set_text(box["value"], stat, font_pt=28, bold=True)
        if "description" in box:
            _set_text(box["description"], rest, font_pt=12)
        filled += 1
    return filled


def _render_hcltech_native_slide(
    slide,
    prs: Presentation,
    title: str,
    bullets,
    diagram=None,
    diagram_bytes: Optional[bytes] = None,
    archetype: str = "Content",
    detailed_points=None,
    page_no: Optional[int] = None,
    key_message: Optional[str] = None,
    cards=None,
    comparison=None,
    kpis=None,
    table=None,
) -> None:
    """Populate an official HCLTech layout while preserving its master design."""
    archetype_l = (archetype or "").strip().lower()
    plain_text = False
    # `layout_hint` is not passed separately into this function, so use a marker
    # embedded in archetype by the page preparation path.
    if archetype_l.endswith("|plaintext"):
        archetype_l = archetype_l.replace("|plaintext", "").strip()
        archetype = archetype_l.title()
        plain_text = True
    if plain_text:
        _render_plain_text_slide(
            slide,
            prs,
            title,
            bullets,
            _theme_for("Content"),
            detailed_points=detailed_points,
            page_no=page_no,
            key_message=key_message,
            cards=cards,
            comparison=comparison,
            kpis=kpis,
        )
        return
    _set_native_title(slide, title, archetype)
    _fill_footer_placeholders(slide, page_no)

    cards = [c for c in (cards or []) if getattr(c, "heading", "").strip()]

    if archetype_l == "team" and _is_org_chart_layout(slide):
        role_points = [
            point for point in (detailed_points or [])
            if (getattr(point, "text", "") or "").strip()
        ]
        if _fill_org_chart_slide(slide, cards, role_points):
            return

    visual_archetypes = {
        "architecture",
        "deployment architecture",
        "high availability & dr",
        "timeline",
        "team",
        "delivery plan",
        "solution overview",
    }

    if archetype_l == "title":
        # Prefer the slide's key_message (the deck's win-thesis statement,
        # written as one complete sentence) over the first bullet — the cover
        # slide's bullets are often a list of terse noun-phrase candidates
        # ("Professional Assessments applications and platforms.") that read
        # as an unfinished fragment, not a tagline a reader should see first.
        subtitle = (key_message or "").strip() or next(
            (b.strip() for b in (bullets or []) if (b or "").strip()), ""
        )
        # The cover subtitle is a tagline, not a paragraph. A long value
        # proposition overflows the cover, so keep the first sentence within a
        # readable cap (cut at a word boundary).
        subtitle = re.sub(r"\s+", " ", subtitle).strip()
        if len(subtitle) > 165:
            first = re.split(r"(?<=[.;:])\s+", subtitle, maxsplit=1)[0]
            subtitle = first if 40 <= len(first) <= 165 else subtitle[:162].rsplit(" ", 1)[0] + "…"
        for shape in _body_placeholders(slide):
            if "subtitle" in _placeholder_name(shape).lower() or "SUBTITLE" in _placeholder_type_name(shape):
                _set_text(shape, subtitle, font_pt=None, bold=None)
                return
        return

    if archetype_l == "agenda":
        items = [t for lvl, t in _normalize_body_lines(bullets, detailed_points) if lvl == 0] or [
            t for _, t in _normalize_body_lines(bullets, detailed_points)
        ]
        if not _render_agenda_table(slide, items):
            filled = _fill_numbered_or_keypoint_placeholders(slide, items[:6])
        else:
            filled = len(items)
        if not filled:
            _fill_native_body(slide, prs, [f"- {item}" for item in items])
        return

    if archetype_l == "divider":
        subtitle = next((b.strip() for b in (bullets or []) if (b or "").strip()), "")
        _fill_divider_slide(slide, title, subtitle, prs, page_no)
        return

    if archetype_l == "win theme":
        quote = (key_message or "").strip() or next(
            (b.strip() for b in (bullets or []) if (b or "").strip()), ""
        )
        _fill_quote_slide(slide, quote, "HCLTech")
        return

    if table:
        sidebar = next(
            (
                shape for shape in _body_placeholders(slide)
                if "sidebar" in _semantic_placeholder_name(slide, shape)
            ),
            None,
        )
        if sidebar is not None:
            sidebar_lines = _native_slide_lines(
                bullets,
                detailed_points=detailed_points,
                key_message=key_message,
            )
            if sidebar_lines:
                box_w = float(sidebar.width) / EMU_PER_INCH
                box_h = float(sidebar.height) / EMU_PER_INCH
                fit_lines = [
                    (0, line.replace("- ", "").strip()) for line in sidebar_lines
                    if line.strip()
                ]
                _set_text(
                    sidebar,
                    "\n".join(sidebar_lines),
                    font_pt=_fit_font_for_box(fit_lines, box_w, box_h),
                    bold=None,
                )
        table_box = _placeholder_box(slide, "TABLE")
        if table_box is not None:
            _render_table(slide, prs, table, *table_box)
        else:
            x, y, w, h = _native_content_box(prs)
            if lines := _native_slide_lines(bullets, detailed_points=detailed_points, key_message=key_message):
                y += 0.45
                h = max(1.0, h - 0.45)
                _set_text(
                    slide.shapes.add_textbox(Inches(x), Inches(y - 0.45), Inches(w), Inches(0.35)),
                    " | ".join(line.replace("- ", "") for line in lines[:2]),
                    font_pt=10,
                )
            _render_table(slide, prs, table, x, y, w, h)
        return

    # Two-column comparison -> two key-point boxes (never a flattened textbox
    # that collides with the title).
    if comparison is not None:
        comp_cards = _comparison_to_cards(comparison)
        if comp_cards and _fill_card_placeholders(slide, comp_cards, key_message=key_message):
            return
        lines = _native_slide_lines(
            bullets, detailed_points=detailed_points,
            key_message=key_message, comparison=comparison,
        )
        _fill_native_body(slide, prs, lines)
        return

    if cards:
        stat_kpis = _stat_shaped_kpi_list(kpis)
        if stat_kpis and _is_two_stat_infographic_layout(slide):
            filled = _fill_infographic_two_key_point_slots(
                slide, [_format_card_text(c) for c in cards], key_message=key_message,
            )
            if filled:
                _fill_infographic_stats(slide, stat_kpis)
                return
        filled = _fill_card_placeholders(slide, cards, key_message=key_message)
        if filled:
            return
        _fill_native_body(slide, prs, [_format_card_text(c) for c in cards])
        return

    has_visual_diagram = bool(
        diagram
        and getattr(diagram, "approved", False)
        and archetype_l in visual_archetypes
    )
    detailed_points = [
        point for point in (detailed_points or [])
        if (getattr(point, "text", "") or "").strip()
    ]
    if detailed_points and not has_visual_diagram:
        stat_kpis = _stat_shaped_kpi_list(kpis)
        if stat_kpis and _is_two_stat_infographic_layout(slide):
            filled = _fill_infographic_two_key_point_slots(
                slide, [_format_detailed_point(p) for p in detailed_points], key_message=key_message,
            )
            if filled:
                _fill_infographic_stats(slide, stat_kpis)
                return
        filled = _fill_detailed_point_placeholders(
            slide,
            detailed_points,
            key_message=key_message,
        )
        if filled:
            return

    if not has_visual_diagram and comparison is None:
        flat_items = [
            text for level, text in _normalize_body_lines(bullets, None)
            if level == 0 and (text or "").strip()
        ]
        if _fill_flat_point_placeholders(slide, flat_items, key_message=key_message):
            return

    lines = _native_slide_lines(
        bullets,
        detailed_points=detailed_points,
        key_message=key_message,
        comparison=comparison,
    )
    if kpis:
        lines.extend(["", "Key metrics"])
        lines.extend(f"- {kpi}" for kpi in kpis if (kpi or "").strip())

    # Insert the diagram for visual archetypes, for image-only pages, and for a
    # captioned single-diagram slide (approved diagram + no cards/points, only a
    # key-message caption) regardless of archetype.
    should_insert_diagram = bool(
        diagram
        and getattr(diagram, "approved", False)
        and comparison is None
        and (archetype_l in visual_archetypes or not lines or (not cards and not detailed_points))
    )
    inserted_diagram = (
        _insert_native_diagram(slide, prs, diagram, diagram_bytes)
        if should_insert_diagram
        else None
    )

    if inserted_diagram:
        if lines:
            _fill_diagram_side_text(slide, inserted_diagram, lines, key_message)
        return
    _fill_native_body(slide, prs, lines)


# ------------------------------------------------------------------
# Slide renderers
# ------------------------------------------------------------------
def _render_title_slide(slide, prs: Presentation, title: str, bullets, theme: dict) -> None:
    """Navy cover slide: left accent stripe, dot-grid motif, hero title, footer."""
    w_in = float(prs.slide_width) / EMU_PER_INCH
    h_in = float(prs.slide_height) / EMU_PER_INCH

    # Left cyan stripe + top-right dot-grid decoration (cover motif).
    _add_rect(slide, 0, 0, max(0.16, w_in * 0.018), h_in, theme["accent"])
    grid_cols, grid_rows = 6, 4
    gap = 0.45
    grid_w = grid_cols * gap
    _dot_grid(slide, w_in - grid_w - 0.25, h_in * 0.09, grid_cols, grid_rows, gap_in=gap)

    lm = max(0.7, w_in * 0.06)
    box_w = w_in - lm - max(0.7, w_in * 0.06)

    title_pt = _fit_title_font(title, box_w, max_lines=3,
                               start_pt=FONT_TITLE_SLIDE_PT, min_pt=FONT_TITLE_SLIDE_PT - 12)

    tb = slide.shapes.add_textbox(Inches(lm), Inches(h_in * 0.30), Inches(box_w), Inches(h_in * 0.30))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(theme["title"])
    _set_run_font(run, FONT_NAME_HEAVY)
    p.alignment = PP_ALIGN.LEFT

    _add_rect(slide, lm, h_in * 0.62, min(box_w, 2.6), 0.06, theme["accent"])

    subtitle = next((b.strip() for b in (bullets or []) if (b or "").strip()), "")
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(lm), Inches(h_in * 0.66), Inches(box_w), Inches(h_in * 0.16))
        stf = sb.text_frame
        stf.clear()
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(16)
        srun.font.color.rgb = RGBColor.from_string(theme["body"])
        _set_run_font(srun)
        sp.alignment = PP_ALIGN.LEFT

    # Footer label on the cover.
    fb = slide.shapes.add_textbox(Inches(lm), Inches(h_in - 0.45), Inches(box_w), Inches(0.3))
    ftf = fb.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    frun = fp.add_run()
    frun.text = FOOTER_TEXT
    frun.font.size = Pt(10)
    frun.font.color.rgb = RGBColor.from_string(theme["body"])
    _set_run_font(frun)
    fp.alignment = PP_ALIGN.LEFT


def _render_plain_text_slide(
    slide,
    prs: Presentation,
    title: str,
    bullets,
    theme: dict,
    detailed_points=None,
    page_no: Optional[int] = None,
    key_message: Optional[str] = None,
    cards=None,
    comparison=None,
    kpis=None,
) -> None:
    """Controlled text slide that avoids fragile native placeholders."""
    if _title_placeholders(slide):
        _set_native_title(slide, title, "Content")
        content_top = 1.55
        _add_rect(slide, 0.75, content_top - 0.12, 1.6, 0.055, theme["accent"])
    else:
        content_top = _draw_header(slide, prs, title, theme)
    content_bottom = _draw_footer(slide, prs, page_no)
    w_in = float(prs.slide_width) / EMU_PER_INCH
    pad = max(0.5, w_in * 0.055)
    body_w = w_in - 2 * pad
    body_h = max(1.0, content_bottom - content_top - 0.05)

    cards = [card for card in (cards or []) if getattr(card, "heading", "").strip()]
    if cards:
        # Preserve the native title/footer while keeping structured proposal
        # content as cards instead of flattening it into one overflowing box.
        _render_cards_slide(slide, prs, pad, content_top, body_w, content_bottom, cards)
        return

    lines = _native_slide_lines(
        bullets,
        detailed_points=detailed_points,
        key_message=key_message,
        comparison=comparison,
    )
    if kpis:
        lines.extend(["", "Key metrics"])
        lines.extend(f"- {kpi}" for kpi in kpis if (kpi or "").strip())

    body_lines = [
        (1 if line.startswith("  - ") else 0, line[4:].strip() if line.startswith("  - ") else line[2:].strip() if line.startswith("- ") else line.strip())
        for line in lines
        if (line or "").strip()
    ]
    if not body_lines:
        return
    font_pt = _fit_font_for_box(body_lines, body_w, body_h)
    _add_body(
        slide,
        pad,
        content_top,
        body_w,
        body_h,
        body_lines,
        font_pt,
        color_hex=theme["body"],
        bullet_hex=theme["accent"],
    )


def _render_section_slide(slide, prs: Presentation, title: str, bullets, theme: dict,
                          detailed_points=None, page_no: Optional[int] = None) -> None:
    """Dark closing/section slide: navy bg, dot grid, white title, light bullets."""
    w_in = float(prs.slide_width) / EMU_PER_INCH
    h_in = float(prs.slide_height) / EMU_PER_INCH

    _add_rect(slide, 0, 0, max(0.16, w_in * 0.018), h_in, theme["accent"])
    grid_cols, grid_rows = 5, 3
    gap = 0.45
    _dot_grid(slide, w_in - grid_cols * gap - 0.25, h_in * 0.10, grid_cols, grid_rows, gap_in=gap)

    lm = max(0.7, w_in * 0.06)
    box_w = w_in - 2 * lm

    title_pt = _fit_title_font(title, box_w, max_lines=2, start_pt=32, min_pt=22)
    tb = slide.shapes.add_textbox(Inches(lm), Inches(h_in * 0.16), Inches(box_w), Inches(h_in * 0.22))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(theme["title"])
    _set_run_font(run, FONT_NAME_HEAVY)

    _add_rect(slide, lm, h_in * 0.40, min(box_w, 2.2), 0.05, theme["accent"])

    body_lines = _normalize_body_lines(bullets, detailed_points)
    if body_lines:
        by = h_in * 0.46
        bh = h_in - by - 0.5
        font_pt = _fit_font_for_box(body_lines, box_w, bh)
        _add_body(slide, lm, by, box_w, bh, body_lines, font_pt,
                  color_hex=theme["body"], bullet_hex=theme["accent"])


def _render_agenda_slide(slide, prs: Presentation, title: str, bullets, theme: dict,
                         detailed_points=None, page_no: Optional[int] = None) -> None:
    """Agenda slide: numbered navy chips down the left, item text beside each."""
    content_top = _draw_header(slide, prs, title or "Agenda", theme)
    content_bottom = _draw_footer(slide, prs, page_no)

    w_in = float(prs.slide_width) / EMU_PER_INCH
    pad = max(0.4, w_in * 0.04)

    items = [t for lvl, t in _normalize_body_lines(bullets, detailed_points) if lvl == 0] or \
            [t for _, t in _normalize_body_lines(bullets, detailed_points)]
    items = [t for t in items if t]
    if not items:
        return

    n = len(items)
    avail_h = content_bottom - content_top
    row_h = min(0.62, avail_h / n)
    chip = min(0.45, row_h * 0.8)
    gap = (avail_h - n * row_h) / max(1, n) if n * row_h < avail_h else 0.0
    y = content_top + max(0.0, gap / 2.0)
    font_pt = 16 if row_h > 0.5 else 13

    for i, text in enumerate(items, start=1):
        _add_rect(slide, pad, y, chip, chip, COLOR_PRIMARY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Number inside the chip.
        nb = slide.shapes.add_textbox(Inches(pad), Inches(y), Inches(chip), Inches(chip))
        ntf = nb.text_frame
        ntf.clear()
        ntf.word_wrap = False
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        npp = ntf.paragraphs[0]
        nrun = npp.add_run()
        nrun.text = f"{i:02d}"
        nrun.font.size = Pt(13)
        nrun.font.bold = True
        nrun.font.color.rgb = RGBColor.from_string(COLOR_WHITE)
        _set_run_font(nrun)
        npp.alignment = PP_ALIGN.CENTER
        # Item label beside the chip.
        lb = slide.shapes.add_textbox(
            Inches(pad + chip + 0.25), Inches(y), Inches(w_in - 2 * pad - chip - 0.25), Inches(chip)
        )
        ltf = lb.text_frame
        ltf.clear()
        ltf.word_wrap = True
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lpp = ltf.paragraphs[0]
        lrun = lpp.add_run()
        lrun.text = text
        lrun.font.size = Pt(font_pt)
        lrun.font.color.rgb = RGBColor.from_string(theme["body"])
        _set_run_font(lrun)
        lpp.alignment = PP_ALIGN.LEFT
        y += row_h + gap


def _render_cards_slide(slide, prs: Presentation, x_in: float, top_in: float,
                        w_in: float, bottom_in: float, cards) -> None:
    """Lay a list of Card objects out as a responsive 2-up grid."""
    cards = [c for c in (cards or []) if getattr(c, "heading", "").strip()]
    if not cards:
        return
    n = len(cards)
    cols = 1 if n == 1 else 2
    rows = math.ceil(n / cols)
    gutter = 0.30
    vgap = 0.24
    card_w = (w_in - gutter * (cols - 1)) / cols
    avail_h = bottom_in - top_in
    card_h = (avail_h - vgap * (rows - 1)) / rows

    for i, card in enumerate(cards):
        r, c = divmod(i, cols)
        cx = x_in + c * (card_w + gutter)
        cy = top_in + r * (card_h + vgap)
        accent = _resolve_accent(getattr(card, "accent", None), i)
        _add_card(
            slide, cx, cy, card_w, card_h,
            heading=getattr(card, "heading", ""),
            body=getattr(card, "body", ""),
            bullets=getattr(card, "bullets", None),
            accent_hex=accent,
        )


def _render_comparison_block(slide, x_in: float, top_in: float, w_in: float,
                             bottom_in: float, comparison) -> None:
    """Two-column comparison: coloured header band + icon-bullet list per side."""
    gutter = 0.40
    col_w = (w_in - gutter) / 2.0
    sides = [
        (comparison.left, x_in, "challenge"),
        (comparison.right, x_in + col_w + gutter, "goal"),
    ]
    head_h = 0.46
    for col, cx, default_key in sides:
        if col is None:
            continue
        accent = _resolve_accent(getattr(col, "accent", None) or default_key)
        # Column header band.
        _add_rect(slide, cx, top_in, col_w, head_h, accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        hb = slide.shapes.add_textbox(Inches(cx + 0.16), Inches(top_in), Inches(col_w - 0.32), Inches(head_h))
        htf = hb.text_frame
        htf.clear()
        htf.word_wrap = True
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE
        hp = htf.paragraphs[0]
        hr = hp.add_run()
        hr.text = (getattr(col, "heading", "") or "").strip().upper()
        hr.font.size = Pt(12)
        hr.font.bold = True
        hr.font.color.rgb = RGBColor.from_string(COLOR_WHITE)
        _set_run_font(hr, FONT_NAME_HEAVY)
        hp.alignment = PP_ALIGN.LEFT
        # Items as icon bullets.
        _add_icon_bullets(
            slide, cx, top_in + head_h + 0.18, col_w,
            max(0.6, bottom_in - top_in - head_h - 0.18),
            getattr(col, "items", None), accent_hex=accent,
        )


def _render_consulting_slide(
    slide,
    prs: Presentation,
    title: str,
    bullets,
    diagram=None,
    diagram_bytes: Optional[bytes] = None,
    archetype: str = "Content",
    detailed_points=None,
    page_no: Optional[int] = None,
    key_message: Optional[str] = None,
    cards=None,
    comparison=None,
    kpis=None,
    table=None,
):
    """Render one slide, dispatching on archetype and available structures."""
    theme = _theme_for(archetype)
    _set_slide_background(slide, theme["bg"])

    if theme["kind"] == "title":
        _render_title_slide(slide, prs, title, bullets, theme)
        return
    if theme["kind"] == "section":
        _render_section_slide(slide, prs, title, bullets, theme,
                              detailed_points=detailed_points, page_no=page_no)
        return
    if theme["kind"] == "agenda":
        _render_agenda_slide(slide, prs, title, bullets, theme,
                             detailed_points=detailed_points, page_no=page_no)
        return

    # ---- Standard content slide: clean header + body region ----
    content_top = _draw_header(slide, prs, title, theme)
    content_bottom = _draw_footer(slide, prs, page_no)

    w_in = float(prs.slide_width) / EMU_PER_INCH
    pad = max(0.4, w_in * 0.04)
    inner_w = w_in - 2 * pad

    # Optional key-message banner directly under the header.
    if key_message and (key_message or "").strip():
        content_top = _add_key_message(slide, pad, content_top, inner_w, key_message)

    # KPI chip row reserves space at the bottom.
    kpis = [k for k in (kpis or []) if (k or "").strip()]
    if kpis:
        kpi_y = content_bottom - 0.52
        _add_kpi_row(slide, pad, kpi_y, inner_w, kpis, accent_hex=COLOR_BLUE)
        content_bottom = kpi_y - 0.18

    body_lines = _normalize_body_lines(bullets, detailed_points)

    has_diagram = bool(
        diagram
        and getattr(diagram, "approved", False)
        and (diagram_bytes is not None or getattr(diagram, "image_path", None))
    )
    image_path = getattr(diagram, "image_path", None) if diagram else None

    # 0) Native editable table.
    if table:
        if body_lines:
            note_h = 0.45
            _add_body(slide, pad, content_top, inner_w, note_h, body_lines[:2], 10,
                      color_hex=theme["body"], bullet_hex=theme["accent"])
            content_top += note_h + 0.12
        _render_table(slide, prs, table, pad, content_top, inner_w, max(0.8, content_bottom - content_top))
        return

    # 1) Two-column comparison (problem vs. goal) takes precedence.
    if comparison is not None:
        _render_comparison_block(slide, pad, content_top, inner_w, content_bottom, comparison)
        return

    # 2) Card grid.
    cards = [c for c in (cards or []) if getattr(c, "heading", "").strip()]
    if cards:
        _render_cards_slide(slide, prs, pad, content_top, inner_w, content_bottom, cards)
        return

    # 3) Diagram slide.
    if has_diagram:
        if body_lines:
            # Image left, text right.
            box = _split_box(prs, content_top, content_bottom)
            ix, iy, iw, ih = box["image"]
            bx, by, bw, bh = box["body"]
            if diagram_bytes is not None:
                _place_image_contain(slide, diagram_bytes, ix, iy, iw, ih, inset_in=0.05)
            elif image_path is not None:
                img = Path(str(image_path))
                if img.exists():
                    _place_image_contain(slide, img, ix, iy, iw, ih, inset_in=0.05)
            font_pt = _fit_font_for_box(body_lines, bw, bh)
            _add_body(slide, bx, by, bw, bh, body_lines, font_pt,
                      color_hex=theme["body"], bullet_hex=theme["accent"])
        else:
            # Full-bleed visual (text lives on its own slide).
            ih = max(0.6, content_bottom - content_top)
            if diagram_bytes is not None:
                _place_image_contain(slide, diagram_bytes, pad, content_top, inner_w, ih, inset_in=0.05)
            elif image_path is not None:
                img = Path(str(image_path))
                if img.exists():
                    _place_image_contain(slide, img, pad, content_top, inner_w, ih, inset_in=0.05)
        return

    # 4) Plain bulleted body.
    box = _content_box(prs, content_top, content_bottom)
    bx, by, bw, bh = box["body"]
    font_pt = _fit_font_for_box(body_lines, bw, bh)
    _add_body(slide, bx, by, bw, bh, body_lines, font_pt,
              color_hex=theme["body"], bullet_hex=theme["accent"])


def _find_blank_layout(prs: Presentation):
    """Pick a layout with minimal placeholders."""
    bad_tokens = ("end plate", "closing", "thank", "contact", "appendix", "divider", "section")
    for layout in prs.slide_layouts:
        name = _layout_name(layout).lower()
        if any(token in name for token in bad_tokens):
            continue
        if len(getattr(layout, "placeholders", [])) == 0:
            return layout
    best = None
    best_n = 10**9
    for layout in prs.slide_layouts:
        name = _layout_name(layout).lower()
        if any(token in name for token in bad_tokens):
            continue
        n = len(getattr(layout, "placeholders", []))
        if n < best_n:
            best = layout
            best_n = n
    return best or prs.slide_layouts[0]


def _content_capacity_for_plan() -> tuple[float, float]:
    """Conservative body capacity before the concrete template is opened."""
    # 16:9 content area after title/footer/chrome. Split pages are rendered as
    # plain body slides, so use broad body capacity instead of small key-point
    # card capacity.
    return 10.4, 4.4


def _object_text_lines(value) -> list[tuple[int, str]]:
    if hasattr(value, "sub_points"):
        return _normalize_body_lines([], [value])
    if hasattr(value, "bullets"):
        return [(0, line) for line in _format_card_text(value).splitlines() if line.strip()]
    return [(0, str(value or "").strip())] if str(value or "").strip() else []


def _split_values_by_text_capacity(values: list, max_items: int, w_in: float, h_in: float) -> list[list]:
    pages: list[list] = []
    current: list = []
    for value in values:
        candidate = current + [value]
        lines: list[tuple[int, str]] = []
        for item in candidate:
            lines.extend(_object_text_lines(item))
        if current and (len(candidate) > max_items or not _text_fits_box(lines, w_in, h_in)):
            pages.append(current)
            current = [value]
        else:
            current = candidate
    if current:
        pages.append(current)
    return pages


def _cards_fit_grid(cards: list, w_in: float, h_in: float) -> bool:
    """Check card copy against the actual 1-up/2-up grid geometry."""
    if not cards:
        return True
    cols = 1 if len(cards) == 1 else 2
    rows = math.ceil(len(cards) / cols)
    card_w = (w_in - 0.30 * (cols - 1)) / cols
    card_h = (h_in - 0.24 * (rows - 1)) / rows
    return all(
        _text_fits_box(_object_text_lines(card), card_w - 0.32, card_h - 0.24, margin=0.90)
        for card in cards
    )


def _split_cards_by_grid_capacity(cards: list, w_in: float, h_in: float) -> list[list]:
    pages: list[list] = []
    current: list = []
    for card in cards:
        candidate = current + [card]
        if current and (len(candidate) > 4 or not _cards_fit_grid(candidate, w_in, h_in)):
            pages.append(current)
            current = [card]
        else:
            current = candidate
    if current:
        pages.append(current)
    return pages


def _rebalance_sparse_page_groups(groups: list[list]) -> list[list]:
    """Avoid 3+1/4+1 continuation splits when a balanced split is possible."""
    balanced = [list(group) for group in groups if group]
    if len(balanced) >= 2 and len(balanced[-1]) == 1 and len(balanced[-2]) >= 3:
        balanced[-1].insert(0, balanced[-2].pop())
    return balanced


def _split_bullets_by_text_capacity(bullets: list[str], max_items: int, w_in: float, h_in: float) -> list[list[str]]:
    pages: list[list[str]] = []
    current: list[str] = []
    for bullet in [b for b in bullets if (b or "").strip()]:
        candidate = current + [bullet]
        lines = [(0, item) for item in candidate]
        if current and (len(candidate) > max_items or not _text_fits_box(lines, w_in, h_in)):
            pages.append(current)
            current = [bullet]
        else:
            current = candidate
    if current:
        pages.append(current)
    return pages


def _word_count(text: str | None) -> int:
    return len(re.sub(r"\s+", " ", text or "").strip().split())


def _body_bullets_with_key_message(slide_spec, bullets: list[str]) -> tuple[list[str], Optional[str]]:
    key_message = (getattr(slide_spec, "key_message", None) or "").strip()
    clean_bullets = [b for b in bullets if (b or "").strip()]
    if key_message and _word_count(key_message) > 26:
        return [key_message] + clean_bullets, None
    return clean_bullets, key_message or getattr(slide_spec, "key_message", None)


def _diagram_prompt_fallback_bullets(slide_spec, bullets: list[str]) -> list[str]:
    """Convert missing-image prompt/explanation text into readable bullets."""
    cleaned: list[str] = []
    for bullet in bullets:
        if re.match(r"^\s*diagram\s+explanation\s*:", bullet or "", flags=re.I):
            continue
        text = (bullet or "").strip()
        if text:
            cleaned.append(text)
    if not cleaned:
        key_message = (getattr(slide_spec, "key_message", None) or "").strip()
        if key_message:
            cleaned = [key_message]

    units: list[str] = []
    for text in cleaned:
        parts = [
            part.strip(" .;:")
            for part in re.split(
                r"(?:\.\s+|;\s+|\n+|,\s+(?=(?:show|include|use|apply|route|connect|protect|support|provide|retain|centralize|centralise|validate|publish|store|monitor)\b))",
                text,
                flags=re.I,
            )
            if part.strip(" .;:")
        ]
        units.extend(parts or [text])

    bullets_out: list[str] = []
    for text in units:
        text = re.sub(r"^(title|layout|left|right|top|bottom|centre|center)\s*:\s*", "", text, flags=re.I).strip()
        if not text:
            continue
        # Keep complete clauses. Very long prompt instructions are omitted
        # rather than displayed as visibly truncated customer-facing copy.
        if len(text) > 190:
            clauses = [
                clause.strip(" ,.;:")
                for clause in re.split(r",\s+|\s+(?:and|with|while)\s+", text, flags=re.I)
                if clause.strip(" ,.;:")
            ]
            text = next((clause for clause in clauses if 35 <= len(clause) <= 190), "")
        if not text:
            continue
        if text not in bullets_out:
            bullets_out.append(text)
        if len(bullets_out) >= 5:
            break
    return bullets_out or ["The approved diagram asset is unavailable. Regenerate this diagram before customer delivery."]


def _sanitize_render_bullets(slide_spec, bullets: list[str]) -> list[str]:
    if any(re.match(r"^\s*diagram\s+explanation\s*:", bullet or "", flags=re.I) for bullet in bullets):
        return _diagram_prompt_fallback_bullets(slide_spec, bullets)
    return bullets


def _diagram_companion_title(title: str) -> str:
    clean = re.sub(r"\s*\(\d+\s+of\s+\d+\)\s*$", "", title or "", flags=re.I).strip()
    title_l = clean.lower()
    patterns = (
        (("migration", "cutover"), "How migration protects operational continuity"),
        (("testing", "quality", "acceptance"), "How testing builds release confidence"),
        (("ams", "warranty", "support"), "How warranty and AMS support operate"),
        (("deployment", "resilience", "dr"), "How deployment and resilience work"),
        (("roadmap", "timeline"), "How the delivery roadmap unfolds"),
        (("agile delivery", "delivery increments"), "How phased delivery protects continuity"),
        (("governance", "team", "squad"), "How governance and delivery teams operate"),
        (("data model", "information model"), "How the information model is organised"),
        (("data domain",), "How the operational data domains work together"),
        (("flow", "process"), "How the end-to-end operational flow works"),
        (("integration", "source and consumer"), "How integrations connect sources, controls and consumers"),
        (("architecture",), "How the proposed architecture works"),
        (("solution",), "How the proposed solution works"),
    )
    for tokens, companion_title in patterns:
        if any(token in title_l for token in tokens):
            return companion_title
    return "How to read the proposed design"


# Verbs that reliably begin a new parallel clause in these proposal sentences.
# Used (together with any gerund, ``*ing``) as a lookahead so clause splitting
# never fractures enumerations such as "FIH, ELP, GP4, and SAP".
_CLAUSE_VERBS = (
    "use|uses|provide|provides|support|supports|enable|enables|ensure|ensures|"
    "deliver|delivers|reduce|reduces|improve|improves|replace|replaces|"
    "validate|validates|normalize|normalizes|normalise|normalises|"
    "generate|generates|integrate|integrates|post|posts|retain|retains|"
    "monitor|monitors|automate|automates|capture|captures|consolidate|consolidates|"
    "derive|derives|protect|protects|connect|connects|route|routes|apply|applies|"
    "publish|publishes|store|stores|standardise|standardises|standardize|standardizes|"
    "centralise|centralises|centralize|centralizes|keep|keeps|drive|drives|"
    "align|aligns|maintain|maintains|track|tracks|expose|exposes"
)


def _key_message_to_points(text: str) -> list[str]:
    """Decompose a rich 'so what' sentence into 2-4 crisp parallel points.

    Thin diagram slides (architecture / solution overview) often carry only a
    single dense key-message sentence. Splitting it on genuine clause
    boundaries — a gerund or a known verb after a comma/semicolon — lets the
    companion fill a multi-box key-point layout instead of leaving empty boxes
    and an orphan number marker. Enumerations ("FIH, ELP, GP4, and SAP") are
    preserved because plain nouns are not treated as clause starters.
    """
    text = re.sub(r"\s+", " ", (text or "").strip())
    if not text:
        return []
    # Drop a leading framing clause ("... by/through:" style) so the parallel
    # list reads cleanly; only anchors unlikely to fire mid-sentence are used.
    body = text
    match = re.match(r"^(.{0,120}?(?:\bby\b|\bthrough\b|:))\s+(.+)$", text)
    if match and len(match.group(2)) > 55:
        body = match.group(2)

    clause_starter = rf"(?:[a-z]+ing|{_CLAUSE_VERBS})"
    split_re = (
        rf";\s+"
        rf"|,\s+then\s+|,\s+followed\s+by\s+"  # temporal sequences
        rf"|,\s+and\s+(?={clause_starter}\b)|,\s+(?={clause_starter}\b)"
    )
    parts = [
        part.strip(" ,.;:")
        for part in re.split(split_re, body)
        if part and part.strip(" ,.;:")
    ]
    parts = [p for p in parts if len(p.split()) >= 2]
    if len(parts) < 2:
        return [text]
    points = parts[:4]
    return [p[0].upper() + p[1:] if p else p for p in points]


def _complete_proposal_assertion(text: str) -> str:
    """Turn fragments and requirement clauses into customer-facing assertions."""
    clean = re.sub(r"\s+", " ", (text or "").strip()).rstrip(" .")
    if not clean:
        return ""
    clean = re.sub(
        r"^(?:the\s+)?(?:platform|system|solution|project)\s+(?:shall|should|must)\s+",
        "The proposed solution will ",
        clean,
        flags=re.I,
    )
    if re.match(r"^(?:explain|show|describe|confirm|validate|use|provide|ensure)\b", clean, flags=re.I):
        return clean
    if not re.search(
        r"\b(?:is|are|will|uses|provides|supports|enables|ensures|delivers|connects|"
        r"protects|creates|consolidates|routes|applies|includes|combines|separates)\b",
        clean,
        flags=re.I,
    ):
        clean = "The proposed solution supports " + clean[0].lower() + clean[1:]
    return clean + "."


def _diagram_companion_bullets(slide_spec, bullets: list[str]) -> list[str]:
    """Keep customer-facing authored explanation; never expose image prompts."""
    authored = [
        _complete_proposal_assertion(bullet)
        for bullet in bullets
        if (bullet or "").strip()
    ]
    if len(authored) >= 2:
        return authored[:5]

    detailed = [
        text
        for _, text in _normalize_body_lines([], getattr(slide_spec, "detailed_points", None))
        if (text or "").strip()
    ]
    if len(detailed) >= 2:
        return detailed[:5]

    # Thin slide: expand the single key message (or lone bullet) into parallel
    # points so the companion fills its layout cleanly.
    seed = (authored or detailed or [""])[0]
    key_message = (getattr(slide_spec, "key_message", None) or "").strip()
    expanded = _key_message_to_points(seed if len(seed) > len(key_message) else key_message)
    if len(expanded) >= 2:
        return [_complete_proposal_assertion(item) for item in expanded]
    return [_complete_proposal_assertion(seed or key_message)] if (seed or key_message) else []


def _dedupe_page_groups(groups: list[list]) -> list[list]:
    """Remove duplicate/prefix continuation groups caused by model repetition."""
    result: list[list] = []
    seen: set[tuple[str, ...]] = set()
    for group in groups:
        normalized = tuple(re.sub(r"\s+", " ", str(item)).strip().lower() for item in group)
        if not normalized or normalized in seen:
            continue
        if any(
            len(normalized) <= len(prev)
            and normalized == prev[: len(normalized)]
            for prev in seen
        ):
            continue
        seen.add(normalized)
        result.append(group)
    return result


def _normalized_slide_title(title: str | None) -> str:
    text = re.sub(r"\s+", " ", title or "").strip().lower()
    return re.sub(r"\s*\(\d+\s+of\s+\d+\)\s*$", "", text)


def _normalized_slide_text(text: str | None) -> str:
    text = re.sub(r"\s+", " ", text or "").strip().lower()
    text = re.sub(r"^[\-•\d\.\)\s]+", "", text)
    return text


def _slide_text_units(slide_spec) -> list[str]:
    units: list[str] = []
    for value in getattr(slide_spec, "bullets", None) or []:
        if normalized := _normalized_slide_text(value):
            units.append(normalized)
    for point in getattr(slide_spec, "detailed_points", None) or []:
        if normalized := _normalized_slide_text(getattr(point, "text", "")):
            units.append(normalized)
        for sub in getattr(point, "sub_points", None) or []:
            if normalized := _normalized_slide_text(sub):
                units.append(normalized)
    for card in getattr(slide_spec, "cards", None) or []:
        if normalized := _normalized_slide_text(getattr(card, "heading", "")):
            units.append(normalized)
        for bullet in getattr(card, "bullets", None) or []:
            if normalized := _normalized_slide_text(bullet):
                units.append(normalized)
    table = getattr(slide_spec, "table", None)
    if isinstance(table, dict):
        for row in table.get("rows") or []:
            if normalized := _normalized_slide_text(" | ".join(str(cell) for cell in row)):
                units.append(normalized)
    diagram = getattr(slide_spec, "diagram", None)
    if diagram is not None:
        prompt = _normalized_slide_text(getattr(diagram, "prompt", ""))
        kind = _normalized_slide_text(getattr(diagram, "kind", ""))
        if prompt or kind:
            units.append(f"diagram:{kind}:{prompt}")
    return units


def _slide_content_key(slide_spec) -> tuple[str, str, tuple[str, ...]]:
    return (
        (getattr(slide_spec, "archetype", "") or "").strip().lower(),
        _normalized_slide_title(getattr(slide_spec, "title", "")),
        tuple(_slide_text_units(slide_spec)),
    )


def _is_prefix_or_subset_repeat(previous, current) -> bool:
    prev_title = _normalized_slide_title(getattr(previous, "title", ""))
    cur_title = _normalized_slide_title(getattr(current, "title", ""))
    if prev_title != cur_title:
        return False
    prev_units = _slide_text_units(previous)
    cur_units = _slide_text_units(current)
    if not prev_units or not cur_units:
        return False
    if len(cur_units) <= len(prev_units) and cur_units == prev_units[: len(cur_units)]:
        return True
    return set(cur_units).issubset(set(prev_units))


def _dedupe_render_slides(slides: list) -> list:
    """Remove final exact/prefix duplicates before PPTX creation."""
    deduped: list = []
    seen: set[tuple[str, str, tuple[str, ...]]] = set()
    for slide in slides:
        key = _slide_content_key(slide)
        if key in seen:
            log.info(
                "Dropping duplicate rendered slide %s (%r)",
                getattr(slide, "slide_id", "?"),
                getattr(slide, "title", ""),
            )
            continue
        if deduped and _is_prefix_or_subset_repeat(deduped[-1], slide):
            log.info(
                "Dropping prefix/subset duplicate rendered slide %s (%r)",
                getattr(slide, "slide_id", "?"),
                getattr(slide, "title", ""),
            )
            continue
        seen.add(key)
        deduped.append(slide)
    return deduped


def _clean_continuation_title(title: str | None) -> str:
    return re.sub(r"\s*\(1\s+of\s+2\)\s*$", "", title or "", flags=re.I).strip()


def _normalize_singleton_continuation_titles(slides: list) -> list:
    """Remove stale '(1 of 2)' titles when page 2 was de-duplicated."""
    base_counts: dict[str, int] = {}
    for slide in slides:
        base = _normalized_slide_title(getattr(slide, "title", ""))
        base_counts[base] = base_counts.get(base, 0) + 1
    for slide in slides:
        base = _normalized_slide_title(getattr(slide, "title", ""))
        if base_counts.get(base, 0) == 1:
            cleaned = _clean_continuation_title(getattr(slide, "title", ""))
            if cleaned:
                slide.title = cleaned
    return slides


# Usable text area (inches, after insets) inside each key-point box, keyed by
# how many boxes share the slide. Fewer boxes per slide => wider, taller boxes.
# Calibrated to the real HCLTech key-point box geometry so a dense card lands on
# the widest layout that still holds it at 14pt (avoids over-splitting).
_KEYPOINT_BOX_14 = {
    1: (11.4, 5.0),
    2: (5.6, 3.45),
    3: (3.4, 3.05),
    4: (2.4, 3.05),
    5: (1.85, 3.05),
}


def _item_fits_at_pt(item, w_in: float, h_in: float, font_pt: int = 14) -> bool:
    lines = _object_text_lines(item)
    if not lines:
        return True
    return _estimated_required_height_in(lines, w_in, font_pt) <= h_in * 0.94


def _split_items_to_fit_14pt(items: list) -> list[list]:
    """Split a card / detailed-point set so every item renders at >=14pt.

    Picks the largest boxes-per-slide count at which every item still fits at
    14pt (fewer boxes => wider boxes), then splits into balanced pages. No text
    is trimmed — density is solved with more slides, per the user's choice.
    """
    n = len(items)
    if n <= 1:
        return [items] if items else []
    # Keep three concise proof points in one native composition. Rendering at
    # 11pt is preferable to a visually unrelated 2+1 continuation.
    if n == 3:
        w_in, h_in = _KEYPOINT_BOX_14[3]
        if all(_item_fits_at_pt(item, w_in, h_in, 11) for item in items):
            return [items]
    # Never explode to one card per slide: a lone genuinely-oversized card would
    # spawn a slide each and still not reach 14pt. Two-up (wide boxes) is the
    # densest split we allow; the rare oversized card then renders at ~13pt.
    per_page = 2
    for k in range(min(n, 5), 1, -1):
        w_in, h_in = _KEYPOINT_BOX_14[k]
        if all(_item_fits_at_pt(item, w_in, h_in, 14) for item in items):
            per_page = k
            break
    if per_page >= n:
        return [items]
    page_count = math.ceil(n / per_page)
    base, extra = divmod(n, page_count)
    groups: list[list] = []
    idx = 0
    for i in range(page_count):
        size = base + (1 if i < extra else 0)
        groups.append(items[idx:idx + size])
        idx += size
    return groups


def _is_exec_summary_spec(slide_spec) -> bool:
    title = (getattr(slide_spec, "title", "") or "").strip().lower()
    return "executive summary" in title or "executive overview" in title


def _render_pages_for_slide(
    slide_spec,
    diagram_images: Optional[dict[str, bytes]] = None,
    native: bool = False,
) -> list:
    """Split text-heavy slides before rendering so PowerPoint never clips copy.

    On native HCLTech layouts (``native=True``) the template's key-point layouts
    hold whole card/point sets in tall boxes, so small sets are kept on a single
    slide. Pagination only kicks in when content genuinely overflows even the
    largest matching layout (e.g. long tables), matching the "clean continuation
    slide only when needed" policy.
    """
    table = getattr(slide_spec, "table", None)
    rows = list((table or {}).get("rows") or []) if isinstance(table, dict) else []
    cards = list(getattr(slide_spec, "cards", None) or [])
    detailed_points = list(getattr(slide_spec, "detailed_points", None) or [])
    bullets, adjusted_key_message = _body_bullets_with_key_message(
        slide_spec,
        list(getattr(slide_spec, "bullets", None) or []),
    )
    bullets = _sanitize_render_bullets(slide_spec, bullets)
    archetype = (getattr(slide_spec, "archetype", "") or "").strip().lower()
    has_diagram = _has_renderable_diagram(slide_spec, diagram_images)

    if has_diagram and archetype not in {"title", "agenda"}:
        image_page = slide_spec.model_copy(deep=True)
        image_page.bullets = []
        image_page.detailed_points = []
        image_page.cards = []
        image_page.comparison = None
        image_page.kpis = []
        image_page.layout_hint = None

        authored_points = [
            point for point in (getattr(slide_spec, "detailed_points", None) or [])
            if (getattr(point, "text", "") or "").strip()
        ]
        # Only authored explanation earns a companion slide. Earlier versions
        # synthesized explanation from the diagram prompt itself, which doubled
        # visual sections and turned a focused plan into a deterministic long
        # deck. Prompt text remains available in notes/fallback handling, but it
        # no longer creates an extra audience-facing page by itself.
        companion_points = [] if len(authored_points) >= 2 else list(bullets)

        # A separate "How X works" companion earns its own slide only when there
        # is enough explanatory content to fill it. Otherwise a lone-sentence
        # companion "doesn't make sense" — merge that sentence onto the diagram
        # slide as a caption and keep it to a single slide.
        has_companion_content = len(authored_points) >= 2 or len(companion_points) >= 2
        if not has_companion_content:
            image_page.key_message = adjusted_key_message  # caption beside the diagram
            return [image_page]

        image_page.key_message = None  # the companion carries the explanation
        companion = slide_spec.model_copy(deep=True)
        companion.slide_id = f"{slide_spec.slide_id}_interpretation"
        companion.title = _diagram_companion_title(slide_spec.title)
        companion.archetype = "Content"
        companion.layout_hint = None
        companion.diagram = None
        companion.cards = []
        companion.comparison = None
        companion.kpis = []
        companion.key_message = adjusted_key_message
        if len(authored_points) >= 2:
            companion.bullets = []
            companion.detailed_points = authored_points
        else:
            companion.bullets = companion_points
            companion.detailed_points = []
        # Run the companion back through pagination so its explanation is split
        # to render at 14pt too (it has no diagram, so this recurses once).
        companion_pages = _render_pages_for_slide(companion, None, native=native)
        return [image_page] + companion_pages

    if _has_approved_diagram(slide_spec) and not has_diagram and archetype not in {"title", "agenda"}:
        page = slide_spec.model_copy(deep=True)
        page.diagram = None
        page.layout_hint = None
        page.bullets = _diagram_prompt_fallback_bullets(slide_spec, bullets)
        page.detailed_points = []
        page.cards = []
        page.comparison = None
        page.kpis = []
        page.key_message = None
        page.archetype = "Content"
        return [page]

    collection = None
    page_groups: list[list] = []
    # Native HCLTech key-point layouts provide 2-5 tall boxes; whole small sets
    # fit one slide, so only very large content paginates on the native path.
    if native:
        table_page_size = 5 if archetype == "software bill of materials" else 10
    else:
        table_page_size = 6 if archetype == "software bill of materials" else 8
    if len(rows) > table_page_size:
        collection = "table"
        page_groups = [rows[idx : idx + table_page_size] for idx in range(0, len(rows), table_page_size)]
    elif not has_diagram and cards:
        if native:
            # Keep the Executive Summary on one slide; otherwise split dense card
            # sets so every card renders at 14pt (more slides, no trimming).
            # A "team" slide with no diagram is the org-chart-grid case
            # (_fill_org_chart_slide) — splitting its roster across pages
            # produces two sparse, disconnected half-charts instead of one
            # coherent org chart, so it stays on one slide too.
            page_groups = (
                [cards] if _is_exec_summary_spec(slide_spec) or archetype == "team"
                else _split_items_to_fit_14pt(cards)
            )
        else:
            w_in, h_in = _content_capacity_for_plan()
            if adjusted_key_message:
                h_in -= 0.45
            page_groups = _split_cards_by_grid_capacity(cards, w_in, h_in)
        if len(page_groups) > 1:
            collection = "cards"
    elif not has_diagram and detailed_points:
        if native:
            page_groups = (
                [detailed_points] if _is_exec_summary_spec(slide_spec) or archetype == "team"
                else _split_items_to_fit_14pt(detailed_points)
            )
        else:
            w_in, h_in = _content_capacity_for_plan()
            if adjusted_key_message:
                h_in -= 0.45
            total_words = sum(
                _word_count(getattr(point, "text", ""))
                + sum(_word_count(sub) for sub in (getattr(point, "sub_points", None) or []))
                for point in detailed_points
            )
            if len(detailed_points) <= 4 and total_words <= 180:
                page_groups = [detailed_points]
            else:
                page_groups = _split_values_by_text_capacity(detailed_points, 3, w_in, h_in)
        if len(page_groups) > 1:
            collection = "detailed_points"
    elif not has_diagram and archetype not in {"title", "agenda"} and bullets:
        if native and archetype == "next steps":
            page_groups = [bullets]  # keep the numbered next-steps flow on one slide
        elif native:
            # Split so every bullet renders at 14pt in its key-point box.
            page_groups = _split_items_to_fit_14pt(bullets)
        else:
            w_in, h_in = _content_capacity_for_plan()
            if adjusted_key_message:
                h_in -= 0.45
            max_bullets = 5 if archetype == "next steps" else 4
            page_groups = _split_bullets_by_text_capacity(bullets, max_bullets, w_in, h_in)
        if len(page_groups) > 1:
            collection = "bullets"

    if collection is None:
        if adjusted_key_message != getattr(slide_spec, "key_message", None):
            page = slide_spec.model_copy(deep=True)
            page.key_message = adjusted_key_message
            page.bullets = bullets
            page.layout_hint = None
            return [page]
        if not has_diagram and not rows and archetype not in {"title", "agenda"}:
            page = slide_spec.model_copy(deep=True)
            page.layout_hint = None
            return [page]
        return [slide_spec]

    page_groups = _rebalance_sparse_page_groups(_dedupe_page_groups(page_groups))
    if len(page_groups) <= 1:
        return [slide_spec]

    pages = []
    page_count = len(page_groups)
    for page_idx, page_values in enumerate(page_groups):
        page = slide_spec.model_copy(deep=True)
        page.slide_id = f"{slide_spec.slide_id}_page_{page_idx + 1}"
        page.title = f"{slide_spec.title} ({page_idx + 1} of {page_count})"
        if collection == "table":
            page.table["rows"] = page_values
            # A long table pages onto many consecutive slides that would
            # otherwise look identical. Alternating in the sidebar variant
            # every other page breaks up the run, and a short "part N of M"
            # note fills the sidebar when the slide has no caption of its own.
            if page_count > 1 and page_idx % 2 == 1:
                page.layout_hint = _TABLE_LAYOUTS[1]
                if not (page.key_message or "").strip():
                    page.key_message = f"Continued — part {page_idx + 1} of {page_count}"
        else:
            setattr(page, collection, page_values)
            # Split pages should use a broad body layout, not tiny numbered
            # key-point boxes. Keep the original slide's text, but render it
            # through the plain content path.
            if collection == "detailed_points":
                page.bullets = []
            elif collection == "bullets":
                page.key_message = adjusted_key_message if page_idx == 0 else None
            if collection != "cards":
                page.cards = []
            page.comparison = None
            page.layout_hint = None
            if collection in {"cards", "detailed_points"} and page_idx > 0:
                page.key_message = None
            if page_idx > 0:
                # kpis are a headline stat for the slide as a whole; without
                # this every continuation page would repeat the same stat
                # callout (and could itself re-trigger the infographic layout).
                page.kpis = []
        pages.append(page)
    return pages


def _template_is_native_hcltech(template_pptx: Union[Path, bytes, None]) -> bool:
    """Detect the official HCLTech expanded template from a path or bytes."""
    if template_pptx is None:
        return False
    try:
        prs = (
            Presentation(BytesIO(template_pptx))
            if isinstance(template_pptx, (bytes, bytearray))
            else Presentation(str(template_pptx))
        )
        return _is_hcltech_template(prs)
    except Exception:
        return False


def rendered_slide_count(
    deck_plan: DeckPlan,
    diagram_images: Optional[dict[str, bytes]] = None,
    template_pptx: Union[Path, bytes, None] = None,
) -> int:
    """Return the final PPTX page count after diagram companions and pagination."""
    native = _template_is_native_hcltech(template_pptx)
    render_slides = []
    for slide_spec in deck_plan.slides:
        render_slides.extend(_render_pages_for_slide(slide_spec, diagram_images, native=native))
    return len(_normalize_singleton_continuation_titles(_dedupe_render_slides(render_slides)))


def _add_customer_logo(slide, prs: Presentation, logo_bytes: bytes | None) -> None:
    """Overlay customer branding without modifying the template or slide master."""
    if not logo_bytes:
        return
    slide_w = float(prs.slide_width) / EMU_PER_INCH
    slide_h = float(prs.slide_height) / EMU_PER_INCH
    max_w = 1.20
    max_h = 0.28
    hcltech_reserved_w = 1.15
    right_pad = 0.42
    gap = 0.16
    right_edge = slide_w - right_pad - hcltech_reserved_w - gap
    try:
        picture = slide.shapes.add_picture(
            BytesIO(logo_bytes),
            Inches(max(0.4, right_edge - max_w)),
            Inches(slide_h - 0.57),
            height=Inches(max_h),
        )
    except Exception as exc:
        raise ValueError("The uploaded customer logo is not a valid PNG or JPEG image.") from exc
    width_in = float(picture.width) / EMU_PER_INCH
    if width_in > max_w:
        scale = max_w / width_in
        picture.width = Inches(max_w)
        picture.height = int(picture.height * scale)
    picture.left = Inches(right_edge - float(picture.width) / EMU_PER_INCH)
    picture.top = Inches(slide_h - 0.43 - float(picture.height) / EMU_PER_INCH / 2)
    picture.name = "Customer Logo"


_CONTINUATION_PAGE_SUFFIX_RE = re.compile(r"_page_\d+$")


def _continuation_group_key(slide_id: str) -> str:
    """The pre-pagination slide_id shared by every page split off one slide.

    ``_render_pages_for_slide`` names split pages ``f"{slide_id}_page_{n}"``
    (and recurses for a diagram companion's own split under
    ``f"{slide_id}_interpretation_page_{n}"``) — stripping that suffix groups
    a "(1 of N)"/"(2 of N)" run back together without conflating an unrelated
    diagram image page with its text companion (those already have distinct
    base ids and are correctly treated as separate slides).
    """
    return _CONTINUATION_PAGE_SUFFIX_RE.sub("", slide_id or "")


def _lock_continuation_layouts(prs: Presentation, render_slides: list) -> dict:
    """Choose one native layout per original slide and stamp it onto every
    page pagination split off from it, via ``layout_hint``.

    Without this, each split page independently calls ``_choose_hcltech_layout``
    and the shared layout-variety usage counter (see ``_pick_varied_layout``)
    makes page 2 avoid whatever sibling layout page 1 just picked — guaranteed
    alternation within a 2-member family, and (when a later page also lost its
    ``kpis``/other fields during splitting) sometimes a completely different
    layout family. Locking the choice to the page carrying the fullest content
    (page 1) and propagating it via the existing ``layout_hint`` fast-path
    (``_choose_hcltech_layout`` line ~1160) keeps continuation pages visually
    identical while leaving variety rotation between *distinct* slides intact.
    ``_hint_is_compatible`` still guards each page, so a page whose own content
    genuinely doesn't fit the locked layout falls back to normal selection
    instead of forcing a bad fit.

    Table pagination is deliberately excluded: ``_render_pages_for_slide``
    already alternates a long table's odd pages onto ``_TABLE_LAYOUTS[1]`` on
    purpose, so a run of many consecutive table slides doesn't look identical
    (see the "Alternating in the sidebar variant" comment there). Locking
    those pages to one layout here would undo that intentional design.
    """
    groups: dict[str, list] = {}
    for slide_spec in render_slides:
        key = _continuation_group_key(getattr(slide_spec, "slide_id", ""))
        groups.setdefault(key, []).append(slide_spec)

    usage: dict = {}
    for pages in groups.values():
        if any(_table_has_content(getattr(page, "table", None)) for page in pages):
            continue
        primary = pages[0]
        layout = _choose_hcltech_layout(prs, primary, usage)
        if layout is None:
            continue
        for page in pages:
            page.layout_hint = layout.name
    return usage


def render_deck_from_template(
    deck_plan: DeckPlan,
    template_pptx: Union[Path, bytes],
    out_path: Optional[Path] = None,
    diagram_images: Optional[dict[str, bytes]] = None,
    customer_logo: Optional[bytes] = None,
) -> Union[Path, bytes]:
    """Render a new PPTX from a template and a deck plan.

    Customer branding is added as a slide-level overlay; the source template,
    masters, layouts, and existing HCLTech branding remain unchanged.
    """
    if isinstance(template_pptx, (bytes, bytearray)):
        prs = Presentation(BytesIO(template_pptx))
    else:
        prs = Presentation(str(template_pptx))
    use_hcltech_native_layouts = _is_hcltech_template(prs)

    render_slides = []
    for slide_spec in deck_plan.slides:
        prepared = slide_spec.model_copy(deep=True)
        if (
            (prepared.archetype or "").strip().lower() == "title"
            and (prepared.title or "").strip().lower() in {"title", "proposal", "proposal title"}
            and (deck_plan.deck_title or "").strip()
        ):
            prepared.title = deck_plan.deck_title.strip()
        render_slides.extend(
            _render_pages_for_slide(prepared, diagram_images, native=use_hcltech_native_layouts)
        )
    render_slides = _dedupe_render_slides(render_slides)
    render_slides = _normalize_singleton_continuation_titles(render_slides)

    n_slides = len(render_slides)
    n_images = len(diagram_images) if diagram_images else 0
    log.info("Rendering deck: %d slides, %d diagram image(s)", n_slides, n_images)
    blank_layout = _find_blank_layout(prs)
    if use_hcltech_native_layouts:
        log.info("Detected HCLTech expanded template; rendering with native layouts.")

    _remove_presentation_sections(prs)

    # remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId  # pylint: disable=protected-access
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]  # pylint: disable=protected-access

    layout_usage: dict = (
        _lock_continuation_layouts(prs, render_slides) if use_hcltech_native_layouts else {}
    )
    for idx, slide_spec in enumerate(render_slides):
        layout = (
            _choose_hcltech_layout(prs, slide_spec, layout_usage)
            if use_hcltech_native_layouts
            else blank_layout
        )
        # Layout selection must always yield a usable layout; never crash on None.
        slide = prs.slides.add_slide(layout or blank_layout)

        archetype = getattr(slide_spec, "archetype", "Content") or "Content"
        is_plain_text = False
        if is_plain_text:
            removed_shapes = _clear_slide_shapes(slide, preserve_title=True)
            if removed_shapes:
                log.debug(
                    "Cleared %d inherited shape(s) for plain text slide_id=%s.",
                    removed_shapes,
                    getattr(slide_spec, "slide_id", "?"),
                )
        # No page number on the cover slide.
        page_no = None if (archetype or "").strip().lower() == "title" else idx + 1

        diagram_bytes = None
        if diagram_images is not None:
            diagram_bytes = _diagram_image_for_slide(slide_spec, diagram_images)
        try:
            render_func = (
                _render_hcltech_native_slide
                if use_hcltech_native_layouts
                else _render_consulting_slide
            )
            render_archetype = archetype
            if is_plain_text:
                render_archetype = f"{archetype}|PlainText"
            render_func(
                slide,
                prs,
                slide_spec.title or "",
                slide_spec.bullets or [],
                diagram=getattr(slide_spec, "diagram", None),
                diagram_bytes=diagram_bytes,
                archetype=render_archetype,
                detailed_points=getattr(slide_spec, "detailed_points", None),
                page_no=page_no,
                key_message=getattr(slide_spec, "key_message", None),
                cards=getattr(slide_spec, "cards", None),
                comparison=getattr(slide_spec, "comparison", None),
                kpis=getattr(slide_spec, "kpis", None),
                table=getattr(slide_spec, "table", None),
            )
        except Exception:
            log.exception(
                "Failed to render slide_id=%s (title=%r)",
                getattr(slide_spec, "slide_id", "?"),
                slide_spec.title,
            )
            raise

        if use_hcltech_native_layouts:
            _remove_overlapping_generated_pictures(slide, slide_spec)
            removed_placeholders = _remove_unused_placeholders(slide)
            if removed_placeholders:
                log.debug(
                    "Removed %d unused placeholder(s) from slide_id=%s.",
                    removed_placeholders,
                    getattr(slide_spec, "slide_id", "?"),
                )
            _validate_rendered_native_slide(slide, slide_spec)
        _add_customer_logo(slide, prs, customer_logo)
        _set_speaker_notes(slide, getattr(slide_spec, "notes", None))

    if out_path is not None:
        prs.save(str(out_path))
        log.info("Deck rendered to %s", out_path)
        return out_path
    stream = BytesIO()
    prs.save(stream)
    data = stream.getvalue()
    log.info("Deck rendered to in-memory PPTX (%d bytes)", len(data))
    return data
