from __future__ import annotations

import unittest
import base64
import json
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from types import SimpleNamespace
from unittest.mock import Mock, patch

from docx import Document
from openpyxl import Workbook
from openai import OpenAIError

from rfp2deck.agent.prompts import RFP_UNDERSTAND_PROMPT
from rfp2deck.agent.nodes import _extract_evidence_chunk, extract_source_evidence
from rfp2deck.agent.nodes import (
    enrich_understanding_risks,
    enforce_slide_density,
    order_deck,
    _compact_deck_plan_prompt,
    _chunked_plan_input,
    _fallback_deck_plan,
    _fallback_visual_briefs,
    _source_grounded_technical_architecture_elements,
    _proposal_section_skeleton,
    _chunked_deck_plan,
    _risk_detailed_points,
    _ai_ml_opportunities,
    _source_grounded_technology_table,
    _technology_recommendation_table,
    _align_recommendations_to_customer_platform,
    _source_grounded_technology_fallback,
    build_solution_brief,
    _sbom_table,
    _build_diagram_prompt,
    _testing_proposal_points,
    _ams_proposal_points,
    _assumptions_dependency_points,
    _data_domain_points,
    ensure_diagrams_for_key_slides,
    enrich_slide_detail,
    ensure_required_slides,
    consulting_grade_proposal_polish,
    prune_profile_misaligned_slides,
    prune_redundant_storyline_slides,
    plan_deck,
    derive_technology_recommendations,
    generate_notes,
)
from rfp2deck.agent.state import AgentState
from rfp2deck.diagrams import generator as diagram_generator
from rfp2deck.agent.evidence import (
    merge_evidence_batches,
    render_evidence_for_prompt,
    split_source_document,
)
from rfp2deck.core.schemas import (
    BulletPoint,
    DeckPlan,
    DeckNotes,
    DiagramSpec,
    DiagramBrief,
    EngagementProfile,
    LifecycleStageAssessment,
    RFPUnderstanding,
    Requirement,
    SlideSpec,
    SlideNote,
    SourceDocument,
    SourceEvidenceBatch,
    SolutionComponentDecision,
    TechnologyRecommendation,
    TechnologyRecommendationSet,
)
from rfp2deck.ingestion.docx_parser import parse_docx
from rfp2deck.ingestion.source_package import (
    build_source_reconciliation,
    classify_source,
    parse_source_document,
    render_source_package,
    sort_sources_by_precedence,
)
from rfp2deck.llm.structured import _dereference, _is_proxy_block_error, _make_strict
from rfp2deck.qa.coverage import build_traceability_report
from rfp2deck.rendering.pptx_renderer import (
    _concise_key_message,
    _dedupe_render_slides,
    _normalize_singleton_continuation_titles,
    _remove_overlapping_generated_pictures,
    _render_pages_for_slide,
    _repair_title_only_slide,
    _choose_hcltech_layout,
    render_deck_from_template,
    rendered_slide_count,
)


class SourcePackageTests(unittest.TestCase):
    def test_customer_logo_is_added_to_every_slide_without_master_changes(self) -> None:
        from pptx import Presentation

        template = Presentation()
        template_stream = BytesIO()
        template.save(template_stream)
        logo = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
        )
        deck = DeckPlan(deck_title="Test", slides=[
            SlideSpec(slide_id="title", title="Test proposal", archetype="Title"),
            SlideSpec(slide_id="content", title="A proposal point", archetype="Content", bullets=["A complete point"]),
        ])

        output = render_deck_from_template(
            deck,
            template_stream.getvalue(),
            customer_logo=logo,
        )
        rendered = Presentation(BytesIO(output))

        self.assertEqual(len(rendered.slides), 2)
        for slide in rendered.slides:
            logos = [shape for shape in slide.shapes if shape.name == "Customer Logo"]
            self.assertEqual(len(logos), 1)
            self.assertLess(logos[0].left + logos[0].width, rendered.slide_width)

    def test_compact_deck_plan_prompt_is_bounded(self) -> None:
        understanding = RFPUnderstanding(
            summary="A" * 3000,
            project_scope="B" * 3000,
            requirements=[
                Requirement(id=f"R-{i}", text="Requirement text " * 60, priority="must")
                for i in range(80)
            ],
            risks=["Risk text " * 40 for _ in range(20)],
            assumptions=["Assumption text " * 40 for _ in range(20)],
        )

        prompt = _compact_deck_plan_prompt(
            layout_names=[f"Layout {i}" for i in range(100)],
            understanding=understanding,
            narrative=None,
        )

        self.assertLess(len(prompt), 16000)
        self.assertIn("strict JSON", prompt)
        self.assertIn("top_requirements", prompt)

    def test_diagram_generator_uses_explicit_timeout_and_no_sdk_retries(self) -> None:
        png_b64 = base64.b64encode(b"fake-png").decode("ascii")
        client = SimpleNamespace(
            images=SimpleNamespace(
                generate=Mock(
                    return_value=SimpleNamespace(
                        data=[SimpleNamespace(b64_json=png_b64)]
                    )
                )
            )
        )
        settings = SimpleNamespace(
            image_model="gpt-image-2",
            image_timeout_s=321,
            image_retry_attempts=1,
        )

        with patch.object(diagram_generator, "settings", settings), patch.object(
            diagram_generator, "get_client", return_value=client
        ) as get_client:
            png = diagram_generator.generate_diagram_png(
                "draw architecture",
                out_path=None,
                model=None,
            )

        self.assertEqual(png, b"fake-png")
        get_client.assert_called_once_with(timeout=321.0, max_retries=0)
        client.images.generate.assert_called_once_with(
            model="gpt-image-2",
            prompt="draw architecture",
            size="auto",
            quality="auto",
        )

    def test_diagram_generator_retries_cloudflare_520_using_server_delay(self) -> None:
        png_b64 = base64.b64encode(b"retry-png").decode("ascii")
        transient = OpenAIError("Cloudflare 520")
        transient.status_code = 520
        transient.body = {"retryable": True, "retry_after": 60}
        success = SimpleNamespace(data=[SimpleNamespace(b64_json=png_b64)])
        client = SimpleNamespace(
            images=SimpleNamespace(generate=Mock(side_effect=[transient, success]))
        )
        settings = SimpleNamespace(
            image_model="gpt-image-2",
            image_timeout_s=300,
            image_retry_attempts=3,
            openai_retry_base_wait_s=5,
            openai_retry_max_wait_s=90,
        )

        with patch.object(diagram_generator, "settings", settings), patch.object(
            diagram_generator, "get_client", return_value=client
        ), patch.object(diagram_generator.time, "sleep") as sleep:
            png = diagram_generator.generate_diagram_png(
                "draw deployment architecture",
                out_path=None,
            )

        self.assertEqual(png, b"retry-png")
        self.assertEqual(client.images.generate.call_count, 2)
        sleep.assert_called_once_with(60.0)

    def test_risk_detailed_points_have_default_mitigations_when_no_risks_exist(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="Customer",
            opportunity_title="Opportunity",
            summary="Summary",
            project_scope="Scope",
            risks=[],
        )

        points = _risk_detailed_points(understanding)

        self.assertGreaterEqual(len(points), 3)
        self.assertTrue(all(point.sub_points for point in points))

    def test_understanding_risks_are_inferred_from_delivery_signals(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="Customer",
            opportunity_title="Data Hub",
            summary="Build a centralized data hub.",
            project_scope="Integrate source systems, extract data, validate reporting, and deploy to production.",
            in_scope_work=["API integration", "Data quality validation", "Production deployment"],
            risks=[],
        )

        enriched = enrich_understanding_risks(understanding)

        self.assertGreaterEqual(len(enriched.risks), 3)
        self.assertTrue(any("integration" in risk.lower() for risk in enriched.risks))
        self.assertTrue(any("data" in risk.lower() for risk in enriched.risks))
        self.assertTrue(any("deployment" in risk.lower() for risk in enriched.risks))

    def test_fallback_deck_plan_contains_core_storyline(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            opportunity_title="Catering Uplift Data Hub",
            summary="SATS needs a centralized operational data hub.",
            project_scope="Design and build the data hub.",
            requirements=[
                Requirement(id="R-1", text="Automate ELP extraction", priority="must")
            ],
        )

        deck = _fallback_deck_plan(
            title="Catering Uplift Data Hub",
            understanding=understanding,
            narrative=None,
        )
        archetypes = [slide.archetype for slide in deck.slides]

        self.assertIn("Customer Context", archetypes)
        self.assertIn("Architecture", archetypes)
        self.assertNotIn("Deployment Architecture", archetypes)
        self.assertIn("Next Steps", archetypes)

    def test_adaptive_skeleton_adds_data_integration_and_support_sections(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            opportunity_title="Catering Uplift Data Hub",
            summary="Build a data hub with ELP extraction, SAP and KSMS API integration, AMS support, and ICCMS migration.",
            project_scope="Data hub, reporting, migration, warranty and support.",
        )

        sections = _proposal_section_skeleton(understanding)
        ids = {section["slide_id"] for section in sections}

        self.assertIn("sk_integration", ids)
        self.assertIn("sk_technical_arch", ids)
        self.assertIn("sk_data_model", ids)
        self.assertIn("sk_reporting", ids)
        self.assertIn("sk_ai_opportunities", ids)
        self.assertIn("sk_migration", ids)
        self.assertIn("sk_ams", ids)
        data_model_section = next(
            section for section in sections if section["slide_id"] == "sk_data_model"
        )
        self.assertEqual(data_model_section["diagram_kind"], "data_model")
        technical_section = next(
            section for section in sections if section["slide_id"] == "sk_technical_arch"
        )
        self.assertEqual(technical_section["diagram_kind"], "technical_architecture")
        self.assertGreaterEqual(len(sections), 24)

    def test_managed_operations_skeleton_excludes_unsupported_sdlc_sections(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="Pearson",
            opportunity_title="Technology Operations Support",
            summary=(
                "Provide a managed technology operations service for the existing "
                "production application and infrastructure estate."
            ),
            project_scope=(
                "Deliver ITIL-aligned incident, problem and change management, monitoring, "
                "service levels, governance, staffing, transition, service reporting and "
                "continuous improvement. Include client references and a commercial proposal. "
                "Phase 2 expansion is optional."
            ),
            in_scope_work=[
                "Service transition and knowledge transfer",
                "Incident, problem and change management",
                "Operational monitoring and service reporting",
            ],
            out_of_scope_work=[
                "Application development",
                "Solution architecture",
                "Production deployment",
            ],
        )

        ids = {section["slide_id"] for section in _proposal_section_skeleton(understanding)}

        self.assertEqual(understanding.engagement_profile.primary_type, "managed_service_operations")
        self.assertTrue({
            "sk_operating_model", "sk_service_lifecycle", "sk_roadmap",
            "sk_service_measures", "sk_improvement", "sk_governance", "sk_staffing",
            "sk_expansion", "sk_references", "sk_commercials",
        }.issubset(ids))
        self.assertTrue({
            "sk_arch", "sk_technical_arch", "sk_integration", "sk_deployment",
            "sk_testing", "sk_tech", "sk_ai_opportunities",
        }.isdisjoint(ids))

    def test_application_development_skeleton_selects_technical_lifecycle_sections(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="Customer",
            opportunity_title="Digital Service Development",
            summary="Design and build a secure web application and API-based digital service.",
            project_scope=(
                "Software development includes UX and solution design, API integration, "
                "automated system testing, user acceptance testing, CI/CD and production deployment."
            ),
            in_scope_work=["Application development", "API integration", "Production deployment"],
        )

        ids = {section["slide_id"] for section in _proposal_section_skeleton(understanding)}

        self.assertEqual(understanding.engagement_profile.primary_type, "application_development")
        self.assertTrue({
            "sk_flow", "sk_arch", "sk_technical_arch", "sk_integration",
            "sk_deployment", "sk_testing", "sk_roadmap", "sk_tech",
        }.issubset(ids))
        self.assertTrue({"sk_operating_model", "sk_service_lifecycle", "sk_service_measures"}.isdisjoint(ids))

    def test_hybrid_profile_selects_only_the_declared_lifecycle_sections(self) -> None:
        understanding = RFPUnderstanding(
            summary="Operate the existing service and build a bounded customer application enhancement.",
            engagement_profile=EngagementProfile(
                primary_type="hybrid",
                secondary_types=["managed_service_operations", "application_development"],
                delivery_mode="hybrid",
                lifecycle_stages=[
                    LifecycleStageAssessment(stage="configure_build", in_scope=True, confidence=0.9),
                    LifecycleStageAssessment(stage="test_validate", in_scope=True, confidence=0.9),
                    LifecycleStageAssessment(stage="mobilize_transition", in_scope=True, confidence=0.9),
                    LifecycleStageAssessment(stage="operate_support", in_scope=True, confidence=0.9),
                ],
                confidence=0.9,
            ),
        )

        ids = {section["slide_id"] for section in _proposal_section_skeleton(understanding)}

        self.assertTrue({"sk_operating_model", "sk_service_lifecycle", "sk_arch", "sk_testing"}.issubset(ids))
        self.assertNotIn("sk_deployment", ids)

    def test_explicitly_unsupported_topics_override_classified_stages(self) -> None:
        understanding = RFPUnderstanding(
            summary="Design and build a web application.",
            engagement_profile=EngagementProfile(
                primary_type="application_development",
                delivery_mode="project_delivery",
                lifecycle_stages=[
                    LifecycleStageAssessment(stage="configure_build", in_scope=True, confidence=0.9),
                    LifecycleStageAssessment(stage="test_validate", in_scope=True, confidence=0.9),
                    LifecycleStageAssessment(stage="deploy_release", in_scope=True, confidence=0.9),
                ],
                explicitly_unsupported_topics=["testing strategy", "deployment architecture"],
                confidence=0.9,
            ),
        )

        ids = {section["slide_id"] for section in _proposal_section_skeleton(understanding)}

        self.assertIn("sk_arch", ids)
        self.assertNotIn("sk_testing", ids)
        self.assertNotIn("sk_deployment", ids)

    def test_profile_pruning_removes_sdlc_slides_from_managed_operations_plan(self) -> None:
        understanding = RFPUnderstanding(
            summary="Managed technology operations and service management for existing systems.",
            project_scope=(
                "Transition and operate the live service using incident, problem and change management, "
                "service levels, governance, operational reporting and continuous improvement."
            ),
        )
        deck = DeckPlan(
            deck_title="Operations proposal",
            slides=[
                SlideSpec(slide_id="sk_operating_model", title="Operating model", archetype="Solution Overview"),
                SlideSpec(slide_id="generated_arch", title="Solution architecture", archetype="Architecture"),
                SlideSpec(slide_id="generated_deploy", title="Deployment architecture", archetype="Deployment Architecture"),
                SlideSpec(slide_id="generated_test", title="Testing strategy", archetype="Delivery Plan"),
                SlideSpec(slide_id="generated_stack", title="Technology stack", archetype="Software Bill of Materials"),
            ],
        )

        pruned = prune_profile_misaligned_slides(deck, understanding)
        remaining = {slide.slide_id for slide in pruned.slides}

        self.assertEqual(remaining, {"sk_operating_model"})

    def test_ai_opportunities_are_grounded_in_data_hub_scope(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Centralize flight, ELP, productivity, accuracy and SLA data for reporting and AMS support.",
            project_scope="Ingest email and ELP files, validate operational data, report exceptions, and support uplift operations.",
        )

        opportunities = _ai_ml_opportunities(understanding)

        self.assertGreaterEqual(len(opportunities), 3)
        self.assertTrue(any("Anomaly" in item["name"] for item in opportunities))
        self.assertTrue(all("GPU" not in item["approach"] for item in opportunities))

    def test_assumptions_cover_four_proposal_control_categories(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Azure data hub with interfaces, UAT, cutover, warranty and AMS.",
            assumptions=["Representative source data will be available during mobilisation."],
            requirements=[
                Requirement(id="R-1", text="Customer provides interface access and sample ELP files.", priority="must"),
                Requirement(id="R-2", text="Confirm Azure hosting, identity and security controls.", priority="must"),
                Requirement(id="R-3", text="SATS owns UAT approval and production cutover sign-off.", priority="must"),
            ],
        )

        points = _assumptions_dependency_points(understanding)

        self.assertEqual(len(points), 4)
        self.assertEqual(
            [point.text for point in points],
            [
                "Scope and design baseline to validate",
                "Customer inputs and access dependencies",
                "Platform, security and environment decisions",
                "Acceptance and operational readiness dependencies",
            ],
        )
        self.assertTrue(all(point.sub_points for point in points))

    def test_platform_open_items_move_to_dependencies(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Deploy the digital catalogue on Azure.",
        )
        recommendations = TechnologyRecommendationSet(
            hosting_model="public-cloud",
            selected_platform="Microsoft Azure",
            platform_assumptions=[
                "SATS confirms the Azure landing-zone subscription and private-connectivity pattern.",
                "SATS confirms environment promotion approvals and production support ownership.",
            ],
            recommendations=[
                TechnologyRecommendation(
                    architecture_layer="Recovery objectives",
                    proposed_technology="Policy baseline",
                    technology_category="operational decision",
                    role="Agree recovery targets and evidence before production readiness",
                    status="customer-decision",
                    rationale="Business criticality drives recovery commitments",
                ),
            ],
        )

        visual_briefs = [
            DiagramBrief(
                slide_id="sk_deployment",
                visual_type="deployment",
                open_assumptions=[
                    "Runtime environment count and release approvals are not specified.",
                    "Enterprise integration targets and methods are not identified.",
                    "Backup frequency, retention, RPO and RTO are not specified.",
                    "Exact AWS hosting is not specified.",
                ],
            )
        ]
        points = _assumptions_dependency_points(
            understanding,
            recommendations,
            visual_briefs,
        )
        platform_points = next(
            point.sub_points
            for point in points
            if point.text == "Platform, security and environment decisions"
        )
        all_dependency_text = " ".join(
            item for point in points for item in point.sub_points
        )

        self.assertTrue(any("landing-zone" in item for item in platform_points))
        self.assertIn("Runtime environment count", all_dependency_text)
        self.assertIn("Enterprise integration targets", all_dependency_text)
        self.assertIn("Backup frequency", all_dependency_text)
        self.assertIn("Recovery objectives", all_dependency_text)
        self.assertNotIn("AWS", all_dependency_text)

    def test_data_domain_companion_content_is_authored_not_prompt_instructions(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Operational data domains for the catering uplift hub.",
            requirements=[
                Requirement(id="R-1", text="Ingest FIH flight and SQ SFTP ELP source records.", priority="must"),
                Requirement(id="R-2", text="Track catering uplift productivity, accuracy and SLA milestones.", priority="must"),
                Requirement(id="R-3", text="Validate and reconcile records with complete audit lineage.", priority="must"),
                Requirement(id="R-4", text="Publish governed Power BI operational reports.", priority="must"),
            ],
        )

        points = _data_domain_points(understanding)

        self.assertEqual(len(points), 4)
        self.assertTrue(any("FIH" in " ".join(point.sub_points) for point in points))
        self.assertTrue(any("Power BI" in " ".join(point.sub_points) for point in points))

    def test_required_slides_add_ai_opportunity_and_cost_controls(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Operational data hub for flights, ELP files, reporting, SLA exceptions, and support.",
            project_scope="Ingest email and files, validate data, support forecasting, reporting, and AMS.",
        )
        deck = DeckPlan(
            deck_title="Test",
            slides=[SlideSpec(slide_id="title", title="Test", archetype="Title")],
        )

        enriched = ensure_required_slides(deck, understanding=understanding)
        ai_slides = [slide for slide in enriched.slides if "ai-assisted" in slide.title.lower()]

        self.assertEqual(len(ai_slides), 1)
        self.assertGreaterEqual(len(ai_slides[0].cards), 3)

    def test_incomplete_technology_table_is_replaced_without_inventing_ai_product(self) -> None:
        understanding = RFPUnderstanding(
            summary="Data hub for flight files, SLA exceptions, reporting and support.",
            project_scope="Validate operational data, forecast demand, and support analytics.",
        )
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="tech",
                    title="Technology stack",
                    archetype="Software Bill of Materials",
                    table={"headers": ["Phase", "Toolset", "Purpose", "Basis"], "rows": [["Build", "TBC", "Build", "Scope"]]},
                )
            ],
        )

        enriched = ensure_required_slides(deck, understanding=understanding)
        rows = next(slide.table["rows"] for slide in enriched.slides if slide.slide_id == "tech")

        self.assertTrue(any("No product selected by fallback" in row for row in rows))
        self.assertFalse(any("Azure AI" in " ".join(row) or "Bedrock" in " ".join(row) for row in rows))

    def test_power_bi_does_not_imply_fabric_or_azure_stack(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Central catering uplift data hub with reporting and SLA analytics.",
            project_scope="Ingest ELP files, validate canonical data, integrate APIs, and serve Power BI.",
            solution_technologies=["Power BI", "Microsoft Entra ID", "Microsoft Sentinel", "SAP Ariba"],
            procurement_or_submission_tools=["SAP Ariba"],
        )

        table = _source_grounded_technology_table(understanding)
        text = " ".join(str(cell) for row in table["rows"] for cell in row)

        self.assertIn("Power BI", text)
        self.assertNotIn("Fabric Data Factory pipelines", text)
        self.assertNotIn("Microsoft Fabric OneLake", text)
        self.assertNotIn("Azure Functions", text)
        self.assertNotIn("Ariba", text)
        self.assertEqual(table["headers"][0], "Source classification")

    def test_procurement_tools_are_excluded_from_sbom_and_diagrams(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Azure-hosted operational data hub with Power BI reporting.",
            project_scope="Ingest, validate, store, and report operational data.",
            solution_technologies=["Power BI", "SAP Ariba"],
            procurement_or_submission_tools=["Ariba"],
        )

        sbom_text = " ".join(str(cell) for row in _sbom_table(understanding)["rows"] for cell in row)
        diagram_prompt = _build_diagram_prompt("architecture", understanding)

        self.assertNotIn("Ariba", sbom_text)
        self.assertNotIn("Ariba", diagram_prompt)
        self.assertIn("Power BI", diagram_prompt)
        self.assertNotIn("Fabric Data Factory pipelines", diagram_prompt)

    def test_explicit_aws_signal_does_not_propose_fabric(self) -> None:
        understanding = RFPUnderstanding(
            summary="AWS-hosted operational data platform.",
            project_scope="Ingest files, expose APIs, store curated data, and provide reporting.",
            solution_technologies=["AWS", "Amazon S3"],
        )

        text = " ".join(str(cell) for row in _source_grounded_technology_table(understanding)["rows"] for cell in row)

        self.assertIn("AWS", text)
        self.assertIn("Amazon S3", text)
        self.assertNotIn("AWS Glue", text)
        self.assertNotIn("Microsoft Fabric", text)

    def test_visual_briefs_do_not_cross_route_between_slide_types(self) -> None:
        slides = [
            SlideSpec(slide_id="sk_exec", title="Executive Summary", archetype="Solution Overview", diagram=DiagramSpec(kind="architecture", prompt="wrong")),
            SlideSpec(slide_id="sk_technical_arch", title="Layered technical architecture connects systems and products", archetype="Architecture"),
            SlideSpec(slide_id="sk_data_model", title="Core data domains and ownership", archetype="Content"),
            SlideSpec(slide_id="sk_deployment", title="Deployment and resilience protect operations", archetype="Deployment Architecture"),
            SlideSpec(slide_id="auto_hr_dr", title="HA and DR protect continuity", archetype="High Availability & DR"),
            SlideSpec(slide_id="sk_testing", title="Acceptance evidence proves readiness", archetype="Delivery Plan"),
            SlideSpec(slide_id="sk_ams", title="AMS protects the live service", archetype="Delivery Plan"),
        ]
        briefs = [
            DiagramBrief(slide_id="architecture", visual_type="architecture", entities=["App", "API", "DB", "User"], flows=["User -> App", "App -> DB"]),
            DiagramBrief(slide_id="sk_technical_arch", visual_type="technical_architecture", entities=["Source systems", "COTS products", "Custom services", "Data services"], flows=["Sources -> Integration", "Integration -> Services"], controls=["Security"], evidence_refs=["R-1"]),
            DiagramBrief(slide_id="sk_data_model", visual_type="data_model", entities=["Catalogue", "Customer", "Solution", "Validation"], flows=["Catalogue -> Solution", "Customer -> Solution"], controls=["Ownership", "Quality"], evidence_refs=["R-1"]),
            DiagramBrief(slide_id="sk_deployment", visual_type="deployment", entities=["Build", "UAT", "Production", "DR"], flows=["Build -> UAT", "UAT -> Production"]),
            DiagramBrief(slide_id="auto_hr_dr", visual_type="hadr", entities=["Primary", "Standby", "Backup", "Monitor"], flows=["Primary -> Standby", "Primary -> Backup"]),
            DiagramBrief(slide_id="sk_testing", visual_type="testing", entities=["API tests", "Data reconciliation", "UAT", "Evidence"], flows=["Tests -> Evidence", "Evidence -> UAT"]),
            DiagramBrief(slide_id="sk_ams", visual_type="ams", entities=["Telemetry", "Alert", "Runbook", "Resolver"], flows=["Telemetry -> Alert", "Alert -> Resolver"]),
        ]
        understanding = RFPUnderstanding(
            summary="Design and build a governed data platform application with warranty and AMS support.",
            project_scope=(
                "Application development includes API integration, system testing, production deployment, "
                "backup and disaster recovery, governance, reporting and live-service support."
            ),
        )
        deck = ensure_diagrams_for_key_slides(
            DeckPlan(deck_title="Test", slides=slides), understanding, briefs
        )
        by_id = {slide.slide_id: slide for slide in deck.slides}
        self.assertIsNone(by_id["sk_exec"].diagram)
        self.assertEqual(by_id["sk_technical_arch"].diagram.kind, "technical_architecture")
        self.assertEqual(by_id["sk_data_model"].diagram.kind, "data_model")
        self.assertEqual(by_id["sk_deployment"].diagram.kind, "deployment")
        self.assertEqual(by_id["auto_hr_dr"].diagram.kind, "hadr")
        self.assertEqual(by_id["sk_testing"].diagram.kind, "testing")
        self.assertIsNone(by_id["sk_ams"].diagram)

    def test_required_visual_sections_receive_grounded_supplemental_briefs(self) -> None:
        understanding = RFPUnderstanding(
            summary="Design and build a governed catalogue application and data platform.",
            project_scope=(
                "Application development includes source API integration, system testing, production deployment, "
                "governance, backup, disaster recovery and live-service support."
            ),
            in_scope_work=["Catalogue consolidation", "API integration", "Search and approved publishing"],
            requirements=[
                Requirement(id="R-1", text="Integrate source APIs and files with validation and audit."),
                Requirement(id="R-2", text="Provide backup, disaster recovery, monitoring and controlled access."),
            ],
        )
        slides = [
            SlideSpec(slide_id="sk_solution", title="Proposed solution at a glance", archetype="Solution Overview"),
            SlideSpec(slide_id="sk_technical_arch", title="Layered technical architecture connects systems and products", archetype="Architecture"),
            SlideSpec(slide_id="sk_integration", title="Integration architecture connects source and consumer systems", archetype="Architecture"),
            SlideSpec(slide_id="sk_data_model", title="Core data domains and ownership", archetype="Content"),
            SlideSpec(slide_id="sk_deployment", title="Deployment and resilience protect operations", archetype="Deployment Architecture"),
            SlideSpec(slide_id="auto_ha_and_dr_protect_business_continuity", title="HA and DR protect business continuity", archetype="High Availability & DR"),
            SlideSpec(slide_id="sk_roadmap", title="Agile roadmap releases value through increments", archetype="Timeline"),
            SlideSpec(slide_id="sk_governance", title="Product-aligned squads combine business and engineering ownership", archetype="Team"),
        ]
        planned = ensure_diagrams_for_key_slides(
            DeckPlan(deck_title="Test", slides=slides),
            understanding,
            _fallback_visual_briefs(understanding),
        )
        expected = {
            "sk_technical_arch": "technical_architecture",
            "sk_integration": "architecture",
            "sk_data_model": "data_model",
            "sk_deployment": "deployment",
            "auto_ha_and_dr_protect_business_continuity": "hadr",
            "sk_roadmap": "timeline",
            "sk_governance": "org",
        }
        for slide in planned.slides:
            if slide.slide_id == "sk_solution":
                self.assertIsNone(slide.diagram)
                continue
            self.assertIsNotNone(slide.diagram)
            self.assertEqual(slide.diagram.kind, expected[slide.slide_id])
            self.assertGreaterEqual(slide.diagram.grounding_score, 0.45)

    def test_data_model_content_slide_gets_grounded_fallback_diagram(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Build a governed digital catalogue with shared data ownership.",
            project_scope="Consolidate product and service catalogues, customer briefs, validation outcomes and pricing decisions.",
            requirements=[
                Requirement(id="R-1", text="Consolidate product, service and SKU catalogue records."),
                Requirement(id="R-2", text="Capture customer requirements and solution shortlists."),
                Requirement(id="R-3", text="Retain compliance validation and pricing decision evidence."),
                Requirement(id="R-4", text="Assign domain ownership, stewardship, lineage and audit controls."),
            ],
        )
        deck = DeckPlan(deck_title="Test", slides=[
            SlideSpec(
                slide_id="sk_data_model",
                title="Core data domains and ownership",
                archetype="Content",
            )
        ])

        enrich_slide_detail(deck, understanding)
        ensure_diagrams_for_key_slides(deck, understanding, visual_briefs=None)

        slide = deck.slides[0]
        self.assertIsNotNone(slide.diagram)
        self.assertEqual(slide.diagram.kind, "data_model")
        self.assertGreaterEqual(slide.diagram.grounding_score, 0.45)
        self.assertIn("conceptual core data-domain and ownership map", slide.diagram.prompt)
        self.assertIn("Grounded domain requirements", slide.diagram.prompt)
        self.assertTrue(slide.detailed_points)

    def test_layered_technical_architecture_uses_build_buy_and_data_source_decisions(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Build a governed digital catalogue on Microsoft Azure.",
            project_scope=(
                "Consolidate product data and approved images, optimise pricing, integrate warehouse "
                "availability, and publish customer-facing catalogues."
            ),
            requirements=[
                Requirement(id="R-1", text="Master product, SKU, dietary and packaging attributes."),
                Requirement(id="R-2", text="Store approved product documents and images."),
                Requirement(id="R-3", text="Use raw-material cost and market inputs for pricing optimisation."),
                Requirement(id="R-4", text="Integrate cold and ambient warehouse inventory availability."),
            ],
        )
        recommendations = TechnologyRecommendationSet(
            hosting_model="public-cloud",
            selected_platform="Microsoft Azure",
            component_decisions=[
                SolutionComponentDecision(
                    capability="Product master and catalogue stewardship",
                    recommendation="TIBCO EBX",
                    sourcing_model="COTS/SaaS",
                    role="Govern product hierarchy, attributes and stewardship workflow",
                    system_of_record="Regional product owners and approved enterprise sources",
                    data_inputs=["Product and SKU attributes", "Dietary and packaging attributes"],
                    data_outputs=["Mastered catalogue records"],
                    decision_status="recommended",
                    rationale="Configurable mastering and stewardship avoid rebuilding commodity MDM controls",
                    alternatives_considered=["Informatica MDM", "Semarchy xDM"],
                ),
                SolutionComponentDecision(
                    capability="Approved digital assets",
                    recommendation="Azure Data Lake Storage Gen2",
                    sourcing_model="managed-cloud",
                    role="Store documents and images while catalogue records retain governed links and metadata",
                    system_of_record="Approved content owners",
                    data_inputs=["Product images", "Technical documents"],
                    data_outputs=["Versioned asset URIs and metadata"],
                    decision_status="recommended",
                    rationale="Object storage fits binary assets better than the master-data repository",
                ),
                SolutionComponentDecision(
                    capability="Pricing optimisation",
                    recommendation="Custom pricing optimisation service",
                    sourcing_model="custom-build",
                    role="Apply proposal-specific pricing constraints and approval logic",
                    system_of_record="ERP or procurement pricing source",
                    data_inputs=["Raw-material costs", "Market and customer constraints"],
                    data_outputs=["Recommended price and decision evidence"],
                    decision_status="recommended",
                    rationale="The differentiated optimisation and approval rules are not catalogue-master functions",
                    alternatives_considered=["Dedicated price optimisation SaaS", "ERP pricing extension"],
                ),
                SolutionComponentDecision(
                    capability="Warehouse availability",
                    recommendation="Warehouse Management System integration",
                    sourcing_model="integration-only",
                    role="Reuse cold and ambient inventory without duplicating stock ownership",
                    system_of_record="Existing warehouse management systems",
                    data_inputs=["Inventory and facility availability"],
                    data_outputs=["Availability indicators"],
                    decision_status="recommended",
                    rationale="The catalogue consumes availability while the WMS remains authoritative",
                ),
                SolutionComponentDecision(
                    capability="Optional enterprise DAM",
                    recommendation="Optional DAM vendor",
                    sourcing_model="customer-decision",
                    role="Potential future enterprise media governance",
                    decision_status="customer-decision",
                ),
            ],
            recommendations=[
                TechnologyRecommendation(
                    architecture_layer="Object storage",
                    proposed_technology="Azure Data Lake Storage Gen2",
                    technology_category="object store",
                    role="Store approved documents and images",
                    status="recommended",
                    rationale="Managed Azure storage fit",
                    sourcing_model="managed-cloud",
                    build_vs_buy_rationale="Use managed object storage instead of building binary storage",
                )
            ],
        )
        deck = DeckPlan(deck_title="Test", slides=[
            SlideSpec(
                slide_id="sk_technical_arch",
                title="Layered technical architecture",
                archetype="Architecture",
            )
        ])

        enrich_slide_detail(deck, understanding, recommendations)
        ensure_diagrams_for_key_slides(
            deck,
            understanding,
            visual_briefs=None,
            technology_recommendations=recommendations,
        )

        diagram = deck.slides[0].diagram
        self.assertIsNotNone(diagram)
        self.assertEqual(diagram.kind, "technical_architecture")
        self.assertGreaterEqual(diagram.grounding_score, 0.45)
        self.assertIn("layered technical architecture", diagram.prompt.lower())
        self.assertIn("TIBCO EBX", diagram.prompt)
        self.assertIn("Azure Data Lake Storage Gen2", diagram.prompt)
        self.assertIn("Warehouse Management System integration", diagram.prompt)
        self.assertIn("Raw-material costs", diagram.prompt)
        self.assertNotIn("Optional DAM vendor", diagram.prompt)

    def test_technology_recommendations_require_concrete_products(self) -> None:
        recommendations = TechnologyRecommendationSet(recommendations=[
            TechnologyRecommendation(architecture_layer="API/backend", proposed_technology="Java 21 with Spring Boot", technology_category="API framework", role="Implement catalogue APIs", status="recommended", rationale="Strong transactional and integration fit", evidence_refs=["PARAGRAPH 79"], alternatives_considered=[".NET 8", "Python FastAPI"]),
            TechnologyRecommendation(architecture_layer="Data store", proposed_technology="PostgreSQL", technology_category="relational SQL database", role="Store governed catalogue and workflow data", status="recommended", rationale="Relational integrity and flexible JSON support", alternatives_considered=["SQL Server", "MongoDB"]),
            TechnologyRecommendation(architecture_layer="Search", proposed_technology="OpenSearch", technology_category="search index", role="Faceted catalogue search", status="recommended", rationale="Search and filtering fit", alternatives_considered=["Elasticsearch"]),
            TechnologyRecommendation(architecture_layer="Testing", proposed_technology="JUnit 5; REST Assured; Playwright", technology_category="automated test toolchain", role="Automate unit, API, and UI acceptance", status="recommended", rationale="Matches the proposed implementation layers", alternatives_considered=["Cypress"]),
        ])
        table = _technology_recommendation_table(recommendations)
        text = " ".join(str(cell) for row in table["rows"] for cell in row)
        self.assertIn("PostgreSQL", text)
        self.assertIn("Java 21", text)
        self.assertNotIn("FS Digital Catalogue", text)
        self.assertNotIn("PARAGRAPH", text)
        self.assertNotIn("Alternatives considered", text)

    def test_customer_preferred_platform_overrides_conflicting_provider_draft(self) -> None:
        recommendations = TechnologyRecommendationSet(
            hosting_model="public-cloud",
            selected_platform="Amazon Web Services",
            deployment_rationale="Use AWS managed services",
            primary_region_strategy="AWS multi-region deployment",
            platform_assumptions=["Confirm the AWS landing zone."],
            component_decisions=[
                SolutionComponentDecision(
                    capability="Object storage",
                    recommendation="Amazon S3",
                    sourcing_model="managed-cloud",
                    role="Store catalogue assets on AWS",
                    decision_status="recommended",
                )
            ],
            recommendations=[
                TechnologyRecommendation(architecture_layer="Application runtime", proposed_technology="Amazon ECS on AWS Fargate", technology_category="managed container runtime", role="Run application services", status="recommended", rationale="Managed runtime"),
                TechnologyRecommendation(architecture_layer="Application UI", proposed_technology="React 19", technology_category="web framework", role="Build the user interface", status="recommended", rationale="Application fit"),
            ],
        )

        aligned = _align_recommendations_to_customer_platform(
            recommendations,
            {"platform": "Microsoft Azure", "status": "Customer-preferred"},
        )

        self.assertEqual(aligned.selected_platform, "Microsoft Azure")
        self.assertTrue(all("AWS" not in item.proposed_technology for item in aligned.recommendations))
        self.assertEqual([item.proposed_technology for item in aligned.recommendations], ["React 19"])
        self.assertEqual(aligned.component_decisions, [])
        self.assertNotIn("AWS", aligned.deployment_rationale)
        self.assertEqual(aligned.primary_region_strategy, "")
        self.assertEqual(aligned.platform_assumptions, [])

    def test_customer_platform_overrides_solution_brief_cloud_inference(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            opportunity_title="FS Digital Catalogue",
            summary="Deploy the catalogue on AWS using Lambda services.",
            solution_technologies=["Amazon Web Services", "AWS Lambda"],
        )
        customer_context = {
            "platform": "Microsoft Azure",
            "status": "Customer-preferred",
            "details": "Use the approved Azure estate.",
        }

        brief = build_solution_brief(
            understanding,
            None,
            customer_context,
        )
        plan_input = json.loads(
            _chunked_plan_input(understanding, None, customer_context)
        )

        self.assertEqual(brief.target_cloud, "azure")
        self.assertEqual(plan_input["solution_brief"]["target_cloud"], "azure")
        self.assertEqual(
            plan_input["customer_technology_context"]["platform"],
            "Microsoft Azure",
        )

    def test_supporting_reference_survives_direct_path_as_advisory_context(self) -> None:
        reference = SourceDocument(
            document_id="doc-architecture-reference",
            name="Catalogue Architecture Research.docx",
            document_type="supporting_reference",
            authority="contextual",
            text=(
                "Use an Experience Layer, API Layer and Business Service Layer. "
                "Consider TIBCO MDM, Azure API Management and Data Lake Storage Gen2."
            ),
        )
        state = AgentState(
            rfp_text="Short authoritative RFP plus supporting material",
            template_info={},
            source_documents=[reference],
        )
        fake_settings = SimpleNamespace(
            understanding_direct_max_chars=180000,
            contextual_reference_max_chars=18000,
        )

        with patch("rfp2deck.agent.nodes.settings", fake_settings):
            result = extract_source_evidence(state)

        self.assertIn("ADVISORY SUPPORTING REFERENCE CONTEXT", result["contextual_reference_context"])
        self.assertIn("TIBCO MDM", state.contextual_reference_context)
        self.assertIsNone(result["evidence_text"])

    def test_fallback_does_not_invent_azure_stack_or_region_pair(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="A Singapore digital catalogue integrates pricing, inventory, documents and mastered product data.",
            project_scope="Provide catalogue search, workflows, reporting and Azure deployment.",
            solution_technologies=["Customer Catalogue SDK"],
        )
        recommendations = _source_grounded_technology_fallback(
            understanding,
            {
                "platform": "Microsoft Azure",
                "status": "Existing estate",
                "details": "Singapore deployment with customer landing-zone controls.",
            },
            "Advisory research discusses TIBCO MDM as one product option. "
            "Primary region is Advisory Alpha and the recovery region is Advisory Beta.",
        )
        technology_text = " ".join(
            item.proposed_technology for item in recommendations.recommendations
        )
        self.assertIn("Customer Catalogue SDK", technology_text)
        self.assertNotIn("React", technology_text)
        self.assertNotIn("ASP.NET Core", technology_text)
        self.assertNotIn("Azure API Management", technology_text)
        self.assertNotIn("TIBCO MDM", technology_text)
        self.assertEqual(recommendations.component_decisions, [])
        self.assertEqual(recommendations.selected_platform, "Microsoft Azure")
        self.assertEqual(recommendations.primary_region_strategy, "")

    def test_technology_node_failure_preserves_sources_without_default_stack(self) -> None:
        state = AgentState(
            rfp_text="",
            template_info={},
            understanding=RFPUnderstanding(
                summary="Singapore digital catalogue on the existing Azure estate.",
                project_scope="Build web, API, data, reporting and deployment capabilities.",
                solution_technologies=["Customer UI Framework"],
            ),
            customer_technology_context={
                "platform": "Microsoft Azure",
                "status": "Existing estate",
                "details": "Singapore deployment.",
            },
        )
        fake_settings = SimpleNamespace(
            reasoning_effort_medium="medium",
            deck_plan_timeout_s=30,
        )

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema",
            side_effect=RuntimeError("proxy timeout"),
        ):
            result = derive_technology_recommendations(state)

        technologies = " ".join(
            item.proposed_technology
            for item in result["technology_recommendations"].recommendations
        )
        self.assertIn("Customer UI Framework", technologies)
        self.assertNotIn("React", technologies)
        self.assertNotIn("ASP.NET Core", technologies)
        self.assertNotIn("Azure Container Apps", technologies)
        self.assertEqual(result["technology_recommendations"].selected_platform, "Microsoft Azure")
        self.assertEqual(result["technology_recommendations"].primary_region_strategy, "")

    def test_incomplete_agent_stack_is_rederived_without_a_fixed_fallback(self) -> None:
        state = AgentState(
            rfp_text="",
            template_info={},
            understanding=RFPUnderstanding(
                summary="Software development will build a customer web portal with APIs, search and governed catalogue data.",
                project_scope=(
                    "Design and build the web application, including API integration, data, "
                    "system testing and production deployment capabilities."
                ),
            ),
            customer_technology_context={"platform": "Customer Cloud", "status": "Customer-mandated"},
        )
        initial = TechnologyRecommendationSet(
            recommendations=[
                TechnologyRecommendation(
                    architecture_layer="Portal",
                    proposed_technology="Customer Web Runtime",
                    technology_category="web runtime",
                    role="Run the portal",
                    status="recommended",
                    rationale="Initial proposal-specific choice",
                )
            ]
        )
        repaired = TechnologyRecommendationSet(
            selected_platform="Customer Cloud",
            recommendations=[
                TechnologyRecommendation(architecture_layer="Portal", proposed_technology="Customer Web Runtime", technology_category="web UI", role="Render journeys", status="recommended", rationale="Customer standard"),
                TechnologyRecommendation(architecture_layer="Application", proposed_technology="Product Service Runtime", technology_category="application framework", role="Run APIs", status="recommended", rationale="Workload fit"),
                TechnologyRecommendation(architecture_layer="Exchange", proposed_technology="Partner Exchange Gateway", technology_category="integration", role="Exchange partner data", status="recommended", rationale="Interface fit"),
                TechnologyRecommendation(architecture_layer="Information", proposed_technology="Catalogue Record Store", technology_category="database", role="Persist governed records", status="recommended", rationale="Data fit"),
                TechnologyRecommendation(architecture_layer="Delivery", proposed_technology="Customer Delivery Toolchain", technology_category="CI/CD and test", role="Build, test and deploy", status="recommended", rationale="Operating-model fit"),
            ],
        )
        fake_settings = SimpleNamespace(reasoning_effort_medium="medium", deck_plan_timeout_s=30)

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema",
            side_effect=[initial, repaired],
        ) as structured_call:
            result = derive_technology_recommendations(state)

        self.assertEqual(structured_call.call_count, 2)
        technologies = {
            item.proposed_technology
            for item in result["technology_recommendations"].recommendations
        }
        self.assertIn("Partner Exchange Gateway", technologies)
        self.assertIn("Customer Delivery Toolchain", technologies)
        self.assertEqual(result["technology_recommendations"].selected_platform, "Customer Cloud")

    def test_explicit_primary_and_recovery_regions_are_preserved_verbatim(self) -> None:
        recommendations = _source_grounded_technology_fallback(
            RFPUnderstanding(summary="Customer-hosted application."),
            {
                "platform": "Customer Cloud",
                "status": "Customer-mandated",
                "details": "Primary region is Region Alpha and the recovery region is Region Beta.",
            },
        )

        self.assertIn("Region Alpha", recommendations.primary_region_strategy)
        self.assertIn("Region Beta", recommendations.primary_region_strategy)

    def test_required_sections_include_reporting_testing_and_roadmap_explanation(self) -> None:
        understanding = RFPUnderstanding(
            summary="Digital catalogue data platform with governance, reporting, testing and incremental delivery.",
            project_scope="Build a data platform that integrates catalogue data, APIs and dashboards, then deploy and support the service.",
        )
        sections = {item["slide_id"]: item for item in _proposal_section_skeleton(understanding)}
        briefs = {item.slide_id: item for item in _fallback_visual_briefs(understanding)}

        self.assertEqual(sections["sk_reporting"]["diagram_kind"], "process")
        self.assertEqual(sections["sk_testing"]["diagram_kind"], "testing")
        self.assertIn("sk_roadmap_detail", sections)
        self.assertEqual(briefs["sk_reporting"].visual_type, "process")
        self.assertEqual(briefs["sk_testing"].visual_type, "testing")

    def test_fallback_technical_architecture_uses_source_content_not_fixed_layers(self) -> None:
        integration_entities, integration_flows = _source_grounded_technical_architecture_elements(
            RFPUnderstanding(
                summary="Exchange partner files through SFTP with validation and audit.",
                project_scope="Integrate external supplier interfaces and monitored file processing.",
                in_scope_work=["Supplier SFTP exchange", "Catalogue validation service"],
                requirements=[
                    Requirement(id="R-1", text="Supplier files enter through SFTP and pass validation before catalogue publication."),
                ],
            )
        )
        application_entities, application_flows = _source_grounded_technical_architecture_elements(
            RFPUnderstanding(
                summary="Build a mobile and web workflow application with APIs and a database.",
                project_scope="Deliver user journeys, business services, persistent data and cloud deployment.",
                in_scope_work=["Mobile approval journey", "Product authoring workflow"],
                solution_technologies=["Customer Design System"],
                requirements=[
                    Requirement(id="R-2", text="The application API stores approved product records in the customer database."),
                ],
            )
        )
        integration_text = " ".join(integration_entities + integration_flows)
        application_text = " ".join(application_entities + application_flows)

        self.assertIn("Supplier SFTP exchange", integration_text)
        self.assertIn("Supplier files enter through SFTP", integration_text)
        self.assertIn("Customer Design System", application_text)
        self.assertIn("Mobile approval journey", application_text)
        for fixed_layer in ("Experience Layer", "API Layer", "Business Service Layer", "Data Layer", "Cloud Layer"):
            self.assertNotIn(fixed_layer, integration_text)
            self.assertNotIn(fixed_layer, application_text)
        self.assertNotEqual(integration_entities, application_entities)

    def test_speaker_notes_are_generated_in_small_batches_with_visual_context(self) -> None:
        slides = [
            SlideSpec(
                slide_id=f"slide_{index}",
                title=f"Decision {index}",
                archetype="Architecture" if index == 0 else "Content",
                bullets=[f"Explain decision {index}"],
                diagram=(
                    DiagramSpec(kind="architecture", prompt="Show React, .NET and Azure APIs")
                    if index == 0 else None
                ),
            )
            for index in range(13)
        ]
        state = AgentState(
            rfp_text="",
            template_info={},
            understanding=RFPUnderstanding(summary="A governed Azure catalogue."),
            deck_plan=DeckPlan(deck_title="Test", slides=slides),
            technology_recommendations=TechnologyRecommendationSet(selected_platform="Microsoft Azure"),
        )
        fake_settings = SimpleNamespace(
            notes_batch_size=6,
            notes_workers=3,
            model_fast="test-fast-model",
            reasoning_effort_low="low",
        )

        def fake_notes(prompt, schema, **kwargs):
            payload = prompt.split("SLIDE BATCH (JSON; includes neighbouring titles, content, visual and table):\n", 1)[1]
            payload = payload.split("\n\nEXECUTIVE NARRATIVE SPINE", 1)[0]
            batch = json.loads(payload)
            return DeckNotes(notes=[
                SlideNote(slide_id=item["slide_id"], notes=f"Narrative for {item['title']} with design rationale and transition.")
                for item in batch
            ])

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema",
            side_effect=fake_notes,
        ) as structured_call:
            generate_notes(state)

        self.assertEqual(structured_call.call_count, 3)
        self.assertTrue(all((slide.notes or "").startswith("Narrative for") for slide in slides))
        first_prompt = structured_call.call_args_list[0].args[0]
        self.assertIn("Show React, .NET and Azure APIs", " ".join(call.args[0] for call in structured_call.call_args_list))
        self.assertIn("Microsoft Azure", first_prompt)

    def test_full_deck_planning_receives_customer_platform_context(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            opportunity_title="FS Digital Catalogue",
            summary="An earlier draft mentions an AWS target runtime.",
        )
        customer_context = {
            "platform": "Microsoft Azure",
            "status": "Customer-preferred",
            "details": "Use the approved Azure estate.",
        }
        state = AgentState(
            rfp_text="",
            template_info={},
            understanding=understanding,
            customer_technology_context=customer_context,
        )
        fake_settings = SimpleNamespace(
            deck_plan_specialists=False,
            deck_plan_chunked=False,
            deck_plan_prompt_max_chars=100000,
            deck_plan_rag_max_chars=18000,
            reasoning_effort_deck_plan="medium",
            deck_plan_timeout_s=30,
        )
        generated = DeckPlan(
            deck_title="FS Digital Catalogue",
            slides=[SlideSpec(slide_id="title", title="FS Digital Catalogue")],
        )

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema",
            return_value=generated,
        ) as structured_call, patch(
            "rfp2deck.agent.nodes._post_process_deck_plan",
            side_effect=lambda deck_plan, **_: deck_plan,
        ) as post_process:
            plan_deck(state)

        prompt = structured_call.call_args.args[0]
        self.assertIn('"platform":"Microsoft Azure"', prompt)
        self.assertEqual(
            post_process.call_args.kwargs["customer_technology_context"],
            customer_context,
        )

    def test_deployment_diagram_uses_selected_platform_services(self) -> None:
        recommendations = TechnologyRecommendationSet(
            hosting_model="public-cloud",
            selected_platform="Amazon Web Services",
            deployment_rationale="Existing AWS operating model and managed-service fit",
            primary_region_strategy="Multi-AZ primary region with a separate recovery region",
            recommendations=[
                TechnologyRecommendation(architecture_layer="Edge and ingress", proposed_technology="Amazon Route 53; AWS WAF; Application Load Balancer", technology_category="network edge", role="Secure and route inbound traffic", status="recommended", rationale="Managed ingress"),
                TechnologyRecommendation(architecture_layer="Application runtime", proposed_technology="Amazon ECS on AWS Fargate", technology_category="managed container runtime", role="Run application and API services", status="recommended", rationale="Low operations overhead"),
                TechnologyRecommendation(architecture_layer="Data store", proposed_technology="Amazon Aurora PostgreSQL", technology_category="relational SQL database", role="Store transactional catalogue data", status="recommended", rationale="Relational integrity"),
                TechnologyRecommendation(architecture_layer="Identity and secrets", proposed_technology="Amazon Cognito; AWS Secrets Manager; AWS KMS", technology_category="identity and security", role="Protect identities, secrets, and keys", status="recommended", rationale="Managed security controls"),
            ],
        )
        deck = DeckPlan(deck_title="Test", slides=[
            SlideSpec(slide_id="deploy", title="Deployment architecture", archetype="Deployment Architecture")
        ])
        understanding = RFPUnderstanding(
            summary="Design and build a secure digital catalogue application.",
            project_scope="Application development includes production deployment and release automation.",
        )

        enrich_slide_detail(deck, understanding, recommendations)
        ensure_diagrams_for_key_slides(
            deck,
            understanding,
            _fallback_visual_briefs(understanding),
            recommendations,
        )

        prompt = deck.slides[0].diagram.prompt
        self.assertIn("Amazon Web Services", prompt)
        self.assertIn("Amazon ECS on AWS Fargate", prompt)
        self.assertIn("AWS WAF", prompt)
        self.assertTrue(any("Amazon Web Services" in item for item in deck.slides[0].bullets))

    def test_deployment_diagram_replaces_stale_provider_and_visible_open_items(self) -> None:
        recommendations = TechnologyRecommendationSet(
            hosting_model="public-cloud",
            selected_platform="Microsoft Azure",
            deployment_rationale="Customer-preferred cloud and existing operating model",
            primary_region_strategy="Azure Southeast Asia with zone redundancy",
            platform_assumptions=["SATS confirms landing-zone connectivity and release approvals."],
            recommendations=[
                TechnologyRecommendation(architecture_layer="Edge and ingress", proposed_technology="Azure Front Door Premium; Azure Web Application Firewall", technology_category="network edge", role="Secure and route inbound traffic", status="recommended", rationale="Managed Azure ingress"),
                TechnologyRecommendation(architecture_layer="Application runtime", proposed_technology="Azure App Service", technology_category="managed application runtime", role="Run application and API services", status="recommended", rationale="Managed runtime"),
                TechnologyRecommendation(architecture_layer="Data store", proposed_technology="Azure Database for PostgreSQL Flexible Server", technology_category="relational SQL database", role="Store transactional catalogue data", status="recommended", rationale="Relational integrity"),
                TechnologyRecommendation(architecture_layer="Identity and secrets", proposed_technology="Microsoft Entra ID; Azure Key Vault", technology_category="identity and security", role="Protect identities, secrets, and keys", status="recommended", rationale="Azure-native controls"),
                TechnologyRecommendation(architecture_layer="Conflicting runtime", proposed_technology="Amazon ECS on AWS Fargate", technology_category="managed container runtime", role="Run a legacy draft runtime", status="recommended", rationale="Stale visual recommendation"),
            ],
        )
        deck = DeckPlan(deck_title="Test", slides=[
            SlideSpec(
                slide_id="sk_deployment",
                title="Deployment and resilience protect operations",
                archetype="Deployment Architecture",
                diagram=DiagramSpec(
                    kind="deployment",
                    prompt="Show an AWS target runtime. Label every box TBC with SATS.",
                ),
            )
        ])

        ensure_diagrams_for_key_slides(
            deck,
            RFPUnderstanding(
                customer_name="SATS",
                summary="Design and build a secure digital catalogue application.",
                project_scope="Application development includes production deployment and release automation.",
            ),
            visual_briefs=None,
            technology_recommendations=recommendations,
        )

        prompt = deck.slides[0].diagram.prompt
        self.assertIn("Microsoft Azure", prompt)
        self.assertIn("Azure App Service", prompt)
        self.assertNotIn("AWS", prompt)
        self.assertNotIn("Amazon ECS", prompt)
        self.assertNotIn("TBC", prompt)
        self.assertNotIn("to be confirmed", prompt.lower())
        self.assertNotIn("landing-zone connectivity", prompt)

    def test_architecture_prompt_keeps_ai_bounded_without_repeated_sidecar(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            summary="Data hub for flight, ELP, SLA, reporting and support data.",
            project_scope=(
                "Design and build the data hub, ingest files, validate operational data, "
                "report exceptions, forecast uplift demand and deploy to production."
            ),
        )
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="arch",
                    title="Target architecture",
                    archetype="Architecture",
                    diagram=DiagramSpec(kind="architecture", prompt="Show the data hub."),
                )
            ],
        )

        enriched = ensure_diagrams_for_key_slides(deck, understanding)
        prompt = enriched.slides[0].diagram.prompt.lower()

        self.assertIn("one bounded component", prompt)
        self.assertIn("do not add a separate ai sidecar", prompt)

    def test_non_data_proposal_does_not_force_ai_slide(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="Client",
            summary="Provide a fixed classroom training schedule.",
            project_scope="Deliver instructor-led training sessions and attendance certificates.",
        )

        self.assertEqual(_ai_ml_opportunities(understanding), [])
        self.assertNotIn(
            "sk_ai_opportunities",
            {section["slide_id"] for section in _proposal_section_skeleton(understanding)},
        )

    def test_chunked_deck_plan_uses_placeholders_when_batch_call_fails(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            opportunity_title="Catering Uplift Data Hub",
            summary="Build a data hub with ELP extraction, SAP and KSMS API integration, AMS support, and ICCMS migration.",
            project_scope="Data hub, reporting, migration, warranty and support.",
        )
        fake_settings = SimpleNamespace(
            deck_plan_batch_size=4,
            reasoning_effort_deck_plan="medium",
            deck_plan_timeout_s=30,
        )

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema",
            side_effect=RuntimeError("blocked"),
        ):
            deck = _chunked_deck_plan(
                title="Catering Uplift Data Hub",
                understanding=understanding,
                narrative=None,
            )

        self.assertGreaterEqual(len(deck.slides), 24)
        self.assertTrue(any(slide.slide_id == "sk_arch" for slide in deck.slides))
        self.assertTrue(any(slide.slide_id == "sk_risks" for slide in deck.slides))

    def test_storyline_order_keeps_context_before_non_exec_solution(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="solution",
                    title="The proposed workstreams form one platform",
                    archetype="Solution Overview",
                    bullets=["Solution point"],
                ),
                SlideSpec(
                    slide_id="context",
                    title="Today's operating model creates avoidable effort",
                    archetype="Customer Context",
                    bullets=["Current challenge"],
                ),
                SlideSpec(
                    slide_id="exec",
                    title="Executive Summary",
                    archetype="Solution Overview",
                    bullets=["Executive thesis"],
                ),
            ],
        )

        ordered = order_deck(deck)
        ids = [slide.slide_id for slide in ordered.slides]

        self.assertLess(ids.index("exec"), ids.index("context"))
        self.assertLess(ids.index("context"), ids.index("solution"))

    def test_storyline_order_recognizes_exec_summary_by_title(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(slide_id="commercials", title="Commercials Structure", archetype="Commercials"),
                SlideSpec(slide_id="exec", title="Executive Summary", archetype="Content", bullets=["Thesis"]),
                SlideSpec(slide_id="context", title="Current Challenges", archetype="Customer Context"),
            ],
        )

        ordered = order_deck(deck)
        ids = [slide.slide_id for slide in ordered.slides]

        self.assertLess(ids.index("exec"), ids.index("context"))
        self.assertLess(ids.index("exec"), ids.index("commercials"))

    def test_redundant_storyline_slides_are_pruned(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(slide_id="title", title="Title", archetype="Title"),
                SlideSpec(slide_id="dep1", title="Deployment architecture", archetype="Deployment Architecture", bullets=["A"]),
                SlideSpec(slide_id="dep2", title="Recovery planning", archetype="High Availability & DR", bullets=["B"]),
                SlideSpec(slide_id="team1", title="Persistent squads", archetype="Team", bullets=["C"]),
                SlideSpec(slide_id="team2", title="Governance should resolve decisions", archetype="Delivery Plan", bullets=["Agile squad backlog sprint"]),
                SlideSpec(slide_id="team3", title="Agile roadmap", archetype="Timeline", bullets=["Sprint release backlog"]),
            ],
        )

        pruned = prune_redundant_storyline_slides(deck)
        ids = [slide.slide_id for slide in pruned.slides]

        self.assertIn("dep1", ids)
        self.assertIn("dep2", ids)
        delivery_count = sum(slide.slide_id.startswith("team") for slide in pruned.slides)
        self.assertGreaterEqual(delivery_count, 1)
        self.assertLessEqual(delivery_count, 2)

    def test_pruning_preserves_risks_technology_and_assumptions(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(slide_id="scope", title="Scope, Dependencies, and Open Points", archetype="Requirements"),
                SlideSpec(slide_id="assumptions", title="Assumptions and Dependencies", archetype="Assumptions & Dependencies"),
                SlideSpec(slide_id="risks", title="Key Risks and Mitigations", archetype="Risks", bullets=["Risk"]),
                SlideSpec(slide_id="tech", title="Proposed technologies span the delivery lifecycle", archetype="Software Bill of Materials", table={"headers": ["A"], "rows": [["B"]]}),
                SlideSpec(slide_id="delivery", title="Agile roadmap", archetype="Timeline", bullets=["Sprint backlog release"]),
            ],
        )

        pruned = prune_redundant_storyline_slides(deck)
        ids = [slide.slide_id for slide in pruned.slides]

        self.assertIn("scope", ids)
        self.assertIn("assumptions", ids)
        self.assertIn("risks", ids)
        self.assertIn("tech", ids)

    def test_pruning_keeps_value_deployment_and_ha_dr_as_distinct_proof_objects(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="value",
                    title="Proposal value and differentiators",
                    archetype="Value & Differentiators",
                    bullets=["Governance and engineering ownership reduce delivery risk."],
                ),
                SlideSpec(
                    slide_id="deployment",
                    title="Deployment topology protects controlled releases",
                    archetype="Deployment Architecture",
                    diagram=DiagramSpec(kind="deployment", prompt="deployment", approved=False),
                ),
                SlideSpec(
                    slide_id="hadr",
                    title="HA and DR protect business continuity",
                    archetype="High Availability & DR",
                    diagram=DiagramSpec(kind="hadr", prompt="recovery", approved=False),
                ),
            ],
        )

        pruned = prune_redundant_storyline_slides(deck)

        self.assertEqual(
            [slide.slide_id for slide in pruned.slides],
            ["value", "deployment", "hadr"],
        )

    def test_consulting_polish_turns_exec_summary_into_cards(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="exec",
                    title="Executive Summary",
                    archetype="Solution Overview",
                    bullets=["Centralise data", "Reduce manual work", "Improve operational control"],
                )
            ],
        )
        understanding = RFPUnderstanding(
            summary="SATS needs a controlled catering uplift data hub.",
            project_scope="Build a centralized data hub for operational reporting.",
            in_scope_work=["Data ingestion", "Validation", "Reporting"],
        )

        polished = consulting_grade_proposal_polish(deck, understanding=understanding, narrative=None)
        slide = polished.slides[0]

        self.assertGreaterEqual(len(slide.cards), 2)
        self.assertEqual(slide.bullets, [])

    def test_consulting_polish_turns_current_state_into_comparison(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="context",
                    title="Current challenges define priorities",
                    archetype="Customer Context",
                    bullets=["Manual consolidation delays reporting", "Interface readiness is a delivery dependency"],
                )
            ],
        )

        polished = consulting_grade_proposal_polish(deck)
        slide = polished.slides[0]

        self.assertIsNotNone(slide.comparison)
        self.assertEqual(slide.bullets, [])

    def test_consulting_polish_keeps_one_complete_key_message_sentence(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="summary",
                    title="Executive Summary",
                    archetype="Solution Overview",
                    key_message=(
                        "Establish one accountable command center for Technology Operations. "
                        "Scale services and automation only when operational evidence supports it."
                    ),
                )
            ],
        )

        polished = consulting_grade_proposal_polish(deck)

        self.assertEqual(
            polished.slides[0].key_message,
            "Establish one accountable command center for Technology Operations.",
        )

    def test_consulting_polish_turns_security_nfr_into_cards(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="nfr",
                    title="Security, observability and NFR controls are built in",
                    archetype="Architecture",
                    detailed_points=[
                        BulletPoint(
                            text="End-to-end secured communication",
                            sub_points=["Encrypt source-to-hub and hub-to-consumer paths"],
                        )
                    ],
                )
            ],
        )

        polished = consulting_grade_proposal_polish(deck)
        slide = polished.slides[0]

        self.assertEqual(len(slide.cards), 1)
        self.assertEqual(slide.detailed_points, [])

    def test_slide_density_bounds_counts_without_truncating_text(self) -> None:
        long_text = " ".join(f"word{i}" for i in range(60))
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(
                    slide_id="dense",
                    title="Dense",
                    archetype="Customer Context",
                    detailed_points=[
                        BulletPoint(
                            text="Current environment and constraints are lengthy",
                            sub_points=[long_text, long_text, long_text, long_text],
                        )
                    ],
                )
            ],
        )

        bounded = enforce_slide_density(deck)
        point = bounded.slides[0].detailed_points[0]

        self.assertEqual(point.text, "Current environment and constraints are lengthy")
        self.assertEqual(len(point.sub_points), 3)
        self.assertTrue(all(sub == long_text for sub in point.sub_points))

    def test_renderer_paginates_long_bullets_by_text_capacity(self) -> None:
        slide = SlideSpec(
            slide_id="long_bullets",
            title="Executive Summary",
            archetype="Content",
            bullets=[
                " ".join([f"Proposal point {idx} contains detailed rationale and customer implications"] * 16)
                for idx in range(4)
            ],
        )

        pages = _render_pages_for_slide(slide)

        self.assertGreater(len(pages), 1)
        self.assertTrue(all(len(page.bullets) < len(slide.bullets) for page in pages))
        self.assertTrue(all(page.archetype == "Content" for page in pages))
        self.assertTrue(all(page.layout_hint is None for page in pages))

    def test_renderer_paginates_solution_stack_at_six_rows(self) -> None:
        slide = SlideSpec(
            slide_id="tech",
            title="Proposed solution stack",
            archetype="Software Bill of Materials",
            table={
                "headers": ["Layer", "Technology", "Role", "Basis"],
                "rows": [[f"Layer {idx}", f"Service {idx}", "Role", "Proposed"] for idx in range(11)],
            },
        )

        pages = _render_pages_for_slide(slide)

        self.assertEqual([len(page.table["rows"]) for page in pages], [6, 5])

    def test_renderer_paginates_long_detailed_points_by_text_capacity(self) -> None:
        slide = SlideSpec(
            slide_id="long_detail",
            title="Risks and Mitigations",
            archetype="Risks",
            detailed_points=[
                BulletPoint(
                    text=f"Risk theme {idx}",
                    sub_points=[
                        " ".join(["This mitigation explanation includes evidence, impact, owner action, and acceptance dependency"] * 14)
                    ],
                )
                for idx in range(3)
            ],
        )

        pages = _render_pages_for_slide(slide)

        self.assertGreater(len(pages), 1)
        self.assertTrue(all(page.archetype == "Risks" for page in pages))
        self.assertTrue(all(page.layout_hint is None for page in pages))
        self.assertTrue(all(page.detailed_points for page in pages))
        self.assertTrue(all(not page.bullets for page in pages))

    def test_renderer_keeps_concise_five_step_action_slide_together(self) -> None:
        slide = SlideSpec(
            slide_id="next",
            title="Next Steps",
            archetype="Next Steps",
            bullets=[f"Confirm action {idx} with the accountable SATS owner." for idx in range(5)],
        )

        pages = _render_pages_for_slide(slide)

        self.assertEqual(len(pages), 1)

    def test_renderer_balances_four_detailed_points_across_continuations(self) -> None:
        slide = SlideSpec(
            slide_id="domains",
            title="Data domains",
            archetype="Content",
            detailed_points=[
                BulletPoint(
                    text=f"Domain {idx}",
                    sub_points=[" ".join(["Grounded domain control and source-of-record detail"] * 12)],
                )
                for idx in range(4)
            ],
        )

        pages = _render_pages_for_slide(slide)

        self.assertEqual([len(page.detailed_points) for page in pages], [2, 2])

    def test_renderer_preserves_cards_on_split_pages(self) -> None:
        from rfp2deck.core.schemas import Card

        slide = SlideSpec(
            slide_id="exec",
            title="Executive Summary",
            archetype="Solution Overview",
            cards=[
                Card(
                    heading=f"Executive theme {idx}",
                    body=" ".join(["Detailed customer implication and proposal response"] * 30),
                )
                for idx in range(3)
            ],
        )

        pages = _render_pages_for_slide(slide)

        self.assertGreater(len(pages), 1)
        self.assertEqual(sum(len(page.cards) for page in pages), 3)

    def test_renderer_concises_multisentence_key_message_as_complete_sentence(self) -> None:
        from rfp2deck.core.schemas import Card

        lead = (
            "Establish and operate a single, accountable Technology Operations function for "
            "Pearson Professional Assessments—combining follow-the-sun 24x7 monitoring, "
            "primarily L3 command-center ownership, disciplined ITIL-aligned processes, and "
            "measurable governance."
        )
        key_message = (
            lead
            + " The immediate priority is operational control and stabilization; expansion, "
            "automation, and AI-assisted capabilities follow only when data and process "
            "maturity support a credible business case."
        )
        slide = SlideSpec(
            slide_id="exec",
            title="Executive Summary",
            archetype="Solution Overview",
            key_message=key_message,
            cards=[
                Card(heading=f"Theme {idx}", body="A complete executive summary point.")
                for idx in range(3)
            ],
        )

        pages = _render_pages_for_slide(slide, native=True)

        self.assertEqual(len(pages), 1)
        self.assertEqual(pages[0].key_message, lead)
        self.assertEqual(pages[0].bullets, [])
        self.assertNotIn("AI-assisted capabilities", pages[0].key_message)

    def test_native_continuations_repeat_one_complete_key_message(self) -> None:
        from rfp2deck.core.schemas import Card

        key_message = (
            "Disciplined service controls connect operational evidence to accountable "
            "decisions and measurable improvement."
        )
        slide = SlideSpec(
            slide_id="controls",
            title="Service controls",
            archetype="Content",
            key_message=key_message,
            cards=[
                Card(
                    heading=f"Control {idx}",
                    body=" ".join(
                        ["Grounded control detail with ownership and customer implications"] * 30
                    ),
                )
                for idx in range(4)
            ],
        )

        pages = _render_pages_for_slide(slide, native=True)

        self.assertGreater(len(pages), 1)
        self.assertTrue(all(page.key_message == key_message for page in pages))
        self.assertTrue(all(page.key_message.endswith(".") for page in pages))

    def test_concise_key_message_uses_complete_clause_instead_of_word_cut(self) -> None:
        key_message = (
            "Accountable operations establish measurable control across the service; "
            + " ".join(["supporting evidence remains available for governance decisions"] * 8)
            + "."
        )

        concise = _concise_key_message(key_message)

        self.assertEqual(
            concise,
            "Accountable operations establish measurable control across the service.",
        )

    def test_renderer_preserves_native_layout_selection_for_regular_text_slides(self) -> None:
        slide = SlideSpec(
            slide_id="exec",
            title="Executive Summary",
            archetype="Solution Overview",
            bullets=["A complete proposal point.", "A second complete proposal point."],
        )

        pages = _render_pages_for_slide(slide)

        self.assertEqual(len(pages), 1)
        self.assertIsNone(pages[0].layout_hint)

    def test_renderer_adds_explanation_after_approved_diagram(self) -> None:
        slide = SlideSpec(
            slide_id="arch",
            title="Architecture",
            archetype="Architecture",
            bullets=["Explain source to target flow", "Explain controls and operations"],
            diagram=DiagramSpec(kind="architecture", prompt="draw architecture", approved=True),
        )

        pages = _render_pages_for_slide(slide, diagram_images={"arch": b"png"})

        self.assertEqual(len(pages), 2)
        self.assertIsNotNone(pages[0].diagram)
        self.assertEqual(pages[0].bullets, [])
        self.assertIsNone(pages[1].diagram)
        self.assertIsNone(pages[1].layout_hint)
        self.assertEqual(
            pages[1].bullets,
            ["Explain source to target flow", "Explain controls and operations"],
        )

    def test_renderer_does_not_create_companion_from_diagram_prompt_alone(self) -> None:
        slide = SlideSpec(
            slide_id="operating_model",
            title="Operating model",
            archetype="Solution Overview",
            diagram=DiagramSpec(
                kind="process",
                prompt="Show service ownership, incident flow, governance and improvement.",
                approved=True,
            ),
        )

        pages = _render_pages_for_slide(
            slide,
            diagram_images={"operating_model": b"png"},
        )

        self.assertEqual(len(pages), 1)
        self.assertIsNotNone(pages[0].diagram)

    def test_hcltech_diagram_uses_full_width_title_only_layout(self) -> None:
        from pptx import Presentation

        template = Path(".data/outputs/latest_user_deck.pptx")
        if not template.exists():
            self.skipTest("Full HCLTech template fixture is not available")
        prs = Presentation(template)
        slide = SlideSpec(
            slide_id="arch",
            title="Architecture",
            archetype="Architecture",
            diagram=DiagramSpec(kind="architecture", prompt="draw architecture", approved=True),
        )

        layout = _choose_hcltech_layout(prs, slide)

        self.assertEqual(layout.name.lower(), "title only")

    def test_rendered_slide_count_includes_diagram_companion(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(slide_id="title", title="Test", archetype="Title"),
                SlideSpec(
                    slide_id="arch",
                    title="Architecture",
                    archetype="Architecture",
                    bullets=[
                        "Explain the design decision.",
                        "Connect the decision to operational risk and customer value.",
                    ],
                    diagram=DiagramSpec(kind="architecture", prompt="draw architecture", approved=True),
                ),
            ],
        )

        self.assertEqual(rendered_slide_count(deck, {"arch": b"png"}), 3)

    def test_renderer_finds_approved_diagram_by_prompt_cache_key(self) -> None:
        slide = SlideSpec(
            slide_id="arch_v2",
            title="Architecture",
            archetype="Architecture",
            bullets=[
                "Explain source to target flow",
                "Explain the controls and operating boundary",
            ],
            diagram=DiagramSpec(kind="architecture", prompt="draw architecture", approved=True),
        )

        pages = _render_pages_for_slide(slide, diagram_images={"draw architecture": b"png"})

        self.assertEqual(len(pages), 2)
        self.assertIsNotNone(pages[0].diagram)
        self.assertEqual(pages[0].bullets, [])
        self.assertIsNone(pages[1].diagram)

    def test_density_preserves_diagram_explanation_for_companion_slide(self) -> None:
        slide = SlideSpec(
            slide_id="arch",
            title="Architecture",
            archetype="Architecture",
            bullets=[
                "Source systems feed a governed ingestion layer.",
                "Controls protect every interface.",
            ],
            diagram=DiagramSpec(kind="architecture", prompt="draw architecture", approved=True),
        )

        bounded = enforce_slide_density(DeckPlan(deck_title="Test", slides=[slide]))

        self.assertEqual(len(bounded.slides[0].bullets), 2)

    def test_renderer_preserves_diagram_picture_during_overlap_cleanup(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Title"
        picture_like = slide.shapes.add_picture(
            BytesIO(base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
            )),
            Inches(0.8),
            Inches(0.8),
            width=Inches(5),
            height=Inches(2),
        )
        spec = SlideSpec(
            slide_id="arch",
            title="Architecture",
            archetype="Architecture",
            diagram=DiagramSpec(kind="architecture", prompt="draw architecture", approved=True),
        )

        removed = _remove_overlapping_generated_pictures(slide, spec)

        self.assertEqual(removed, 0)
        self.assertIn(picture_like, list(slide.shapes))

    def test_renderer_removes_stale_singleton_continuation_title(self) -> None:
        slide = SlideSpec(
            slide_id="exec_page_1",
            title="Executive Summary (1 of 2)",
            archetype="Content",
            bullets=["Executive point"],
        )

        normalized = _normalize_singleton_continuation_titles([slide])

        self.assertEqual(normalized[0].title, "Executive Summary")

    def test_renderer_drops_exact_duplicate_render_slides(self) -> None:
        first = SlideSpec(
            slide_id="exec_1",
            title="Executive Summary",
            archetype="Solution Overview",
            bullets=["Centralise catering uplift data across ICC1 and ICC2."],
        )
        duplicate = SlideSpec(
            slide_id="exec_2",
            title="Executive Summary",
            archetype="Solution Overview",
            bullets=["Centralise catering uplift data across ICC1 and ICC2."],
        )

        slides = _dedupe_render_slides([first, duplicate])

        self.assertEqual([slide.slide_id for slide in slides], ["exec_1"])

    def test_renderer_drops_adjacent_prefix_duplicate_render_slides(self) -> None:
        first = SlideSpec(
            slide_id="scope_1",
            title="Scope & Boundaries",
            archetype="Requirements",
            bullets=[
                "Centralized data hub for catering uplift operations.",
                "Integration and reporting controls for operational users.",
            ],
        )
        prefix_duplicate = SlideSpec(
            slide_id="scope_2",
            title="Scope & Boundaries",
            archetype="Requirements",
            bullets=["Centralized data hub for catering uplift operations."],
        )

        slides = _dedupe_render_slides([first, prefix_duplicate])

        self.assertEqual([slide.slide_id for slide in slides], ["scope_1"])

    def test_renderer_falls_back_to_text_when_approved_diagram_image_is_missing(self) -> None:
        slide = SlideSpec(
            slide_id="solution",
            title="Solution overview",
            archetype="Solution Overview",
            bullets=["Diagram explanation: Create a left-to-right architecture view. Left: source systems. Right: reporting users."],
            diagram=DiagramSpec(kind="generic", prompt="draw solution", approved=True),
        )

        pages = _render_pages_for_slide(slide, diagram_images={})

        self.assertEqual(len(pages), 1)
        self.assertIsNone(pages[0].diagram)
        self.assertIsNone(pages[0].layout_hint)
        self.assertNotIn("Diagram explanation:", " ".join(pages[0].bullets))
        self.assertNotIn("left-to-right architecture view", " ".join(pages[0].bullets).lower())
        self.assertIn("asset is unavailable", pages[0].bullets[0].lower())

    def test_renderer_repairs_title_only_slide_with_diagram_prompt(self) -> None:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        spec = SlideSpec(
            slide_id="sk_solution",
            title="Solution overview",
            archetype="Solution Overview",
            diagram=DiagramSpec(kind="generic", prompt="Show source systems and data hub.", approved=True),
        )

        repaired = _repair_title_only_slide(slide, spec)

        self.assertTrue(repaired)
        self.assertTrue(any("data hub" in (getattr(shape, "text", "") or "").lower() for shape in slide.shapes))

    def test_testing_is_a_required_evidence_visual_but_ams_remains_grounded(self) -> None:
        deck = DeckPlan(
            deck_title="Test",
            slides=[
                SlideSpec(slide_id="test", title="Testing and acceptance create delivery evidence", archetype="Delivery Plan"),
                SlideSpec(slide_id="ams", title="Warranty and AMS sustain the live service", archetype="Delivery Plan"),
            ],
        )

        planned = ensure_diagrams_for_key_slides(deck, understanding=None)

        self.assertEqual(planned.slides[0].diagram.kind, "testing")
        self.assertIn("testing and acceptance evidence map", planned.slides[0].diagram.prompt.lower())
        self.assertIsNone(planned.slides[1].diagram)

    def test_testing_and_ams_are_tailored_to_named_solution_requirements(self) -> None:
        understanding = RFPUnderstanding(
            customer_name="SATS",
            opportunity_title="Catering Uplift Data Hub",
            summary="Replace ICCMS with a secure data hub and provide warranty and AMS support.",
            project_scope="Integrate operational sources, reconcile uplift data, cut over safely, and support the live service.",
            requirements=[
                Requirement(id="R-1", text="Integrate FIH and GP4 feeds, SAP and KSMS APIs, and SQ SFTP ELP files.", priority="must"),
                Requirement(id="R-2", text="Replace ICCMS while preserving operational processes and functional outcomes.", priority="must"),
                Requirement(id="R-3", text="Provide data validation, reconciliation, security audit evidence, and performance assurance.", priority="must"),
                Requirement(id="R-4", text="Provide warranty and AMS support for production integrations, reporting, and incident resolution.", priority="must"),
                Requirement(id="R-5", text="Show inventory availability in catalogue search results.", priority="must"),
            ],
        )

        testing_points = _testing_proposal_points(understanding)
        ams_points = _ams_proposal_points(understanding)
        testing_prompt = _build_diagram_prompt("testing", understanding)
        ams_prompt = _build_diagram_prompt("ams", understanding)

        self.assertTrue(any("FIH" in " ".join(point.sub_points) or "GP4" in " ".join(point.sub_points) for point in testing_points))
        self.assertTrue(any("ICCMS" in " ".join(point.sub_points) for point in testing_points))
        self.assertFalse(any("inventory availability" in " ".join(point.sub_points).lower() for point in testing_points))
        self.assertIn("FIH", testing_prompt)
        self.assertIn("ICCMS", testing_prompt)
        self.assertIn("acceptance owner", testing_prompt.lower())
        self.assertIn("do not show a textbook test pyramid", testing_prompt.lower())
        self.assertTrue(any("SAP" in " ".join(point.sub_points) or "KSMS" in " ".join(point.sub_points) for point in ams_points))
        self.assertIn("business-flow observability", ams_prompt.lower())
        self.assertIn("correction/replay", ams_prompt.lower())
        self.assertNotIn("to be agreed", ams_prompt.lower())
        self.assertIn("assumptions and dependencies slide", ams_prompt.lower())
        self.assertIn("do not use a generic l1/l2/l3 pyramid", ams_prompt.lower())

    def test_corporate_proxy_block_page_is_detected_deep_in_error_body(self) -> None:
        html = (
            "<!DOCTYPE html><html><head><title>Network Error</title></head>"
            "<body>"
            + ("x" * 5000)
            + "ThreatPulse Symantec BlueCoat Access Restricted Website "
            + "Category: Generative AI "
            + "URL: https:&#x2F;&#x2F;api.openai.com&#x2F;v1&#x2F;responses"
            + "</body></html>"
        )
        exc = SimpleNamespace(response=SimpleNamespace(text=html), status_code=500)

        self.assertTrue(_is_proxy_block_error(exc))

    def test_cloudflare_openai_520_is_not_detected_as_proxy_block(self) -> None:
        body = {
            "title": "Error 520: Web server is returning an unknown error",
            "status": 520,
            "zone": "api.openai.com",
            "cloudflare_error": True,
            "retryable": True,
        }
        exc = SimpleNamespace(response=SimpleNamespace(text=str(body)), status_code=520)

        self.assertFalse(_is_proxy_block_error(exc))

    def test_evidence_chunk_results_are_cached_for_downstream_retries(self) -> None:
        document = SourceDocument(
            document_id="doc-cache",
            name="RFP.pdf",
            document_type="base_rfp",
            authority="authoritative",
            text="--- PAGE 1 ---\nA mandatory requirement.",
        )
        chunk = split_source_document(document, 4000)[0]
        with TemporaryDirectory() as temp_dir:
            fake_settings = SimpleNamespace(
                understanding_evidence_cache=True,
                data_dir=Path(temp_dir),
                model_fast="test-fast-model",
                reasoning_effort_low="low",
                understanding_evidence_timeout_s=30,
            )
            with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
                "rfp2deck.agent.nodes.response_as_schema",
                return_value=SourceEvidenceBatch(source_document_id="", chunk_id=""),
            ) as mocked_response:
                first = _extract_evidence_chunk(chunk)
                second = _extract_evidence_chunk(chunk)

        self.assertEqual(mocked_response.call_count, 1)
        self.assertEqual(first.chunk_id, second.chunk_id)

    def test_large_package_node_uses_chunk_extraction_before_understanding(self) -> None:
        document = SourceDocument(
            document_id="doc-large",
            name="Large RFP.pdf",
            document_type="base_rfp",
            authority="authoritative",
            text="--- PAGE 1 ---\n" + ("Requirement text with locator evidence.\n" * 14500),
        )
        state = AgentState(
            rfp_text=document.text,
            template_info={},
            source_documents=[document],
        )
        fake_settings = SimpleNamespace(
            understanding_direct_max_chars=180000,
            understanding_evidence_chunk_chars=55000,
            understanding_evidence_max_chars=180000,
            understanding_evidence_workers=2,
            understanding_evidence_timeout_s=30,
            model_fast="test-fast-model",
            reasoning_effort_low="low",
        )

        def fake_response(prompt, schema, **kwargs):
            return SourceEvidenceBatch(source_document_id="", chunk_id="")

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema", side_effect=fake_response
        ) as mocked_response:
            result = extract_source_evidence(state)

        self.assertGreaterEqual(len(document.text), 533817)
        self.assertGreaterEqual(mocked_response.call_count, 9)
        self.assertLessEqual(mocked_response.call_count, 11)
        self.assertTrue(result["source_evidence"])
        self.assertLessEqual(len(result["evidence_text"]), 180000)

    def test_contextual_reference_chunk_failure_uses_bounded_source_fallback(self) -> None:
        text = (
            "Azure target architecture uses API Management, App Service, PostgreSQL, "
            "Blob Storage, Entra ID, Key Vault, monitoring and backup controls.\n"
            * 55
        )
        document = SourceDocument(
            document_id="doc-reference",
            name="SATS Catalog Solution.docx",
            document_type="supporting_reference",
            authority="contextual",
            text=text,
        )
        state = AgentState(
            rfp_text=text,
            template_info={},
            source_documents=[document],
        )
        fake_settings = SimpleNamespace(
            understanding_direct_max_chars=1000,
            understanding_evidence_chunk_chars=4000,
            understanding_evidence_max_chars=40000,
            understanding_evidence_workers=1,
            understanding_evidence_timeout_s=300,
            understanding_evidence_grace_s=60,
            understanding_contextual_evidence_llm_enabled=True,
            understanding_contextual_evidence_grace_s=30,
            understanding_evidence_cache=False,
            model_fast="test-fast-model",
            reasoning_effort_low="low",
        )

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema",
            side_effect=RuntimeError("background response timed out"),
        ) as structured_call:
            result = extract_source_evidence(state)

        self.assertTrue(result["source_evidence"])
        self.assertIn("Azure target architecture", result["evidence_text"])
        self.assertTrue(document.warnings)
        self.assertTrue(structured_call.call_args_list)
        self.assertTrue(
            all(
                call.kwargs["timeout_seconds"] == 60
                for call in structured_call.call_args_list
            )
        )
        self.assertTrue(
            all(
                call.kwargs["background_grace_seconds"] == 30
                for call in structured_call.call_args_list
            )
        )
        self.assertTrue(
            all(
                call.kwargs["recoverable_failure"] is True
                for call in structured_call.call_args_list
            )
        )

    def test_contextual_reference_uses_local_evidence_without_llm_by_default(self) -> None:
        text = (
            "The target architecture uses Azure API Management, Entra ID, Blob Storage, "
            "master data services and governed catalogue workflows.\n"
            * 35
        )
        document = SourceDocument(
            document_id="doc-local-reference",
            name="Architecture Reference.docx",
            document_type="supporting_reference",
            authority="contextual",
            text=text,
        )
        state = AgentState(
            rfp_text=text,
            template_info={},
            source_documents=[document],
        )
        fake_settings = SimpleNamespace(
            understanding_direct_max_chars=1000,
            understanding_evidence_chunk_chars=4000,
            understanding_evidence_max_chars=40000,
            understanding_evidence_workers=2,
            understanding_evidence_timeout_s=300,
            understanding_evidence_grace_s=60,
            understanding_contextual_evidence_llm_enabled=False,
            understanding_evidence_cache=False,
            model_fast="test-fast-model",
            reasoning_effort_low="low",
        )

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema"
        ) as structured_call:
            result = extract_source_evidence(state)

        structured_call.assert_not_called()
        self.assertIn("Azure API Management", result["evidence_text"])
        self.assertFalse(document.warnings)

    def test_authoritative_chunk_failure_still_stops_evidence_extraction(self) -> None:
        text = "The solution must retain this authoritative requirement.\n" * 100
        document = SourceDocument(
            document_id="doc-authoritative",
            name="RFP.docx",
            document_type="base_rfp",
            authority="authoritative",
            text=text,
        )
        state = AgentState(
            rfp_text=text,
            template_info={},
            source_documents=[document],
        )
        fake_settings = SimpleNamespace(
            understanding_direct_max_chars=1000,
            understanding_evidence_chunk_chars=4000,
            understanding_evidence_max_chars=40000,
            understanding_evidence_workers=1,
            understanding_evidence_timeout_s=30,
            understanding_evidence_grace_s=5,
            understanding_evidence_cache=False,
            model_fast="test-fast-model",
            reasoning_effort_low="low",
        )

        with patch("rfp2deck.agent.nodes.settings", fake_settings), patch(
            "rfp2deck.agent.nodes.response_as_schema",
            side_effect=RuntimeError("background response timed out"),
        ), self.assertRaises(RuntimeError) as raised:
            extract_source_evidence(state)

        self.assertIn("Evidence extraction failed for RFP.docx", str(raised.exception))

    def test_large_source_is_split_into_bounded_locator_aware_chunks(self) -> None:
        text = "--- PAGE 1 ---\n" + ("Requirement A must be retained.\n" * 350)
        document = SourceDocument(
            document_id="doc-large",
            name="Large RFP.pdf",
            document_type="base_rfp",
            authority="authoritative",
            text=text,
        )

        chunks = split_source_document(document, 4000)

        self.assertGreater(len(chunks), 1)
        self.assertTrue(all(len(chunk.text) <= 4000 for chunk in chunks))
        self.assertTrue(all(chunk.document.document_id == "doc-large" for chunk in chunks))
        self.assertIn("LOCATOR CONTEXT", chunks[1].text)

    def test_chunk_evidence_merge_deduplicates_overlap_and_preserves_requirements(self) -> None:
        requirement = Requirement(
            id="R-1",
            text="Integrate all nominated source systems",
            source_ref="doc-rfp / page 20",
            source_document_ids=["doc-rfp"],
        )
        batches = [
            SourceEvidenceBatch(
                source_document_id="doc-rfp",
                chunk_id="chunk-1",
                summary_points=["Integration is in scope"],
                requirements=[requirement],
            ),
            SourceEvidenceBatch(
                source_document_id="doc-rfp",
                chunk_id="chunk-2",
                summary_points=["Integration is in scope"],
                requirements=[requirement.model_copy(deep=True)],
            ),
        ]

        merged = merge_evidence_batches(batches)
        rendered = render_evidence_for_prompt(merged, 40000)

        self.assertEqual(len(merged.requirements), 1)
        self.assertEqual(merged.summary_points, ["Integration is in scope"])
        self.assertIn("doc-rfp / page 20", rendered)

    def test_large_requirement_register_uses_lossless_indexed_compaction(self) -> None:
        requirements = [
            Requirement(
                id=f"REQ-{index:04d}",
                text=(
                    f"Requirement {index} requires validated processing, audit evidence, "
                    "operational monitoring, exception handling, and retained traceability."
                ),
                source_ref=f"doc-rfp / sheet Requirements / row {index + 2}",
                source_refs=[f"doc-rfp / sheet Requirements / row {index + 2}"],
                source_document_ids=["doc-rfp"],
            )
            for index in range(700)
        ]
        evidence = SourceEvidenceBatch(
            source_document_id="rfp-package",
            chunk_id="merged-evidence",
            requirements=requirements,
        )

        rendered = render_evidence_for_prompt(evidence, 180000)

        self.assertLessEqual(len(rendered), 180000)
        self.assertIn("indexed-rfp-evidence-v1", rendered)
        self.assertIn("Requirement 699", rendered)
        self.assertIn("row 701", rendered)

    def test_oversized_requirement_register_uses_budgeted_indexed_compaction(self) -> None:
        requirements = [
            Requirement(
                id=f"REQ-{index:04d}",
                text=(
                    f"Requirement {index} requires validated processing, audit evidence, "
                    "operational monitoring, exception handling, retained traceability, "
                    "implementation controls, support procedures, and reporting."
                ),
                priority="must" if index < 10 else "should",
                source_ref=f"doc-rfp / sheet Requirements / row {index + 2}",
                source_refs=[f"doc-rfp / sheet Requirements / row {index + 2}"],
                source_document_ids=["doc-rfp"],
            )
            for index in range(1200)
        ]
        evidence = SourceEvidenceBatch(
            source_document_id="rfp-package",
            chunk_id="merged-evidence",
            requirements=requirements,
        )

        rendered = render_evidence_for_prompt(evidence, 40000)

        self.assertLessEqual(len(rendered), 40000)
        self.assertIn('"budgeted":true', rendered)
        self.assertIn('"omitted_requirements":', rendered)
        self.assertIn("Requirement 0", rendered)

    def test_base_rfp_is_not_reclassified_by_procurement_wording(self) -> None:
        document_type, authority = classify_source(
            "Request for Tender.pdf",
            "Vendors may submit clarification questions before the deadline. Scope of Work follows.",
            "primary",
        )

        self.assertEqual(document_type, "base_rfp")
        self.assertEqual(authority, "authoritative")

    def test_understanding_schema_can_be_fully_inlined_for_structured_output(self) -> None:
        def contains_ref(value: object) -> bool:
            if isinstance(value, dict):
                return "$ref" in value or any(contains_ref(item) for item in value.values())
            if isinstance(value, list):
                return any(contains_ref(item) for item in value)
            return False

        for model in (RFPUnderstanding, SourceEvidenceBatch):
            schema = model.model_json_schema()
            defs = schema.get("$defs", {})
            inlined = _dereference(schema, defs)
            inlined.pop("$defs", None)
            strict = _make_strict(inlined)
            self.assertFalse(contains_ref(strict), model.__name__)

    def test_understanding_prompt_accepts_source_reconciliation(self) -> None:
        prompt = RFP_UNDERSTAND_PROMPT.format(
            rfp_text="=== SOURCE DOCUMENT base ===\nRequirement text",
            rfp_focus_guide="Scope section detected",
            source_reconciliation='{"precedence_summary": []}',
        )

        self.assertIn("SOURCE_RECONCILIATION", prompt)
        self.assertIn("vendor question is context only", prompt.lower())

    def test_xlsx_clarifications_preserve_rows_and_answer_authority(self) -> None:
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Clarifications"
        sheet.append(["Question ID", "Bidder's Query", "SATS Response"])
        sheet.append(["Q-01", "Must the solution use Ariba?", "No. Ariba is submission-only."])
        sheet.append(["Q-02", "What is the required RTO?", ""])
        stream = BytesIO()
        workbook.save(stream)

        document, clarifications = parse_source_document(
            "Customer Clarifications 2026-07-15.xlsx",
            stream.getvalue(),
            role="clarification",
        )

        self.assertEqual(document.document_type, "customer_clarification")
        self.assertEqual(document.issue_date, "2026-07-15")
        self.assertEqual(document.locator_format, "sheet/row")
        self.assertIn('[SHEET "Clarifications"][ROW 2]', document.text)
        self.assertEqual(len(clarifications), 2)
        self.assertEqual(clarifications[0].authority, "authoritative")
        self.assertEqual(clarifications[0].status, "active")
        self.assertEqual(clarifications[1].authority, "non_authoritative")
        self.assertEqual(clarifications[1].status, "unresolved")

    def test_docx_tables_are_available_as_row_level_evidence(self) -> None:
        document = Document()
        document.add_heading("Scope", level=1)
        document.add_paragraph("Create a centralized data repository.")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Requirement"
        table.cell(0, 1).text = "Priority"
        table.cell(1, 0).text = "Integrate source systems"
        table.cell(1, 1).text = "Mandatory"
        stream = BytesIO()
        document.save(stream)

        parsed = parse_docx(stream.getvalue())

        self.assertEqual(parsed.table_count, 1)
        self.assertIn("[PARAGRAPH 2]", parsed.text)
        self.assertIn("[TABLE 1][ROW 2]", parsed.text)
        self.assertIn("Integrate source systems", parsed.text)

    def test_source_precedence_orders_addenda_before_base_and_supporting(self) -> None:
        sources = [
            SourceDocument(
                document_id="support",
                name="Reference.docx",
                document_type="supporting_reference",
                authority="contextual",
                text="Reference",
            ),
            SourceDocument(
                document_id="base",
                name="RFP.pdf",
                document_type="base_rfp",
                authority="authoritative",
                text="Base",
            ),
            SourceDocument(
                document_id="addendum-old",
                name="Addendum 1.pdf",
                document_type="customer_addendum",
                authority="binding",
                issue_date="2026-07-01",
                text="Old amendment",
            ),
            SourceDocument(
                document_id="addendum-new",
                name="Addendum 2.pdf",
                document_type="customer_addendum",
                authority="binding",
                issue_date="2026-07-15",
                text="New amendment",
            ),
        ]

        ordered = sort_sources_by_precedence(sources)

        self.assertEqual(
            [source.document_id for source in ordered],
            ["addendum-new", "addendum-old", "base", "support"],
        )
        rendered = render_source_package(sources)
        self.assertLess(rendered.index("addendum-new"), rendered.index("base"))
        self.assertIn("vendor questions alone are not requirements", rendered.lower())

    def test_reconciliation_flags_unanswered_questions(self) -> None:
        workbook = Workbook()
        sheet = workbook.active
        sheet.append(["ID", "Question", "Response"])
        sheet.append(["Q1", "Required hosting region?", ""])
        stream = BytesIO()
        workbook.save(stream)
        document, clarifications = parse_source_document(
            "Clarifications.xlsx", stream.getvalue(), role="clarification"
        )

        reconciliation = build_source_reconciliation([document], clarifications)

        self.assertEqual(reconciliation.unresolved_questions, [clarifications[0].source_ref])
        self.assertIn("vendor question", " ".join(reconciliation.precedence_summary).lower())

    def test_requirement_schema_remains_backward_compatible(self) -> None:
        requirement = Requirement(id="R-1", text="Retain data for seven years")

        self.assertEqual(requirement.status, "active")
        self.assertEqual(requirement.source_refs, [])
        self.assertEqual(requirement.confidence, 1.0)

    def test_traceability_report_retains_source_provenance(self) -> None:
        requirement = Requirement(
            id="R-1",
            text="Retain data for seven years",
            priority="must",
            source_refs=["doc-rfp / page 12"],
            source_document_ids=["doc-rfp"],
            authority="binding",
            status="clarified",
        )
        understanding = RFPUnderstanding(summary="Retention requirement", requirements=[requirement])
        deck = DeckPlan(
            deck_title="Proposal",
            slides=[SlideSpec(slide_id="requirements", title="Retention", rfps=["R-1"])],
        )

        report = build_traceability_report(understanding, deck)

        self.assertEqual(report.coverage[0].source_refs, ["doc-rfp / page 12"])
        self.assertEqual(report.coverage[0].source_document_ids, ["doc-rfp"])
        self.assertEqual(report.coverage[0].status, "clarified")


if __name__ == "__main__":
    unittest.main()
